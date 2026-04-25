import asyncio
from pathlib import Path

from pptx import Presentation

from mellow_link.modules.rebuild_assistant.postprocess.consulting_deck import (
    build_consulting_deck,
)
from mellow_link.modules.rebuild_assistant.postprocess.schemas import (
    ConsultingDeck,
    ConsultingDeckChapter,
    ConsultingMinContract,
    ConsultingDeckSection,
)
from mellow_link.modules.rebuild_assistant.postprocess.slide_schema import (
    build_slide_schema,
)
from mellow_link.services import DocumentRequest, DocumentType, create_document_service


def _fifo_contract() -> ConsultingMinContract:
    return ConsultingMinContract(
        as_is=[
            "대표 도메인 범위는 외화 입출금 FIFO 중심으로 정리합니다.",
            "기존 스키마 호환성을 유지해야 하므로 API 백엔드 분리 시 DB 계약을 우선 보존합니다.",
        ],
        gap=[
            "입금 lot 잔량과 출금 lot 소진 순서를 분리해야 동일 거래의 원가 계산과 lot 추적이 흔들리지 않습니다.",
            "GAP_AMT 계산 기준을 같은 정책으로 고정해야 환차손익과 출금 금액 연결이 일관되게 유지됩니다.",
            "전표와 GL 반영 기준번호를 lot 계산 결과와 함께 관리해야 회계 반영 누락과 재처리 오류를 줄일 수 있습니다.",
        ],
        risks=[
            "FIFO lot 소진 순서가 바뀌면 동일 출금 건의 원가와 lot 추적 결과가 달라질 수 있습니다.",
            "GAP_AMT 계산 기준이 흔들리면 환차손익과 전표 금액이 서로 어긋날 수 있습니다.",
            "전표 생성과 GL_INTERFACE 반영 기준번호가 분리되면 회계 연계 누락이 발생할 수 있습니다.",
            "입력 자산이 제한적이므로 제안은 설계 초안 수준이며 추가 파일 확인이 필요합니다.",
        ],
        process_flow=[
            "1주차: 입금 lot 원장과 출금 lot 소진 순서를 구조화합니다.",
            "2주차: FIFO 계산·회계 연계 분리 구조 기준으로 환차손익 계산과 회계 연계 구조를 설계합니다.",
            "3주차: 외화 입출금 FIFO 계산 서비스와 전표 생성 흐름에 핵심 규칙을 반영합니다.",
            "4주차: 외화 입출금 FIFO lot 추적, 환차손익, 전표 정합성을 규칙 기준으로 검증합니다.",
        ],
        actions=[
            "입금 lot 적재와 출금 lot 소진 계산을 별도 FIFO 계산 계층으로 분리합니다.",
            "환차손익 계산을 lot별 취득 환율과 출금 환율 비교 정책으로 고정합니다.",
            "전표 생성과 GL_INTERFACE 적재를 계산 결과와 같은 거래 기준번호로 연결합니다.",
        ],
        rules=[
            "통화 및 계좌 식별값 유지: 통화 코드와 계좌 식별 값은 lot 계산, 전표, GL 반영 전 과정에서 일관되게 유지합니다.",
            "환차손익 계산: lot별 취득 환율과 출금 환율 차이로 환차손익을 계산합니다.",
            "전표 및 GL 반영: 입출금 및 환차손익 결과는 전표와 GL 인터페이스에 동일 기준번호로 반영합니다.",
            "FIFO lot 소진 순서: 외화 출금은 입금 lot 잔량을 FIFO 순서로 차감합니다.",
        ],
    )


def _build_fifo_slide_schema():
    deck = ConsultingDeck.model_validate(
        build_consulting_deck(
            _fifo_contract(),
            project_name="선입선출 외화관리",
            client_name="원자력연료",
            surface_mode="internal",
        )
    )
    return build_slide_schema(deck)


def _first_shape_texts(prs: Presentation) -> list[str]:
    titles: list[str] = []
    for slide in prs.slides:
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and (shape.text or "").strip():
                title = shape.text.strip().splitlines()[0].strip()
                break
        titles.append(title)
    return titles


def _all_slide_text(prs: Presentation, index: int) -> str:
    texts: list[str] = []
    for shape in prs.slides[index].shapes:
        if hasattr(shape, "text") and (shape.text or "").strip():
            texts.append(shape.text.strip())
    return "\n".join(texts)


def test_build_slide_schema_maps_fifo_deck_to_five_slide_types():
    slide_schema = _build_fifo_slide_schema()

    assert slide_schema.schema_version == "slide_schema.v1"
    assert len(slide_schema.slides) == 6
    assert [slide.slide_type for slide in slide_schema.slides] == [
        "overview",
        "as_is_gap",
        "as_is_gap",
        "flow",
        "design",
        "vision",
    ]
    assert [slide.layout_hint for slide in slide_schema.slides] == [
        "overview_context_scope_constraints",
        "three_column_compare_with_risk_strip",
        "gap_risk_continuation",
        "timeline_horizontal",
        "rule_cards_with_meta_sidebar",
        "future_state_pillars",
    ]
    assert slide_schema.slides[0].context_bullets
    assert slide_schema.slides[0].constraint_bullets
    assert slide_schema.slides[1].to_be_bullets
    assert slide_schema.slides[1].risk_bullets
    assert slide_schema.slides[2].is_continuation is True
    assert slide_schema.slides[2].continuation_value == "retain"
    assert slide_schema.slides[2].continuation_value_score >= 2
    assert len(slide_schema.slides[2].gap_bullets) >= 2
    assert slide_schema.slides[3].action_bullets
    assert slide_schema.slides[4].entity_blocks
    assert slide_schema.slides[4].flow_bullets
    assert slide_schema.slides[4].absorbed_summary_text == ""
    assert slide_schema.slides[5].effect_bullets == []
    assert [slide.slide_id for slide in slide_schema.slides if slide.is_continuation] == ["approach_02"]


def test_build_slide_schema_applies_length_limits_and_overflow():
    oversized = ConsultingDeck(
        project_name="과밀 테스트",
        client_name="ACME",
        surface_mode="internal",
        chapters=[
            ConsultingDeckChapter(
                chapter_key="overview",
                title="컨설팅 개요",
                sections=[
                    ConsultingDeckSection(
                        section_key="as_is",
                        title="현행 요약",
                        items=[
                            "외화 입출금 FIFO 분석 범위를 기준으로 기존 인터페이스와 lot 계산 구조를 동시에 재정리해야 하는 현재 상태를 아주 길게 설명하는 문장입니다."
                        ],
                    )
                ],
            ),
            ConsultingDeckChapter(
                chapter_key="approach",
                title="컨설팅 전개",
                sections=[
                    ConsultingDeckSection(
                        section_key="gap",
                        title="AS-IS / TO-BE GAP",
                        items=[
                            f"입금 lot 잔량과 출금 lot 소진 순서를 분리해야 동일 거래의 원가 계산과 lot 추적이 흔들리지 않는 GAP {index}입니다."
                            for index in range(1, 7)
                        ],
                    ),
                    ConsultingDeckSection(
                        section_key="risks",
                        title="검토 포인트",
                        items=[
                            f"전표 생성과 GL 인터페이스 기준번호가 분리되면 회계 연계 누락과 재처리 오류가 함께 발생할 수 있는 위험 {index}입니다."
                            for index in range(1, 9)
                        ],
                    ),
                ],
            ),
            ConsultingDeckChapter(
                chapter_key="implementation",
                title="컨설팅 구현",
                sections=[
                    ConsultingDeckSection(
                        section_key="process_flow",
                        title="단계별 추진 흐름",
                        items=[
                            f"{index}주차: 외화 입출금 FIFO lot 계산과 전표 반영 구조를 아주 길게 설명하는 단계입니다."
                            for index in range(1, 7)
                        ],
                    ),
                    ConsultingDeckSection(
                        section_key="actions",
                        title="중점 실행 과제",
                        items=[
                            f"입금 lot 적재와 출금 lot 소진 계산을 별도 FIFO 계산 계층으로 분리하는 매우 긴 실행 문장 {index}입니다."
                            for index in range(1, 8)
                        ],
                    ),
                ],
            ),
            ConsultingDeckChapter(
                chapter_key="design",
                title="컨설팅 설계",
                sections=[
                    ConsultingDeckSection(
                        section_key="rules",
                        title="핵심 규칙",
                        items=[
                            f"통화 및 계좌 식별값 유지 {index}: 통화 코드와 계좌 식별 값은 lot 계산, 전표, GL 반영 전 과정에서 일관되게 유지해야 한다는 긴 규칙 설명입니다."
                            for index in range(1, 6)
                        ],
                    ),
                    ConsultingDeckSection(
                        section_key="process_flow",
                        title="계산 / 전표 / GL 흐름",
                        items=[
                            f"외화 입출금 FIFO 구조를 계산, 전표, GL 기준으로 정리하는 설계 흐름 {index}입니다."
                            for index in range(1, 5)
                        ],
                    ),
                ],
            ),
            ConsultingDeckChapter(
                chapter_key="vision",
                title="컨설팅 비전",
                sections=[
                    ConsultingDeckSection(
                        section_key="actions",
                        title="적용 방향",
                        items=[
                            f"입금 lot 적재와 출금 lot 소진 계산을 분리 체계로 전환하는 미래 방향 {index}입니다."
                            for index in range(1, 8)
                        ],
                    )
                ],
            ),
        ],
    )

    slide_schema = build_slide_schema(oversized)
    approach_slides = [slide for slide in slide_schema.slides if slide.slide_type == "as_is_gap"]
    flow_slides = [slide for slide in slide_schema.slides if slide.slide_type == "flow"]
    design_slides = [slide for slide in slide_schema.slides if slide.slide_type == "design"]
    vision_slides = [slide for slide in slide_schema.slides if slide.slide_type == "vision"]

    assert len(approach_slides) >= 2
    assert len(flow_slides) >= 2
    assert len(design_slides) >= 2
    assert len(vision_slides) >= 2
    assert approach_slides[1].is_continuation is True
    assert approach_slides[1].layout_hint == "gap_risk_continuation"
    assert approach_slides[1].continuation_value == "retain"
    assert flow_slides[0].layout_hint == "stacked_flow"
    assert flow_slides[1].continuation_reason == "flow_overflow"
    assert flow_slides[1].continuation_value == "absorb_candidate"
    assert design_slides[1].continuation_reason == "design_rule_card_overflow"
    assert design_slides[1].continuation_value == "retain"
    assert vision_slides[1].layout_hint == "vision_continuation"
    assert vision_slides[1].continuation_value == "retain"
    assert all(len(item) <= 66 for slide in approach_slides for item in slide.gap_bullets)
    assert all(len(item) <= 56 for slide in approach_slides for item in slide.risk_bullets)
    assert all(len(step.step_text) <= 44 for slide in flow_slides for step in slide.steps)
    assert all(len(card.title) <= 20 and len(card.body) <= 56 for slide in slide_schema.slides for card in slide.rule_cards)
    assert all(slide.max_text_objects <= 7 for slide in slide_schema.slides)


def test_build_slide_schema_preserves_minimum_gap_and_rule_information():
    sparse = ConsultingDeck(
        project_name="정보 보강 테스트",
        client_name="ACME",
        surface_mode="internal",
        chapters=[
            ConsultingDeckChapter(
                chapter_key="overview",
                title="컨설팅 개요",
                sections=[ConsultingDeckSection(section_key="as_is", title="현행 요약", items=["외화 입출금 FIFO 구조를 검토합니다."])],
            ),
            ConsultingDeckChapter(
                chapter_key="approach",
                title="컨설팅 전개",
                sections=[
                    ConsultingDeckSection(section_key="gap", title="AS-IS / TO-BE GAP", items=["입금 lot와 출금 lot 계산을 분리해야 정합성이 유지됩니다."]),
                    ConsultingDeckSection(section_key="risks", title="검토 포인트", items=["정합성 검토 근거가 부족하면 lot 추적 리스크가 남습니다."]),
                ],
            ),
            ConsultingDeckChapter(
                chapter_key="implementation",
                title="컨설팅 구현",
                sections=[
                    ConsultingDeckSection(section_key="process_flow", title="단계별 추진 흐름", items=["1주차: FIFO 계산 서비스를 구조화합니다."]),
                    ConsultingDeckSection(section_key="actions", title="중점 실행 과제", items=["FIFO 계산 서비스를 구조화합니다."]),
                ],
            ),
            ConsultingDeckChapter(
                chapter_key="design",
                title="컨설팅 설계",
                sections=[
                    ConsultingDeckSection(section_key="rules", title="핵심 규칙", items=["FIFO lot 소진 순서: 외화 출금은 입금 lot 잔량을 FIFO 순서로 차감합니다."]),
                    ConsultingDeckSection(section_key="process_flow", title="계산 / 전표 / GL 흐름", items=["전표와 GL 흐름을 FIFO 계산 구조와 함께 정리합니다."]),
                ],
            ),
            ConsultingDeckChapter(
                chapter_key="vision",
                title="컨설팅 비전",
                sections=[ConsultingDeckSection(section_key="actions", title="적용 방향", items=["FIFO 계산 기준을 단일 구조로 정리합니다."])],
            ),
        ],
    )

    slide_schema = build_slide_schema(sparse)
    approach_slide = next(slide for slide in slide_schema.slides if slide.slide_type == "as_is_gap")
    design_slide = next(slide for slide in slide_schema.slides if slide.slide_type == "design")

    assert len(approach_slide.gap_bullets) >= 2
    assert len(design_slide.rule_cards) >= 2


def test_build_slide_schema_reabsorbs_flow_only_design_continuation():
    absorb_case = ConsultingDeck(
        project_name="흡수 테스트",
        client_name="ACME",
        surface_mode="internal",
        chapters=[
            ConsultingDeckChapter(
                chapter_key="overview",
                title="컨설팅 개요",
                sections=[ConsultingDeckSection(section_key="as_is", title="현행 요약", items=["설계 보강 검토"])],
            ),
            ConsultingDeckChapter(
                chapter_key="approach",
                title="컨설팅 전개",
                sections=[
                    ConsultingDeckSection(section_key="gap", title="AS-IS / TO-BE GAP", items=["설계 기준 정리가 필요합니다."]),
                    ConsultingDeckSection(section_key="risks", title="검토 포인트", items=["기준번호 추적성이 약합니다."]),
                ],
            ),
            ConsultingDeckChapter(
                chapter_key="implementation",
                title="컨설팅 구현",
                sections=[
                    ConsultingDeckSection(section_key="process_flow", title="단계별 추진 흐름", items=["1주차: 설계 흐름을 정리합니다."]),
                    ConsultingDeckSection(section_key="actions", title="중점 실행 과제", items=["설계 기준을 정리합니다."]),
                ],
            ),
            ConsultingDeckChapter(
                chapter_key="design",
                title="컨설팅 설계",
                sections=[
                    ConsultingDeckSection(
                        section_key="rules",
                        title="핵심 규칙",
                        items=[
                            "통화 및 계좌 식별값 유지: 통화 코드와 계좌 식별 값은 lot 계산, 전표, GL 반영 전 과정에서 일관되게 유지합니다.",
                            "환차손익 계산: lot별 취득 환율과 출금 환율 차이로 환차손익을 계산합니다.",
                            "전표 및 GL 반영: 입출금 및 환차손익 결과는 전표와 GL 인터페이스에 동일 기준번호로 반영합니다.",
                            "FIFO lot 소진 순서: 외화 출금은 입금 lot 잔량을 FIFO 순서로 차감합니다.",
                        ],
                    ),
                    ConsultingDeckSection(
                        section_key="process_flow",
                        title="계산 / 전표 / GL 흐름",
                        items=[
                            "lot 원장과 환차손익, 전표, GL 기준을 같은 설계 흐름으로 정렬하고 기준번호 연계를 검토합니다 XXXXX XXXXX XXXXX",
                            "lot 잔량, 통화, 계좌, 기준번호를 같은 구조에서 관리하고 회계 연계 차이를 검증합니다 XXXXX XXXXX XXXXX",
                            "전표와 GL 인터페이스 반영 흐름을 계산 기준과 함께 확인하고 후속 정합성 검토를 이어갑니다 XXXXX XXXXX XXXXX",
                        ],
                    ),
                ],
            ),
            ConsultingDeckChapter(
                chapter_key="vision",
                title="컨설팅 비전",
                sections=[ConsultingDeckSection(section_key="actions", title="적용 방향", items=["적용 방향"])],
            ),
        ],
    )

    slide_schema = build_slide_schema(absorb_case)
    design_slides = [slide for slide in slide_schema.slides if slide.slide_type == "design"]

    assert len(design_slides) == 1
    assert design_slides[0].layout_hint == "rule_cards_with_meta_sidebar"
    assert design_slides[0].flow_bullets
    assert design_slides[0].is_continuation is False
    assert design_slides[0].absorbed_summary_text == ""


def test_document_service_renders_pptx_from_slide_schema(tmp_path: Path):
    slide_schema = _build_fifo_slide_schema()
    service = create_document_service(output_dir=tmp_path)

    async def _render():
        await service.initialize()
        try:
            return await service.generate(
                DocumentRequest(
                    content="",
                    output_type=DocumentType.PPTX,
                    title="선입선출 외화관리",
                    filename="fifo_slide_schema",
                    payload=slide_schema.model_dump(),
                    style_options={"renderer": "slide_schema"},
                )
            )
        finally:
            await service.shutdown()

    result = asyncio.run(_render())
    prs = Presentation(str(result.output_path))
    titles = _first_shape_texts(prs)

    assert len(prs.slides) == 6
    assert titles == [
        "컨설팅 개요",
        "컨설팅 전개",
        "컨설팅 전개 (2)",
        "컨설팅 구현",
        "컨설팅 설계",
        "컨설팅 비전",
    ]
    assert "독립 연속 장표" in _all_slide_text(prs, 2)
    assert len(prs.slides[3].shapes) > 3


def test_document_service_renders_continuation_titles_for_long_case(tmp_path: Path):
    oversized = ConsultingDeck(
        project_name="과밀 테스트",
        client_name="ACME",
        surface_mode="internal",
        chapters=[
            ConsultingDeckChapter(
                chapter_key="overview",
                title="컨설팅 개요",
                sections=[ConsultingDeckSection(section_key="as_is", title="현행 요약", items=["외화 입출금 FIFO 분석 범위를 기준으로 기존 인터페이스와 lot 계산 구조를 동시에 재정리해야 하는 현재 상태를 아주 길게 설명하는 문장입니다."])],
            ),
            ConsultingDeckChapter(
                chapter_key="approach",
                title="컨설팅 전개",
                sections=[
                    ConsultingDeckSection(section_key="gap", title="AS-IS / TO-BE GAP", items=[f"입금 lot 잔량과 출금 lot 소진 순서를 분리해야 동일 거래의 원가 계산과 lot 추적이 흔들리지 않는 GAP {index}입니다." for index in range(1, 7)]),
                    ConsultingDeckSection(section_key="risks", title="검토 포인트", items=[f"전표 생성과 GL 인터페이스 기준번호가 분리되면 회계 연계 누락과 재처리 오류가 함께 발생할 수 있는 위험 {index}입니다." for index in range(1, 9)]),
                ],
            ),
            ConsultingDeckChapter(
                chapter_key="implementation",
                title="컨설팅 구현",
                sections=[
                    ConsultingDeckSection(section_key="process_flow", title="단계별 추진 흐름", items=[f"{index}주차: 외화 입출금 FIFO lot 계산과 전표 반영 구조를 아주 길게 설명하는 단계입니다." for index in range(1, 7)]),
                    ConsultingDeckSection(section_key="actions", title="중점 실행 과제", items=[f"입금 lot 적재와 출금 lot 소진 계산을 별도 FIFO 계산 계층으로 분리하는 매우 긴 실행 문장 {index}입니다." for index in range(1, 8)]),
                ],
            ),
            ConsultingDeckChapter(
                chapter_key="design",
                title="컨설팅 설계",
                sections=[
                    ConsultingDeckSection(section_key="rules", title="핵심 규칙", items=[f"통화 및 계좌 식별값 유지 {index}: 통화 코드와 계좌 식별 값은 lot 계산, 전표, GL 반영 전 과정에서 일관되게 유지해야 한다는 긴 규칙 설명입니다." for index in range(1, 6)]),
                    ConsultingDeckSection(section_key="process_flow", title="계산 / 전표 / GL 흐름", items=[f"외화 입출금 FIFO 구조를 계산, 전표, GL 기준으로 정리하는 설계 흐름 {index}입니다." for index in range(1, 5)]),
                ],
            ),
            ConsultingDeckChapter(
                chapter_key="vision",
                title="컨설팅 비전",
                sections=[ConsultingDeckSection(section_key="actions", title="적용 방향", items=[f"입금 lot 적재와 출금 lot 소진 계산을 분리 체계로 전환하는 미래 방향 {index}입니다." for index in range(1, 8)])],
            ),
        ],
    )
    slide_schema = build_slide_schema(oversized)
    service = create_document_service(output_dir=tmp_path)

    async def _render():
        await service.initialize()
        try:
            return await service.generate(
                DocumentRequest(
                    content="",
                    output_type=DocumentType.PPTX,
                    title="과밀 테스트",
                    filename="long_slide_schema",
                    payload=slide_schema.model_dump(),
                    style_options={"renderer": "slide_schema"},
                )
            )
        finally:
            await service.shutdown()

    result = asyncio.run(_render())
    prs = Presentation(str(result.output_path))
    titles = _first_shape_texts(prs)

    assert "컨설팅 전개 (2)" in titles
    assert "컨설팅 구현 (2)" in titles
    assert "컨설팅 설계 (2)" in titles
    assert "컨설팅 비전 (2)" in titles
    assert "보조 연속 장표" in _all_slide_text(prs, titles.index("컨설팅 구현 (2)"))
    assert "독립 연속 장표" in _all_slide_text(prs, titles.index("컨설팅 설계 (2)"))
