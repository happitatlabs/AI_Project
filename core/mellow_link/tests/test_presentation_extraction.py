from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from mellow_link.modules.rebuild_assistant.service import RebuildAssistantService
from mellow_link.services.anonymization import AnonymizationAsset, AnonymizationRunRequest, AnonymizationService
from mellow_link.services.anonymization.schemas import BundleAssetSummary, CanonicalAnonymizedSource, MaskingLevel, SafeAnalysisBundle
from mellow_link.services.presentation_extraction import extract_presentation_sml
from mellow_link.services.rag_service import extract_text_from_file
from mellow_link.services.refactoring_support_engine.analysis_context_builder import AnalysisContextBuilder


def _build_sample_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Apollo"
    slide.placeholders[1].text = "API /finance/payments\nOwner Jane"
    notes = slide.notes_slide.notes_text_frame
    notes.text = "Speaker note confidential"

    table_shape = slide.shapes.add_table(2, 2, Inches(0.8), Inches(3.2), Inches(5.5), Inches(1.2))
    table = table_shape.table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "ARR"
    table.cell(1, 1).text = "12"

    prs.save(path)


def test_extract_text_from_pptx_returns_slide_sml(tmp_path: Path):
    pptx_path = tmp_path / "sample_deck.pptx"
    _build_sample_pptx(pptx_path)

    extracted = extract_text_from_file(pptx_path)

    assert "[SML v1]" in extracted
    assert "presentation_file: sample_deck.pptx" in extracted
    assert "[SLIDE 1]" in extracted
    assert "title: Project Apollo" in extracted
    assert "- API /finance/payments" in extracted
    assert "tables:" in extracted
    assert "| Metric | Value |" in extracted
    assert "notes:" in extracted
    assert "- Speaker note confidential" in extracted


def test_extract_text_from_legacy_ppt_uses_conversion_fallback(monkeypatch):
    def fake_extract_legacy(file_path: Path, content_bytes: bytes | None = None) -> str:
        assert file_path.suffix.lower() == ".ppt"
        assert content_bytes == b"legacy-binary"
        return "[SML v1]\npresentation_file: legacy.ppt\nslide_count: 1"

    monkeypatch.setattr(
        "mellow_link.services.presentation_extraction._extract_legacy_ppt_sml",
        fake_extract_legacy,
    )

    extracted = extract_text_from_file(Path("legacy.ppt"), b"legacy-binary")

    assert "presentation_file: legacy.ppt" in extracted


def test_extract_text_from_legacy_ppt_uses_structured_storage_fallback(monkeypatch):
    sample_ppt = (
        Path(__file__).resolve().parents[1]
        / "modules"
        / "rebuild_assistant"
        / "samples"
        / "08_consulting_output_reference"
        / "assets"
        / "0_원가계산컨설팅자료_260418"
        / "1부산우유컨설팅개요.ppt"
    )
    if not sample_ppt.exists():
        pytest.skip("sample legacy ppt not available")

    monkeypatch.setattr(
        "mellow_link.services.presentation_extraction._convert_ppt_to_pptx_with_powerpoint",
        lambda *_args, **_kwargs: False,
    )

    extracted = extract_text_from_file(sample_ppt)

    assert "[SML v1]" in extracted
    assert "presentation_file: 1부산우유컨설팅개요.ppt" in extracted
    assert "extraction_mode: legacy_binary_stream_fallback" in extracted
    assert "[SLIDE 1]" in extracted
    assert "Part I. 컨설팅 개요" in extracted
    assert "컨설팅의 배경과 필요성" in extracted


def test_extract_text_from_legacy_ppt_uses_binary_scan_fallback(monkeypatch):
    monkeypatch.setattr(
        "mellow_link.services.presentation_extraction._convert_ppt_to_pptx_with_powerpoint",
        lambda *_args, **_kwargs: False,
    )
    monkeypatch.setattr(
        "mellow_link.services.presentation_extraction._extract_legacy_ppt_sml_from_structured_storage",
        lambda *_args, **_kwargs: "",
    )

    legacy_bytes = (
        b"\x00\x01"
        + "Legacy fallback title".encode("utf-16le")
        + b"\x00\x00"
        + "Project Apollo".encode("utf-16le")
        + b"\x00\x00"
    )

    extracted = extract_text_from_file(Path("legacy_scan.ppt"), legacy_bytes)

    assert "[SML v1]" in extracted
    assert "presentation_file: legacy_scan.ppt" in extracted
    assert "extraction_mode: legacy_binary_scan_fallback" in extracted
    assert "- Legacy fallback title" in extracted
    assert "- Project Apollo" in extracted


def test_pptx_assets_are_classified_as_docs_in_analysis_context():
    safe_bundle = SafeAnalysisBundle(
        bundle_id="bundle_001",
        project_id="proj_001",
        masking_level=MaskingLevel.FULL,
        asset_summary=[
            BundleAssetSummary(
                asset_id="asset_001",
                name="business_review.pptx",
                temp_file_id="temp_001",
                kind_hint="presentation",
            )
        ],
        sources=[
            CanonicalAnonymizedSource(
                asset_id="asset_001",
                content="[SML v1]\npresentation_file: business_review.pptx\nslide_count: 1",
            )
        ],
    )

    bundle = AnalysisContextBuilder().build(project_id="proj_001", run_id="run_001", safe_bundle=safe_bundle)

    assert bundle.assets[0].asset_type == "doc"
    assert bundle.source_blocks[0].asset_type == "doc"


def test_presentation_sml_survives_anonymization_pipeline():
    raw_sml = "\n".join(
        [
            "[SML v1]",
            "presentation_file: client_review.pptx",
            "slide_count: 1",
            "",
            "[SLIDE 1]",
            "title: Project Apollo",
            "texts:",
            "- API /finance/payments",
            "- class RevenueBridge",
        ]
    )
    result = AnonymizationService().run_anonymization_pipeline(
        AnonymizationRunRequest(
            project_id="proj_001",
            assets=[
                AnonymizationAsset(
                    asset_id="asset_001",
                    name="client_review.pptx",
                    temp_file_id="temp_001",
                    kind_hint="presentation",
                    content_text=raw_sml,
                    original_bytes=b"pptx",
                )
            ],
        )
    )

    canonical = result.safe_bundle.sources[0].content

    assert "[SML v1]" in canonical
    assert "[SLIDE 1]" in canonical
    assert "API_001" in canonical
    assert "CLS_001" in canonical


def test_rebuild_service_treats_pptx_as_doc_asset():
    service = RebuildAssistantService()

    assert service._is_doc_asset_name("deck.pptx") is True
    assert service._is_doc_asset_name("deck.ppt") is True
