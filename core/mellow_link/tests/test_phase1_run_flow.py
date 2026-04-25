import json
import shutil
import uuid
from datetime import datetime, timezone
from pathlib import Path

import pytest

MELLOW_LINK_ROOT = Path(__file__).resolve().parents[1]

_SYNTHETIC_UPLOAD_SESSION_PREFIXES = (
    "phase1-",
    "upload_large_structured",
    "upload_customer_intent",
)

_SYNTHETIC_UPLOAD_SESSION_EXACT = {
    "fallback-session",
    "history-session",
    "anonymization-session",
}

try:
    from fastapi.testclient import TestClient
    _has_fastapi = True
except ImportError:
    _has_fastapi = False

_app = None
_skip_reason = None


def _get_app():
    global _app, _skip_reason
    if _app is not None:
        return _app
    if _skip_reason is not None:
        return None
    if not _has_fastapi:
        _skip_reason = "FastAPI/testclient not installed"
        return None
    try:
        from mellow_link.main import app
        _app = app
        return _app
    except Exception as e:
        _skip_reason = f"Could not import app: {e}"
        return None


@pytest.fixture(scope="module")
def client():
    app = _get_app()
    if app is None:
        pytest.skip(_skip_reason or "FastAPI app not available")
    return TestClient(app)


def _synthetic_upload_session_filters(column):
    filters = [column.like(f"{prefix}%") for prefix in _SYNTHETIC_UPLOAD_SESSION_PREFIXES]
    filters.append(column.in_(sorted(_SYNTHETIC_UPLOAD_SESSION_EXACT)))
    return filters


def _purge_synthetic_phase1_rows():
    from sqlalchemy import or_

    from mellow_link.infra import (
        AgentRun,
        AgentRunEvent,
        ModernizationProject,
        ProjectAsset,
        ProjectRunHistory,
        TempResource,
    )
    from mellow_link.infra.database import SessionLocal

    with SessionLocal() as db:
        projects = db.query(ModernizationProject.id, ModernizationProject.run_id).filter(
            or_(*_synthetic_upload_session_filters(ModernizationProject.upload_session_id))
        ).all()
        project_ids = [row.id for row in projects]
        run_ids = [row.run_id for row in projects if row.run_id]

        if project_ids:
            db.query(ProjectRunHistory).filter(
                ProjectRunHistory.project_id.in_(project_ids)
            ).delete(synchronize_session=False)
            db.query(ProjectAsset).filter(ProjectAsset.project_id.in_(project_ids)).delete(
                synchronize_session=False
            )
            db.query(ModernizationProject).filter(
                ModernizationProject.id.in_(project_ids)
            ).delete(synchronize_session=False)

        if run_ids:
            db.query(AgentRunEvent).filter(AgentRunEvent.run_id.in_(run_ids)).delete(
                synchronize_session=False
            )
            db.query(AgentRun).filter(AgentRun.run_id.in_(run_ids)).delete(
                synchronize_session=False
            )

        temp_resource_filters = list(_synthetic_upload_session_filters(TempResource.temp_session_id))
        if project_ids:
            temp_resource_filters.append(TempResource.promoted_to_project_id.in_(project_ids))
        db.query(TempResource).filter(or_(*temp_resource_filters)).delete(
            synchronize_session=False
        )
        db.commit()


def _build_isolated_session_factory(tmp_path: Path):
    from sqlalchemy import create_engine
    from sqlalchemy.orm import sessionmaker

    from mellow_link.infra.database import Base

    db_path = (tmp_path / f"phase1_isolated_{uuid.uuid4().hex}.db").resolve()
    engine = create_engine(
        f"sqlite:///{db_path.as_posix()}",
        connect_args={"check_same_thread": False},
    )
    Base.metadata.create_all(bind=engine)
    return sessionmaker(autocommit=False, autoflush=False, bind=engine), engine


def _seed_isolated_user(db, *, username_prefix: str):
    from mellow_link.infra import User, UserRole

    user = User(
        username=f"{username_prefix}_{uuid.uuid4().hex[:8]}",
        hashed_password="test-hash",
        role=UserRole.USER.value,
    )
    db.add(user)
    db.commit()
    db.refresh(user)
    return user


@pytest.fixture(autouse=True)
def cleanup_synthetic_phase1_rows():
    _purge_synthetic_phase1_rows()
    yield
    _purge_synthetic_phase1_rows()


def _register(client, username_prefix="phase1"):
    from mellow_link.infra.database import SessionLocal
    from mellow_link.infra import User, UserRole, create_default_folders_for_user, create_access_token

    username = f"{username_prefix}_{uuid.uuid4().hex[:8]}"
    with SessionLocal() as db:
        user = User(
            username=username,
            hashed_password="test-hash",
            role=UserRole.USER.value,
        )
        db.add(user)
        db.commit()
        db.refresh(user)
        create_default_folders_for_user(db, user.id, role=UserRole.USER.value)
        token = create_access_token(data={"sub": username}, role=user.role)
    return {
        "username": username,
        "token": token,
        "headers": {"Authorization": f"Bearer {token}", "Accept": "application/json"},
    }


def _register_admin(client, username_prefix="phase1_admin"):
    from mellow_link.infra.database import SessionLocal
    from mellow_link.infra import User, UserRole, create_default_folders_for_user, create_access_token

    username = f"{username_prefix}_{uuid.uuid4().hex[:8]}"
    with SessionLocal() as db:
        user = User(
            username=username,
            hashed_password="test-hash",
            role=UserRole.ADMIN.value,
        )
        db.add(user)
        db.commit()
        db.refresh(user)
        create_default_folders_for_user(db, user.id, role=UserRole.ADMIN.value)
        token = create_access_token(data={"sub": username}, role=user.role)
    return {
        "username": username,
        "token": token,
        "headers": {"Authorization": f"Bearer {token}", "Accept": "application/json"},
    }


def _emit_finished(run_id: str, success: bool = True, summary: str = "done"):
    from mellow_link.infra.run_events import emit_event, EVENT_TYPE_RUN_STARTED, EVENT_TYPE_RUN_FINISHED

    emit_event(run_id, EVENT_TYPE_RUN_STARTED, {"user_input": "phase1 test", "mode": "fast", "session_id": None})
    emit_event(run_id, EVENT_TYPE_RUN_FINISHED, {"success": success, "summary": summary})


def _upload_temp_asset(
    client,
    session_id: str,
    filename: str,
    content: bytes,
    headers: dict | None = None,
    content_type: str = "text/plain",
) -> dict:
    res = client.post(
        "/chat/upload-temp",
        headers=headers or {},
        data={"session_id": session_id},
        files={"file": (filename, content, content_type)},
    )
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["temp_file_id"]
    return body


def _build_review_pptx_bytes(tmp_path: Path) -> bytes:
    from pptx import Presentation

    pptx_path = tmp_path / f"review_{uuid.uuid4().hex[:8]}.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "차세대 ERP 구축 사업"
    slide.placeholders[1].text = "\n".join(
        [
            "고객사: 한국전력공사",
            "기관명: 산업통상자원부",
            "담당자: 홍길동 PM",
            "프로젝트명: Project Apollo Modernization",
            "이메일: hgd@example.com",
            "한국지역난방공사와 공동 추진",
        ]
    )
    notes = slide.notes_slide.notes_text_frame
    notes.text = "계약명: 2026 ERP 통합 고도화 계약"
    prs.save(pptx_path)
    return pptx_path.read_bytes()


def _create_persisted_project(
    client,
    user: dict,
    monkeypatch,
    *,
    upload_session_id: str,
    filename: str,
    content: bytes,
    project_name: str,
    client_name: str,
    goal: str = "",
) -> tuple[str, dict]:
    from mellow_link.routers import projects as projects_router

    monkeypatch.setattr(projects_router, "start_project_wrapped_run", lambda *args, **kwargs: None)
    upload = _upload_temp_asset(
        client,
        session_id=upload_session_id,
        filename=filename,
        content=content,
        headers=user["headers"],
    )
    create_res = client.post(
        "/projects",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "project_name": project_name,
            "client_name": client_name,
            "goal": goal,
            "upload_session_id": upload_session_id,
            "asset_manifest": [{"name": filename, "temp_file_id": upload["temp_file_id"], "size": len(content)}],
            "template_key": "default_modernization_v1",
            "constraints": [],
        },
    )
    assert create_res.status_code == 200, create_res.text
    return create_res.json()["project_id"], upload


def _create_fallback_project(user: dict, *, project_name: str = "Fallback 프로젝트", status: str = "completed") -> str:
    from mellow_link.infra import ModernizationProject, User
    from mellow_link.infra.database import SessionLocal
    from mellow_link.infra.run_events import create_run, emit_event, EVENT_TYPE_RUN_STARTED, EVENT_TYPE_RUN_FINISHED
    from mellow_link.routers.runs import _resolve_run_session_id

    with SessionLocal() as db:
        db_user = db.query(User).filter(User.username == user["username"]).first()
        session_id = _resolve_run_session_id(db, db_user, None)
        run_id = create_run(session_id=session_id, db=db, module_id="rebuild_assistant", run_kind="rebuild_plan")
        if status in ("completed", "failed"):
            emit_event(run_id, EVENT_TYPE_RUN_STARTED, {"user_input": "fallback project", "mode": "fast", "session_id": session_id}, db=db)
            emit_event(run_id, EVENT_TYPE_RUN_FINISHED, {"success": status == "completed", "summary": status}, db=db)
        project = ModernizationProject(
            id=f"proj_{uuid.uuid4().hex[:12]}",
            user_id=db_user.id,
            session_id=session_id,
            run_id=run_id,
            project_name=project_name,
            client_name="OO캐피탈",
            template_key="default_modernization_v1",
            template_mode="recommended",
            constraints_json="[]",
            upload_session_id="fallback-session",
            asset_manifest_json='[{"name":"legacy.jsp","temp_file_id":"legacy-jsp","size":777}]',
            status=status,
        )
        db.add(project)
        db.commit()
        return project.id


def _build_large_rebuild_result_with_polish_bundle() -> tuple[dict, dict]:
    from mellow_link.modules.rebuild_assistant.schemas import RebuildAssetsPayload
    from mellow_link.modules.rebuild_assistant.service import RebuildAssistantService

    service = RebuildAssistantService()
    prepared = service.prepare_input(
        goal="주문 조회와 저장 검증 흐름을 React + REST API 기준으로 재구성해줘",
        assets=RebuildAssetsPayload(
            source_code="""
class OrderService:
    def submit(self, order, repo):
        if not order.amount:
            raise ValueError("required")
        if order.status == "READY":
            repo.save(order)
            return approve(order)

class OrderQueryService:
    def search(self, keyword, status, page, repo):
        return repo.find(keyword=keyword, status=status, page=page)
            """,
            database_schema="""
CREATE TABLE orders (
    id bigint primary key,
    status varchar(20),
    amount numeric(18, 2),
    created_at timestamp
);
CREATE INDEX idx_orders_status ON orders(status);
            """,
            sql_queries="""
SELECT id, status, amount, created_at
FROM orders
WHERE status = :status
  AND (:keyword IS NULL OR id::text LIKE :keyword)
ORDER BY created_at DESC
LIMIT :limit OFFSET :offset;
            """,
            ui_template="""
<form>
  <input name="keyword" />
  <select name="status"></select>
  <button onclick="submitOrder()">submit</button>
</form>
            """,
            framework_info="JSP + Spring MVC + MyBatis",
        ),
        constraints=["기존 DB 계약 유지", "승인 상태 규칙 유지"],
    )
    result = service.build_result(prepared)
    polish_bundle = service.build_polish_bundle(result, audience="manager", delivery_mode="client_report").model_dump()
    payload = {"structured_result": result.model_dump(), "polish_bundle": polish_bundle}
    assert len(json.dumps(payload, ensure_ascii=False)) > 8192
    return result.model_dump(), polish_bundle


def _project_run_history_rows(project_id: str):
    from mellow_link.infra import ProjectRunHistory
    from mellow_link.infra.database import SessionLocal

    with SessionLocal() as db:
        return (
            db.query(ProjectRunHistory)
            .filter(ProjectRunHistory.project_id == project_id)
            .order_by(ProjectRunHistory.sequence_no.asc())
            .all()
        )


def _project_run_id(project_id: str) -> str:
    from mellow_link.infra import ModernizationProject
    from mellow_link.infra.database import SessionLocal

    with SessionLocal() as db:
        project = db.query(ModernizationProject).filter(ModernizationProject.id == project_id).first()
        assert project is not None
        return project.run_id


def _project_result_archive_dir(project_id: str, run_id: str) -> Path:
    return MELLOW_LINK_ROOT.parents[1] / "data" / "outputs" / "final" / "project_results" / project_id / run_id


def _create_project_with_history(
    user: dict,
    *,
    project_name: str = "히스토리 프로젝트",
    latest_status: str = "running",
    historical_manifest: list[dict] | None = None,
    latest_manifest: list[dict] | None = None,
    upload_session_id: str = "history-session",
) -> tuple[str, str, str]:
    from mellow_link.infra import ModernizationProject, ProjectRunHistory, User
    from mellow_link.infra.database import SessionLocal
    from mellow_link.infra.run_events import create_run, emit_event, EVENT_TYPE_RUN_STARTED, EVENT_TYPE_RUN_FINISHED
    from mellow_link.routers.runs import _resolve_run_session_id

    historical_manifest = historical_manifest or [
        {"name": "legacy.jsp", "temp_file_id": "legacy-jsp", "size": 777}
    ]
    latest_manifest = latest_manifest or [
        {"name": "legacy.jsp", "temp_file_id": "legacy-jsp", "size": 777},
        {"name": "schema.sql", "temp_file_id": "schema-sql", "size": 555},
    ]
    structured_result, polish_bundle = _build_large_rebuild_result_with_polish_bundle()

    with SessionLocal() as db:
        db_user = db.query(User).filter(User.username == user["username"]).first()
        session_id = _resolve_run_session_id(db, db_user, None)
        historical_run_id = create_run(session_id=session_id, db=db, module_id="rebuild_assistant", run_kind="rebuild_plan")
        latest_run_id = create_run(session_id=session_id, db=db, module_id="rebuild_assistant", run_kind="rebuild_plan")
        emit_event(historical_run_id, EVENT_TYPE_RUN_STARTED, {"user_input": "history", "mode": "fast", "session_id": session_id}, db=db)
        emit_event(
            historical_run_id,
            EVENT_TYPE_RUN_FINISHED,
            {
                "success": True,
                "summary": "history",
                "structured_result": structured_result,
                "polish_bundle": polish_bundle,
                "primary_feature_mode": "save_validation",
                "module_id": "rebuild_assistant",
                "run_kind": "rebuild_plan",
            },
            db=db,
        )
        emit_event(latest_run_id, EVENT_TYPE_RUN_STARTED, {"user_input": "latest", "mode": "fast", "session_id": session_id}, db=db)
        if latest_status in ("completed", "failed"):
            emit_event(
                latest_run_id,
                EVENT_TYPE_RUN_FINISHED,
                {
                    "success": latest_status == "completed",
                    "summary": latest_status,
                },
                db=db,
            )
        project = ModernizationProject(
            id=f"proj_{uuid.uuid4().hex[:12]}",
            user_id=db_user.id,
            session_id=session_id,
            run_id=latest_run_id,
            project_name=project_name,
            client_name="OO카드",
            template_key="default_modernization_v1",
            template_mode="recommended",
            constraints_json="[]",
            upload_session_id=upload_session_id,
            asset_manifest_json=json.dumps(latest_manifest, ensure_ascii=False),
            status=latest_status,
        )
        db.add(project)
        db.flush()
        db.add_all(
            [
                ProjectRunHistory(
                    id=f"prh_{uuid.uuid4().hex[:12]}",
                    project_id=project.id,
                    run_id=historical_run_id,
                    sequence_no=1,
                    trigger_kind="initial",
                    asset_manifest_json=json.dumps(historical_manifest, ensure_ascii=False),
                ),
                ProjectRunHistory(
                    id=f"prh_{uuid.uuid4().hex[:12]}",
                    project_id=project.id,
                    run_id=latest_run_id,
                    sequence_no=2,
                    trigger_kind="reanalysis",
                    asset_manifest_json=json.dumps(latest_manifest, ensure_ascii=False),
                ),
            ]
        )
        db.commit()
        return project.id, historical_run_id, latest_run_id



def test_root_redirects_to_project_create(client):
    res = client.get("/", follow_redirects=False)
    assert res.status_code in (302, 307)
    assert res.headers["location"] == "/projects/create"


def test_legacy_ui_is_marked_deprecated(client):
    res = client.get("/index.html")
    assert res.status_code == 200
    assert "Deprecated" in res.text
    assert "/runtime-console" in res.text


def test_home_prioritizes_project_entry(client):
    res = client.get("/ui")
    assert res.status_code == 200
    text = res.text
    assert "레거시 현대화 분석" in text
    assert "새 프로젝트" in text


def test_legacy_ui_disables_new_chat_flow(client):
    res = client.get("/index.html")
    assert res.status_code == 200
    text = res.text
    assert "Legacy UI (Deprecated)" in text
    assert "Runtime Console 사용을 권장합니다." in text
    assert 'id="messageInput"' in text and 'disabled' in text
    assert 'id="sendBtn"' in text and 'Legacy chat disabled' in text


def test_user_console_run_redirects_completed_project_to_result_package(client):
    user = _register(client, "phase1_user_console_completed")
    project_id = _create_fallback_project(user, project_name="완료 리디렉션", status="completed")
    run_id = _project_run_id(project_id)

    res = client.get(f"/user-console?run_id={run_id}", follow_redirects=False)

    assert res.status_code in (302, 307)
    assert res.headers["location"] == f"/projects/{project_id}/result"


def test_user_console_run_redirects_running_project_to_workspace(client):
    user = _register(client, "phase1_user_console_running")
    project_id = _create_fallback_project(user, project_name="진행중 리디렉션", status="running")
    run_id = _project_run_id(project_id)

    res = client.get(f"/user-console?run_id={run_id}", follow_redirects=False)

    assert res.status_code in (302, 307)
    assert res.headers["location"] == f"/projects/{project_id}"


def test_user_console_run_redirects_historical_project_run_to_workspace(client):
    from mellow_link.infra import ModernizationProject, ProjectRunHistory, User
    from mellow_link.infra.database import SessionLocal
    from mellow_link.infra.run_events import create_run, emit_event, EVENT_TYPE_RUN_STARTED, EVENT_TYPE_RUN_FINISHED
    from mellow_link.routers.runs import _resolve_run_session_id

    user = _register(client, "phase1_user_console_history")

    with SessionLocal() as db:
        db_user = db.query(User).filter(User.username == user["username"]).first()
        session_id = _resolve_run_session_id(db, db_user, None)
        historical_run_id = create_run(session_id=session_id, db=db, module_id="rebuild_assistant", run_kind="rebuild_plan")
        latest_run_id = create_run(session_id=session_id, db=db, module_id="rebuild_assistant", run_kind="rebuild_plan")
        emit_event(historical_run_id, EVENT_TYPE_RUN_STARTED, {"user_input": "history", "mode": "fast", "session_id": session_id}, db=db)
        emit_event(historical_run_id, EVENT_TYPE_RUN_FINISHED, {"success": True, "summary": "history"}, db=db)
        emit_event(latest_run_id, EVENT_TYPE_RUN_STARTED, {"user_input": "latest", "mode": "fast", "session_id": session_id}, db=db)
        emit_event(latest_run_id, EVENT_TYPE_RUN_FINISHED, {"success": True, "summary": "latest"}, db=db)
        project = ModernizationProject(
            id=f"proj_{uuid.uuid4().hex[:12]}",
            user_id=db_user.id,
            session_id=session_id,
            run_id=latest_run_id,
            project_name="히스토리 리디렉션",
            client_name="OO카드",
            template_key="default_modernization_v1",
            template_mode="recommended",
            constraints_json="[]",
            upload_session_id="history-session",
            asset_manifest_json="[]",
            status="completed",
        )
        db.add(project)
        db.flush()
        db.add_all(
            [
                ProjectRunHistory(
                    id=f"prh_{uuid.uuid4().hex[:12]}",
                    project_id=project.id,
                    run_id=historical_run_id,
                    sequence_no=1,
                    trigger_kind="initial",
                    asset_manifest_json="[]",
                ),
                ProjectRunHistory(
                    id=f"prh_{uuid.uuid4().hex[:12]}",
                    project_id=project.id,
                    run_id=latest_run_id,
                    sequence_no=2,
                    trigger_kind="reanalysis",
                    asset_manifest_json="[]",
                ),
            ]
        )
        db.commit()
        project_id = project.id

    res = client.get(f"/user-console?run_id={historical_run_id}", follow_redirects=False)

    assert res.status_code in (302, 307)
    assert res.headers["location"] == f"/projects/{project_id}"


def test_projects_list_exposes_recent_entries_with_run_specific_result_urls(client):
    user = _register(client, "phase1_recent_entries")
    project_id, historical_run_id, latest_run_id = _create_project_with_history(
        user,
        project_name="최근 목록 실행 분리",
        latest_status="running",
        upload_session_id="real-history-list",
    )

    res = client.get("/projects?format=json", headers=user["headers"])

    assert res.status_code == 200, res.text
    payload = res.json()
    entries = [item for item in payload["recent_entries"] if item["project_id"] == project_id]
    assert [item["run_id"] for item in entries] == [latest_run_id, historical_run_id]
    assert entries[0]["status"] == "running"
    assert entries[0]["result_url"] is None
    assert entries[1]["status"] == "completed"
    assert entries[1]["result_url"] == f"/projects/{project_id}/result?run_id={historical_run_id}"


def test_projects_list_hides_synthetic_test_rows(client, monkeypatch):
    user = _register(client, "phase1_hide_synthetic")
    visible_project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="real-user-list",
        filename="legacy.sql",
        content=b"select 1;",
        project_name="실제 사용자 프로젝트",
        client_name="실사용고객",
    )
    hidden_project_id = _create_fallback_project(user, project_name="숨김 synthetic", status="completed")

    res = client.get("/projects?format=json", headers=user["headers"])

    assert res.status_code == 200, res.text
    payload = res.json()
    project_ids = [item["id"] for item in payload["projects"]]
    recent_project_ids = [item["project_id"] for item in payload["recent_entries"]]
    assert visible_project_id in project_ids
    assert visible_project_id in recent_project_ids
    assert hidden_project_id not in project_ids
    assert hidden_project_id not in recent_project_ids


def test_projects_list_shows_synthetic_test_rows_for_admin(client):
    admin = _register_admin(client, "phase1_show_synthetic_admin")
    hidden_project_id = _create_fallback_project(admin, project_name="관리자 synthetic", status="completed")

    res = client.get("/projects?format=json", headers=admin["headers"])

    assert res.status_code == 200, res.text
    payload = res.json()
    project_ids = [item["id"] for item in payload["projects"]]
    recent_project_ids = [item["project_id"] for item in payload["recent_entries"]]
    assert hidden_project_id in project_ids
    assert hidden_project_id in recent_project_ids


def test_projects_list_marks_unstarted_synthetic_rows_as_pending_for_admin(client):
    from mellow_link.infra import ModernizationProject, User
    from mellow_link.infra.database import SessionLocal
    from mellow_link.infra.run_events import create_run

    admin = _register_admin(client, "phase1_pending_admin")

    with SessionLocal() as db:
        db_user = db.query(User).filter(User.username == admin["username"]).first()
        assert db_user is not None
        run_id = create_run(session_id="sess_customer_intent_pending", db=db, module_id="rebuild_assistant", run_kind="rebuild_plan")
        project = ModernizationProject(
            id=f"proj_customer_intent_pending_{uuid.uuid4().hex[:8]}",
            user_id=db_user.id,
            session_id="sess_customer_intent_pending",
            run_id=run_id,
            project_name="고객 요청 노출 확인",
            client_name="OO",
            template_key="default_modernization_v1",
            template_mode="recommended",
            constraints_json="[]",
            upload_session_id=f"upload_customer_intent_{uuid.uuid4().hex[:8]}",
            asset_manifest_json="[]",
            status="running",
        )
        db.add(project)
        db.commit()
        project_id = project.id

    res = client.get("/projects?format=json", headers=admin["headers"])

    assert res.status_code == 200, res.text
    payload = res.json()
    entry = next(item for item in payload["recent_entries"] if item["project_id"] == project_id)
    assert entry["status"] == "pending"
    assert entry["status_bucket"] == "pending"


def test_project_result_supports_historical_run_id_and_historical_assets(client):
    user = _register(client, "phase1_historical_result")
    project_id, historical_run_id, _latest_run_id = _create_project_with_history(
        user,
        project_name="히스토리 결과 조회",
        latest_status="running",
        historical_manifest=[
            {"name": "legacy.jsp", "temp_file_id": "legacy-jsp", "size": 777},
        ],
        latest_manifest=[
            {"name": "legacy.jsp", "temp_file_id": "legacy-jsp", "size": 777},
            {"name": "schema.sql", "temp_file_id": "schema-sql", "size": 555},
        ],
        upload_session_id="real-history-detail",
    )

    historical = client.get(
        f"/projects/{project_id}/result",
        params={"format": "json", "run_id": historical_run_id},
        headers=user["headers"],
    )
    assert historical.status_code == 200, historical.text
    historical_body = historical.json()
    assert historical_body["provenance"]["run_id"] == historical_run_id
    assert [item["temp_file_id"] for item in historical_body["provenance"]["input_assets"]] == ["legacy-jsp"]

    explanation = client.get(
        f"/projects/{project_id}/result/explanation",
        params={"run_id": historical_run_id, "audience": "manager", "surface_mode": "external"},
        headers=user["headers"],
    )
    assert explanation.status_code == 200, explanation.text

    detail = client.get(f"/projects/{project_id}?format=json", headers=user["headers"])
    assert detail.status_code == 200, detail.text
    assert detail.json()["project"]["status"] == "running"


def test_project_detail_bootstraps_initial_history_for_legacy_project(client):
    user = _register(client, "phase1_detail_history_bootstrap")
    project_id = _create_fallback_project(user, project_name="레거시 상세 부트스트랩", status="completed")

    res = client.get(f"/projects/{project_id}?format=json", headers=user["headers"])

    assert res.status_code == 200, res.text
    history = res.json()["run_history"]
    assert len(history) == 1
    assert history[0]["sequence_no"] == 1
    assert history[0]["trigger_kind"] == "initial"
    assert history[0]["is_latest"] is True


def test_run_surfaces_include_project_redirect_context(client):
    user = _register(client, "phase1_run_surface_project")
    admin = _register_admin(client, "phase1_run_surface_project_admin")
    project_id = _create_fallback_project(user, project_name="표면 컨텍스트", status="completed")
    run_id = _project_run_id(project_id)

    owned_list = client.get("/runs", headers=user["headers"])
    assert owned_list.status_code == 200, owned_list.text
    run_row = next(item for item in owned_list.json()["runs"] if item["run_id"] == run_id)
    assert run_row["project_id"] == project_id
    assert run_row["preferred_user_url"] == f"/projects/{project_id}/result"

    snapshot = client.get(f"/runs/{run_id}", headers=user["headers"])
    assert snapshot.status_code == 200, snapshot.text
    assert snapshot.json()["project_id"] == project_id
    assert snapshot.json()["preferred_user_url"] == f"/projects/{project_id}/result"

    dev_snapshot = client.get(f"/api/dev/runs/{run_id}", headers=admin["headers"])
    assert dev_snapshot.status_code == 200, dev_snapshot.text
    assert dev_snapshot.json()["project_id"] == project_id
    assert dev_snapshot.json()["preferred_user_url"] == f"/projects/{project_id}/result"

def test_ui_exposes_product_create_flow(client):
    res = client.get("/ui")
    assert res.status_code == 200
    text = res.text
    assert "분석 시작" in text
    assert "자산 업로드" in text
    assert "내 프로젝트" in text


def test_projects_create_page_supports_goal_and_constraints_autofill(client):
    res = client.get("/projects/create", headers={"Accept": "text/html"})
    assert res.status_code == 200
    text = res.text
    assert "applyAssetAutofill" in text
    assert "goal.txt" in text
    assert "constraints.txt" in text
    assert 'id="goal"' in text
    assert 'id="customerIntent"' in text
    assert "customer_request.txt" in text
    assert "customer_intent.txt" in text
    assert "preferred_direction.txt" in text
    assert "buildConstraintPayload" in text
    assert 'id="anonymizationReviewSection"' in text
    assert "입력 익명화 리뷰" in text
    assert "refreshAnonymizationReview" in text
    assert "await refreshAnonymizationReview();" in text
    assert "/projects/anonymization-review" in text
    assert "Safe Preview" in text
    assert "Role Token 요약" in text
    assert "감지된 원문 유형" in text
    assert "high_risk_detected" in text
    assert "label_less_risk_count" in text
    assert "display_review_report" in text
    assert "Source Question Guard" in text
    assert 'id="sourceQuestionCandidates"' in text
    assert 'id="guardedUserQuestions"' in text
    assert "/static/anonymization_display.js" in text
    assert "표시 안내: [PERSON], [ORG], [PROJECT], [TERM] 등은 익명 처리된 민감정보입니다." in text
    assert "분석 목표 자동 채움" in text
    assert "제약 조건 자동 병합" in text
    assert "고객 요청 자동 채움" in text


def test_projects_anonymization_review_preview_returns_safe_report(client):
    user = _register(client, "phase1_review")
    session_id = f"phase1-review-{uuid.uuid4().hex[:8]}"
    sml = "\n".join(
        [
            "[SML v1]",
            "presentation_file: review_deck.pptx",
            "slide_count: 2",
            "",
            "[SLIDE 1]",
            "title: 차세대 ERP 구축 사업",
            "texts:",
            "- 고객사: 한국전력공사",
            "- 기관명: 산업통상자원부",
            "- 담당자: 홍길동 PM",
            "- 프로젝트명: Project Apollo Modernization",
            "- 이메일: hgd@example.com",
            "",
            "[SLIDE 2]",
            "title: 정산 통합 사업",
            "texts:",
            "- 한국지역난방공사와 공동 추진",
            "- API /finance/payments",
        ]
    )
    upload = _upload_temp_asset(
        client,
        session_id=session_id,
        filename="review.txt",
        content=sml.encode("utf-8"),
        headers=user["headers"],
    )

    res = client.post(
        "/projects/anonymization-review",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "upload_session_id": session_id,
            "asset_manifest": [
                {
                    "name": "review.txt",
                    "temp_file_id": upload["temp_file_id"],
                    "size": len(sml.encode("utf-8")),
                }
            ],
        },
    )

    assert res.status_code == 200, res.text
    body = res.json()
    assert body["high_risk_detected"] is True
    assert body["label_less_risk_count"] >= 1
    assert body["label_less_warning_count"] >= 1
    assert body["question_guard_summary"]["guard_input_source_count"] >= 1
    assert body["question_guard_summary"]["guard_input_total_chars"] > 0
    assert body["question_guard_summary"]["source_question_candidate_count"] >= 1
    assert body["source_question_candidates"]

    report = body["review_report"]
    display_report = body["display_review_report"]
    assert report is not None
    assert display_report is not None
    assert report["status"] == "blocked"
    preview_text = report["asset_previews"][0]["preview_text"]
    display_preview_text = display_report["asset_previews"][0]["preview_text"]
    assert "[CLIENT]" in preview_text
    assert "[ORG]" in preview_text
    assert "[PERSON]" in preview_text
    assert "[PROJECT]" in preview_text
    assert "[EMAIL]" in preview_text
    assert "[RISK_ORG_CANDIDATE]" in preview_text
    assert "[WARNING_BUSINESS_CANDIDATE]" in preview_text
    assert "[ORG]" in display_preview_text
    assert "[BUSINESS]" in display_preview_text
    assert "[RISK_ORG_CANDIDATE]" not in display_preview_text
    assert "[WARNING_BUSINESS_CANDIDATE]" not in display_preview_text

    serialized = json.dumps(body, ensure_ascii=False)
    for raw_value in (
        "한국전력공사",
        "산업통상자원부",
        "홍길동",
        "Project Apollo Modernization",
        "hgd@example.com",
        "한국지역난방공사",
        "정산 통합 사업",
    ):
        assert raw_value not in serialized


def test_projects_anonymization_review_preview_returns_source_question_guard_summary(client):
    user = _register(client, "phase1_question_guard_preview")
    session_id = f"phase1-question-guard-{uuid.uuid4().hex[:8]}"
    sml = "\n".join(
        [
            "[SML v1]",
            "presentation_file: cost_consulting_deck.pptx",
            "slide_count: 1",
            "",
            "[SLIDE 1]",
            "title: [CLIENT] 원가 컨설팅 개요",
            "texts:",
            "- 현행 원가체계 분석",
            "- 원가분석 및 원가계산 개선 방향",
            "- 재료비, 노무비, 제조경비 배부기준 검토",
            "- 손익분석 확장 검토",
        ]
    )
    upload = _upload_temp_asset(
        client,
        session_id=session_id,
        filename="cost_guard.txt",
        content=sml.encode("utf-8"),
        headers=user["headers"],
    )

    res = client.post(
        "/projects/anonymization-review",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "upload_session_id": session_id,
            "goal": "제품 저장 전 검증 로직을 어떻게 강화할 것인가?",
            "constraints": ["SQL 파라미터 검증을 어떻게 설계할 것인가?"],
            "asset_manifest": [
                {
                    "name": "cost_guard.txt",
                    "temp_file_id": upload["temp_file_id"],
                    "size": len(sml.encode("utf-8")),
                }
            ],
        },
    )

    assert res.status_code == 200, res.text
    body = res.json()
    questions = [item["question"] for item in body["source_question_candidates"]]
    blocked = {(item["question"], item["blocked_reason"]) for item in body["blocked_user_questions"]}

    assert "현행 원가체계의 한계는 무엇인가?" in questions
    assert "문서가 제안하는 원가계산 개선 방향은 무엇인가?" in questions
    assert (
        "제품 저장 전 검증 로직을 어떻게 강화할 것인가?",
        "source_domain_mismatch",
    ) in blocked
    assert (
        "SQL 파라미터 검증을 어떻게 설계할 것인가?",
        "source_domain_mismatch",
    ) in blocked
    assert body["question_guard_summary"]["preferred_question_axis"] == "processing_flow"


def test_projects_anonymization_review_preview_supports_pptx_upload(client, tmp_path: Path):
    user = _register(client, "phase1_pptx_preview")
    session_id = f"phase1-pptx-preview-{uuid.uuid4().hex[:8]}"
    pptx_bytes = _build_review_pptx_bytes(tmp_path)

    upload = _upload_temp_asset(
        client,
        session_id=session_id,
        filename="review_deck.pptx",
        content=pptx_bytes,
        headers=user["headers"],
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    assert "[SML v1]" in upload["extracted_text"]
    assert "presentation_file: review_deck.pptx" in upload["extracted_text"]
    assert "[SLIDE 1]" in upload["extracted_text"]

    res = client.post(
        "/projects/anonymization-review",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "upload_session_id": session_id,
            "asset_manifest": [
                {
                    "name": "review_deck.pptx",
                    "temp_file_id": upload["temp_file_id"],
                    "size": len(pptx_bytes),
                }
            ],
        },
    )

    assert res.status_code == 200, res.text
    body = res.json()
    report = body["review_report"]
    display_report = body["display_review_report"]
    preview_text = report["asset_previews"][0]["preview_text"]
    display_preview_text = display_report["asset_previews"][0]["preview_text"]
    serialized = json.dumps(body, ensure_ascii=False)

    assert body["high_risk_detected"] is True
    assert body["label_less_risk_count"] >= 1
    assert isinstance(body["label_less_warning_count"], int)
    assert body["question_guard_summary"]["guard_input_source_count"] >= 1
    assert body["question_guard_summary"]["sml_text_length"] > 0
    assert body["question_guard_summary"]["source_question_candidate_count"] >= 1
    assert len(body["source_question_candidates"]) >= 1
    assert report["status"] == "blocked"
    assert report["llm_send_allowed"] is False
    assert "[CLIENT]" in preview_text
    assert "[ORG]" in preview_text
    assert "[PERSON]" in preview_text
    assert "[PROJECT]" in preview_text
    assert "[EMAIL]" in preview_text
    assert "[RISK_ORG_CANDIDATE]" in preview_text
    assert "[ORG]" in display_preview_text
    assert "[RISK_ORG_CANDIDATE]" not in display_preview_text

    for raw_value in (
        "한국전력공사",
        "산업통상자원부",
        "홍길동",
        "Project Apollo Modernization",
        "hgd@example.com",
        "한국지역난방공사",
        "차세대 ERP 구축 사업",
        "2026 ERP 통합 고도화 계약",
    ):
        assert raw_value not in preview_text
        assert raw_value not in serialized


def test_projects_anonymization_review_preview_returns_generic_fallback_questions_when_domain_specific_is_absent(client):
    user = _register(client, "phase1_generic_question_preview")
    session_id = f"phase1-generic-question-{uuid.uuid4().hex[:8]}"
    sml = "\n".join(
        [
            "[SML v1]",
            "presentation_file: generic_strategy_deck.pptx",
            "slide_count: 1",
            "",
            "[SLIDE 1]",
            "title: 개선 전략 검토",
            "texts:",
            "- 현행 프로세스의 한계",
            "- 개선 방향 검토",
            "- 추가 확인이 필요한 항목",
            "- 판단 기준 정리",
        ]
    )
    upload = _upload_temp_asset(
        client,
        session_id=session_id,
        filename="generic_strategy.txt",
        content=sml.encode("utf-8"),
        headers=user["headers"],
    )

    res = client.post(
        "/projects/anonymization-review",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "upload_session_id": session_id,
            "asset_manifest": [
                {
                    "name": "generic_strategy.txt",
                    "temp_file_id": upload["temp_file_id"],
                    "size": len(sml.encode("utf-8")),
                }
            ],
        },
    )

    assert res.status_code == 200, res.text
    body = res.json()
    questions = [item["question"] for item in body["source_question_candidates"]]
    summary = body["question_guard_summary"]

    assert "이 문서는 어떤 문제를 해결하려는가?" in questions
    assert "현행 구조의 한계는 무엇인가?" in questions
    assert summary["guard_input_source_count"] >= 1
    assert summary["guard_input_total_chars"] >= 120
    assert "no_domain_terms_detected" in summary["no_candidate_reasons"]


def test_project_create_keeps_high_risk_as_warning_and_passes_anonymized_sml_bundle(
    client,
    monkeypatch,
    tmp_path: Path,
):
    from mellow_link.routers import projects as projects_router

    user = _register(client, "phase1_pptx_create")
    session_id = f"phase1-pptx-create-{uuid.uuid4().hex[:8]}"
    pptx_bytes = _build_review_pptx_bytes(tmp_path)
    upload = _upload_temp_asset(
        client,
        session_id=session_id,
        filename="review_deck.pptx",
        content=pptx_bytes,
        headers=user["headers"],
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    preview_res = client.post(
        "/projects/anonymization-review",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "upload_session_id": session_id,
            "asset_manifest": [
                {
                    "name": "review_deck.pptx",
                    "temp_file_id": upload["temp_file_id"],
                    "size": len(pptx_bytes),
                }
            ],
        },
    )
    assert preview_res.status_code == 200, preview_res.text
    preview_body = preview_res.json()
    assert preview_body["high_risk_detected"] is True

    captured: dict[str, object] = {}

    def _capture_start(**kwargs):
        captured.update(kwargs)

    monkeypatch.setattr(projects_router, "start_project_wrapped_run", _capture_start)

    create_res = client.post(
        "/projects",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "project_name": "PPTX high risk create",
            "client_name": "OO금융",
            "goal": "PPT 기반 레거시 기능 분석",
            "upload_session_id": session_id,
            "asset_manifest": [
                {
                    "name": "review_deck.pptx",
                    "temp_file_id": upload["temp_file_id"],
                    "size": len(pptx_bytes),
                }
            ],
            "template_key": "default_modernization_v1",
            "constraints": [],
        },
    )

    assert create_res.status_code == 200, create_res.text
    assert create_res.json()["status"] == "pending"
    assert captured.get("safe_bundle") is not None

    safe_bundle = captured["safe_bundle"]
    canonical = "\n\n".join(source.content for source in safe_bundle.sources)
    stats = safe_bundle.sources[0].replacement_stats

    assert "[SML v1]" in canonical
    assert "[SLIDE 1]" in canonical
    assert "CLIENT_001" in canonical
    assert "ORG_001" in canonical
    assert "PERSON_001" in canonical
    assert "PROJECT_001" in canonical
    assert "EMAIL_001" in canonical
    assert "[RISK_ORG_CANDIDATE]" in canonical
    assert stats["organization_name"] == 1

    for raw_value in (
        "한국전력공사",
        "산업통상자원부",
        "홍길동",
        "Project Apollo Modernization",
        "hgd@example.com",
        "한국지역난방공사",
        "2026 ERP 통합 고도화 계약",
    ):
        assert raw_value not in canonical



def test_project_create_preserves_consulting_terms_while_masking_label_less_person_candidates(
    client,
    monkeypatch,
):
    from mellow_link.routers import projects as projects_router

    user = _register(client, "phase1_person_filter_create")
    session_id = f"phase1-person-filter-{uuid.uuid4().hex[:8]}"
    sml = "\n".join(
        [
            "[SML v1]",
            "presentation_file: person_filter_deck.pptx",
            "slide_count: 1",
            "",
            "[SLIDE 1]",
            "title: 분석 범위 검토",
            "texts:",
            "- 원가 분석 시스템 구축",
            "- 업무 흐름도 및 데이터 처리 구조",
            "- 담당 홍길동",
            "- 문의: 홍길동 / hgd@example.com",
            "- ABC컨설팅",
            "- 홍길동 분석 시스템 구축",
        ]
    )
    upload = _upload_temp_asset(
        client,
        session_id=session_id,
        filename="person_filter.txt",
        content=sml.encode("utf-8"),
        headers=user["headers"],
    )

    preview_res = client.post(
        "/projects/anonymization-review",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "upload_session_id": session_id,
            "asset_manifest": [
                {
                    "name": "person_filter.txt",
                    "temp_file_id": upload["temp_file_id"],
                    "size": len(sml.encode("utf-8")),
                }
            ],
        },
    )
    assert preview_res.status_code == 200, preview_res.text
    assert preview_res.json()["high_risk_detected"] is True

    captured: dict[str, object] = {}

    def _capture_start(**kwargs):
        captured.update(kwargs)

    monkeypatch.setattr(projects_router, "start_project_wrapped_run", _capture_start)

    create_res = client.post(
        "/projects",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "project_name": "Label-less person filter create",
            "client_name": "OO금융",
            "goal": "문서형 익명화 의미 보존 확인",
            "upload_session_id": session_id,
            "asset_manifest": [
                {
                    "name": "person_filter.txt",
                    "temp_file_id": upload["temp_file_id"],
                    "size": len(sml.encode("utf-8")),
                }
            ],
            "template_key": "default_modernization_v1",
            "constraints": [],
        },
    )

    assert create_res.status_code == 200, create_res.text
    assert create_res.json()["status"] == "pending"

    safe_bundle = captured["safe_bundle"]
    canonical = "\n\n".join(source.content for source in safe_bundle.sources)
    stats = safe_bundle.sources[0].replacement_stats

    assert "원가 분석 시스템 구축" in canonical
    assert "업무 흐름도 및 데이터 처리 구조" in canonical
    assert "홍길동 분석 시스템 구축" in canonical
    assert "[LOW_CONF_TERM_CANDIDATE]" not in canonical
    assert "담당 [RISK_PERSON_CANDIDATE]" in canonical
    assert "문의: [RISK_PERSON_CANDIDATE] / EMAIL_001" in canonical
    assert "[RISK_ORG_CANDIDATE]" in canonical
    assert "ABC컨설팅" not in canonical
    assert stats.get("person_name", 0) == 0
    assert stats.get("organization_name", 0) == 0
    assert stats["email"] == 1


def test_project_create_masks_contextual_org_alias_in_preview_and_safe_bundle(client, monkeypatch):
    from mellow_link.routers import projects as projects_router

    user = _register(client, "phase1_org_alias_create")
    session_id = f"phase1-org-alias-{uuid.uuid4().hex[:8]}"
    sml = "\n".join(
        [
            "[SML v1]",
            "presentation_file: busan_milk_deck.pptx",
            "slide_count: 1",
            "",
            "[SLIDE 1]",
            "title: 부산우유 컨설팅 개요",
            "texts:",
            "- 고객사: 부산우유 주식회사",
            "- 부산우유의 비전",
            "- 부산우유 기초설계서",
            "- 부산우유 원가 시스템 개선",
        ]
    )
    upload = _upload_temp_asset(
        client,
        session_id=session_id,
        filename="busan_milk.txt",
        content=sml.encode("utf-8"),
        headers=user["headers"],
    )

    preview_res = client.post(
        "/projects/anonymization-review",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "upload_session_id": session_id,
            "asset_manifest": [
                {
                    "name": "busan_milk.txt",
                    "temp_file_id": upload["temp_file_id"],
                    "size": len(sml.encode("utf-8")),
                }
            ],
        },
    )
    assert preview_res.status_code == 200, preview_res.text
    preview_body = preview_res.json()
    raw_preview_text = preview_body["review_report"]["asset_previews"][0]["preview_text"]
    preview_text = preview_body["display_review_report"]["asset_previews"][0]["preview_text"]
    assert "부산우유" not in raw_preview_text
    assert "부산우유" not in preview_text
    assert "[CLIENT]" in raw_preview_text
    assert "[CLIENT]" in preview_text

    captured: dict[str, object] = {}

    def _capture_start(**kwargs):
        captured.update(kwargs)

    monkeypatch.setattr(projects_router, "start_project_wrapped_run", _capture_start)

    create_res = client.post(
        "/projects",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "project_name": "부산우유 alias masking",
            "client_name": "부산우유",
            "goal": "명확한 조직 alias masking 확인",
            "upload_session_id": session_id,
            "asset_manifest": [
                {
                    "name": "busan_milk.txt",
                    "temp_file_id": upload["temp_file_id"],
                    "size": len(sml.encode("utf-8")),
                }
            ],
            "template_key": "default_modernization_v1",
            "constraints": [],
        },
    )

    assert create_res.status_code == 200, create_res.text
    assert create_res.json()["status"] == "pending"

    safe_bundle = captured["safe_bundle"]
    canonical = "\n\n".join(source.content for source in safe_bundle.sources)
    stats = safe_bundle.sources[0].replacement_stats

    assert "부산우유" not in canonical
    assert "CLIENT_001" in canonical
    assert "CLIENT_001의 비전" in canonical
    assert "CLIENT_001 기초설계서" in canonical
    assert "CLIENT_001 원가 시스템 개선" in canonical
    assert stats["client_name"] == 1


def test_project_create_uses_source_question_guard_for_analysis_input(client, monkeypatch):
    from mellow_link.modules.rebuild_assistant.service import RebuildAssistantService
    from mellow_link.routers import projects as projects_router

    user = _register(client, "phase1_question_guard_create")
    session_id = f"phase1-question-guard-create-{uuid.uuid4().hex[:8]}"
    sml = "\n".join(
        [
            "[SML v1]",
            "presentation_file: cost_consulting_deck.pptx",
            "slide_count: 1",
            "",
            "[SLIDE 1]",
            "title: [CLIENT] 원가 컨설팅 개요",
            "texts:",
            "- 현행 원가체계 분석",
            "- 원가분석 및 원가계산 개선 방향",
            "- 재료비, 노무비, 제조경비 배부기준 검토",
            "- 손익분석 확장 검토",
        ]
    )
    upload = _upload_temp_asset(
        client,
        session_id=session_id,
        filename="cost_guard_create.txt",
        content=sml.encode("utf-8"),
        headers=user["headers"],
    )

    captured: dict[str, object] = {}

    def _capture_start(**kwargs):
        captured.update(kwargs)

    monkeypatch.setattr(projects_router, "start_project_wrapped_run", _capture_start)

    create_res = client.post(
        "/projects",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "project_name": "질문 가드 create",
            "client_name": "OO우유",
            "goal": "제품 저장 전 검증 로직을 어떻게 강화할 것인가?",
            "upload_session_id": session_id,
            "asset_manifest": [
                {
                    "name": "cost_guard_create.txt",
                    "temp_file_id": upload["temp_file_id"],
                    "size": len(sml.encode("utf-8")),
                }
            ],
            "template_key": "default_modernization_v1",
            "constraints": [
                "SQL 파라미터 검증을 어떻게 설계할 것인가?",
                "전면 재구축해야 하는가?",
            ],
        },
    )

    assert create_res.status_code == 200, create_res.text
    assert create_res.json()["status"] == "pending"
    assert any("질문 중 일부가 소스 근거 부족으로 분석 입력에서 제외되었습니다." in item for item in create_res.json()["warnings"])

    analysis_context = captured["analysis_context"]
    prepared = RebuildAssistantService().prepare_analysis_context_input(analysis_context=analysis_context)

    assert prepared.goal == "현행 원가체계의 한계는 무엇인가?"
    assert prepared.question_axis == "processing_flow"
    assert any(item.blocked_reason == "source_domain_mismatch" for item in prepared.blocked_user_questions)
    assert any(item.blocked_reason == "conclusion_forcing" for item in prepared.review_user_questions)
    assert prepared.question_guard_summary.selected_questions[0] == "현행 원가체계의 한계는 무엇인가?"


def test_runtime_console_pages_are_exposed(client):
    chat_res = client.get("/runtime-console")
    assert chat_res.status_code == 200
    assert "/runtime/turn" in chat_res.text
    assert "/runtime/status" in chat_res.text

    ops_res = client.get("/runtime-operator")
    assert ops_res.status_code == 200
    assert "/runtime/status" in ops_res.text

def test_runs_creation_assigns_default_session_and_list_shows_new_run(client):
    user = _register(client, "phase1_owner")

    create_res = client.post("/runs", headers=user["headers"])
    assert create_res.status_code == 200, create_res.text
    created = create_res.json()
    run_id = created["run_id"]
    assert created["session_id"]

    list_res = client.get("/runs", headers=user["headers"])
    assert list_res.status_code == 200, list_res.text
    runs = list_res.json()["runs"]
    matched = next((r for r in runs if r["run_id"] == run_id), None)
    assert matched is not None
    assert matched["session_id"] == created["session_id"]


def test_project_workspace_survives_refresh(client, monkeypatch):
    user = _register(client, "phase1_project")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-upload",
        filename="legacy.jsp",
        content=b"<% String sql = \"SELECT * FROM orders\"; %>",
        project_name="주문 화면 현대화",
        client_name="OO생명",
    )

    page_res = client.get(f"/projects/{project_id}", headers={"Accept": "text/html"})
    assert page_res.status_code == 200
    assert "분석 워크스페이스" in page_res.text

    snap_res = client.get(f"/projects/{project_id}?format=json", headers=user["headers"])
    assert snap_res.status_code == 200, snap_res.text
    snap = snap_res.json()
    assert snap["project"]["id"] == project_id
    assert snap["project"]["project_name"] == "주문 화면 현대화"


def test_project_result_markdown_download(client, monkeypatch):
    user = _register(client, "phase1_markdown")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-md",
        filename="legacy.sql",
        content=b"SELECT * FROM orders;",
        project_name="결과 패키지",
        client_name="OO카드",
    )
    res = client.get(f"/projects/{project_id}/result?format=md", headers=user["headers"])
    assert res.status_code == 200, res.text
    assert "text/markdown" in res.headers.get("content-type", "")
    assert "filename*=UTF-8''%EA%B2%B0%EA%B3%BC_%ED%8C%A8%ED%82%A4%EC%A7%80_result.md" in res.headers.get("content-disposition", "")


def test_project_result_markdown_download_external_surface(client, monkeypatch):
    user = _register(client, "phase1_markdown_external")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-md-external",
        filename="legacy.sql",
        content=b"SELECT * FROM orders;",
        project_name="외부 설명 패키지",
        client_name="OO카드",
    )
    res = client.get(f"/projects/{project_id}/result?format=md&surface_mode=external", headers=user["headers"])
    assert res.status_code == 200, res.text
    assert "text/markdown" in res.headers.get("content-type", "")
    assert res.text.startswith("# 컨설팅 결과 - 외부 설명 패키지")
    assert "## 컨설팅 개요" in res.text
    assert "## 컨설팅 비전" in res.text
    assert "## 핵심 판단" not in res.text
    assert "Review Diff" not in res.text
    assert "synthetic_signal_detected" not in res.text
    assert "severity" not in res.text.lower()
    assert "blast radius" not in res.text.lower()
    assert "filename*=UTF-8''%EC%99%B8%EB%B6%80_%EC%84%A4%EB%AA%85_%ED%8C%A8%ED%82%A4%EC%A7%80_external_result.md" in res.headers.get("content-disposition", "")


def test_project_result_json_external_surface_strips_review_diff(client, monkeypatch):
    user = _register(client, "phase1_json_external_surface")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-json-external",
        filename="legacy.sql",
        content=b"SELECT * FROM orders;",
        project_name="외부 JSON 패키지",
        client_name="OO카드",
    )

    res = client.get(f"/projects/{project_id}/result?format=json&surface_mode=external", headers=user["headers"])

    assert res.status_code == 200, res.text
    payload = res.json()
    assert "review_diff" not in (payload.get("extensions") or {})
    assert "decision_governance" not in (payload.get("extensions") or {})
    assert payload["provenance"]["surface_access"]["access_profile"] == "external_basic"
    assert payload["provenance"]["surface_access"]["field_visibility"]["review_diff"] == "absent"
    assert payload["provenance"]["surface_access"]["field_visibility"]["decision_governance"] == "absent"


def test_project_result_surface_modes_keep_same_canonical_decisions(client, monkeypatch):
    user = _register(client, "phase1_surface_consistency")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-surface-consistency",
        filename="legacy.sql",
        content=b"SELECT * FROM orders;",
        project_name="표면 일관성 패키지",
        client_name="OO카드",
    )

    internal = client.get(f"/projects/{project_id}/result?format=json&surface_mode=internal", headers=user["headers"])
    external = client.get(f"/projects/{project_id}/result?format=json&surface_mode=external", headers=user["headers"])

    assert internal.status_code == 200, internal.text
    assert external.status_code == 200, external.text
    internal_body = internal.json()
    external_body = external.json()
    assert internal_body["structural_judgment"] == external_body["structural_judgment"]
    assert internal_body["template_judgment"] == external_body["template_judgment"]
    assert internal_body["primary_judgment"] == external_body["primary_judgment"]
    assert internal_body["authoritative_payload"]["decision_summary"] == external_body["authoritative_payload"]["decision_summary"]
    assert internal_body["provenance"]["surface_access"]["access_profile"] == "internal_full"
    assert "review_diff" not in (external_body.get("extensions") or {})


def test_detect_domain_mismatch_warning_for_claim_goal_with_order_assets():
    from mellow_link.routers.projects import _detect_domain_mismatch_warning

    warnings = _detect_domain_mismatch_warning(
        project_name="청구 조정 기능을 현대적인 서비스 구조 재구성",
        constraints=["FRAUD 규칙 유지", "CLAIM_AUDIT 전담 조건 유지"],
        asset_names=["legacy.jsp", "OrderCloseService.java", "query.sql", "schema.sql"],
        asset_texts=[
            'if ("VIP".equals(order.getCustomerGrade())) return "vip_night_block";',
            'if ("Y".equals(order.getDeliveryHoldFlag())) return "delivery_hold_release_required";',
            'order.setStatus("REVIEW_REQUIRED");',
        ],
    )

    assert warnings
    assert "프로젝트 목표와 업로드 자산의 도메인 축이 일치하지 않을 가능성이 있습니다." in warnings[0]
    assert "'청구 조정'" in warnings[0]
    assert "'주문 마감'" in warnings[0]


def test_create_project_returns_domain_mismatch_warning(client, monkeypatch):
    from mellow_link.routers import projects as projects_router

    user = _register(client, "phase1_mismatch")
    monkeypatch.setattr(projects_router, "start_project_wrapped_run", lambda *args, **kwargs: None)

    upload_session_id = f"phase1-mismatch-{uuid.uuid4().hex[:8]}"
    uploads = [
        ("goal.txt", "청구 조정 기능을 현대적인 서비스 구조 재구성".encode("utf-8")),
        ("constraints.txt", "FRAUD 규칙 유지\nCLAIM_AUDIT 전담 조건 유지".encode("utf-8")),
        ("legacy.jsp", b'<button>close</button>'),
        ("OrderCloseService.java", b'if ("VIP".equals(order.getCustomerGrade())) return "vip_night_block";'),
        ("query.sql", b"SELECT * FROM sales_order WHERE status IN ('PAID','READY','REVIEW_REQUIRED')"),
        ("schema.sql", b"CREATE TABLE sales_order (status varchar(20), delivery_hold_flag varchar(1));"),
    ]
    asset_manifest = []
    for filename, content in uploads:
        uploaded = _upload_temp_asset(client, upload_session_id, filename, content, headers=user["headers"])
        asset_manifest.append({"name": filename, "temp_file_id": uploaded["temp_file_id"], "size": len(content)})

    res = client.post(
        "/projects",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "project_name": "청구 조정 기능을 현대적인 서비스 구조 재구성",
            "client_name": "SI센터",
            "upload_session_id": upload_session_id,
            "asset_manifest": asset_manifest,
            "template_key": "default_modernization_v1",
            "constraints": ["FRAUD 규칙 유지", "CLAIM_AUDIT 전담 조건 유지"],
        },
    )
    assert res.status_code == 200, res.text
    payload = res.json()
    assert payload["warnings"]

    detail = client.get(f"/projects/{payload['project_id']}?format=json", headers=user["headers"])
    assert detail.status_code == 200, detail.text
    detail_payload = detail.json()
    assert detail_payload["warnings"]


def test_create_project_preserves_explicit_goal_and_reuses_it_for_rerun_paths(client, monkeypatch):
    from mellow_link.routers import projects as projects_router

    captured_calls = []

    def _capture_start(*args, **kwargs):
        captured_calls.append(dict(kwargs))

    monkeypatch.setattr(projects_router, "start_project_wrapped_run", _capture_start)

    user = _register(client, "phase1_goal_preserve")
    upload_session_id = f"phase1-goal-preserve-{uuid.uuid4().hex[:8]}"
    upload = _upload_temp_asset(
        client,
        upload_session_id,
        "legacy_proc.sql",
        b"CREATE OR REPLACE PROCEDURE p_fx_fifo AS BEGIN NULL; END;",
        headers=user["headers"],
    )
    explicit_goal = "이 SQL/프로시저가 실제로 어떤 처리 흐름과 계산 규칙으로 동작하는지 분석해줘."

    create_res = client.post(
        "/projects",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "project_name": "선입선출 분석",
            "client_name": "OO은행",
            "goal": explicit_goal,
            "upload_session_id": upload_session_id,
            "asset_manifest": [{"name": "legacy_proc.sql", "temp_file_id": upload["temp_file_id"], "size": 58}],
            "template_key": "default_modernization_v1",
            "constraints": [],
        },
    )
    assert create_res.status_code == 200, create_res.text
    project_id = create_res.json()["project_id"]
    assert captured_calls
    assert captured_calls[-1]["goal"] == explicit_goal

    detail_res = client.get(f"/projects/{project_id}?format=json", headers=user["headers"])
    assert detail_res.status_code == 200, detail_res.text
    assert detail_res.json()["project"]["goal"] == explicit_goal

    captured_calls.clear()
    rerun_res = client.post(f"/projects/{project_id}/run", headers=user["headers"])
    assert rerun_res.status_code == 200, rerun_res.text
    assert captured_calls
    assert captured_calls[-1]["goal"] == explicit_goal

    captured_calls.clear()
    reanalysis_res = client.post(
        f"/projects/{project_id}/reanalysis",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={"new_asset_manifest": []},
    )
    assert reanalysis_res.status_code == 200, reanalysis_res.text
    assert captured_calls
    assert captured_calls[-1]["goal"] == explicit_goal


def test_create_project_uses_goal_txt_when_inline_goal_is_blank(client, monkeypatch):
    from mellow_link.routers import projects as projects_router

    captured_calls = []

    def _capture_start(*args, **kwargs):
        captured_calls.append(dict(kwargs))

    monkeypatch.setattr(projects_router, "start_project_wrapped_run", _capture_start)

    user = _register(client, "phase1_goal_txt")
    upload_session_id = f"phase1-goal-txt-{uuid.uuid4().hex[:8]}"
    goal_text = "이 SQL/프로시저의 FIFO lot 처리 흐름과 계산 규칙을 분석해줘."
    uploads = [
        ("goal.txt", goal_text.encode("utf-8")),
        ("legacy_proc.sql", b"CREATE OR REPLACE PROCEDURE p_fifo AS BEGIN NULL; END;"),
    ]
    asset_manifest = []
    for filename, content in uploads:
        uploaded = _upload_temp_asset(client, upload_session_id, filename, content, headers=user["headers"])
        asset_manifest.append({"name": filename, "temp_file_id": uploaded["temp_file_id"], "size": len(content)})

    create_res = client.post(
        "/projects",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "project_name": "FIFO 분석",
            "client_name": "OO증권",
            "goal": "",
            "upload_session_id": upload_session_id,
            "asset_manifest": asset_manifest,
            "template_key": "default_modernization_v1",
            "constraints": [],
        },
    )
    assert create_res.status_code == 200, create_res.text
    project_id = create_res.json()["project_id"]
    assert captured_calls
    assert captured_calls[-1]["goal"] == goal_text

    detail_res = client.get(f"/projects/{project_id}?format=json", headers=user["headers"])
    assert detail_res.status_code == 200, detail_res.text
    assert detail_res.json()["project"]["goal"] == goal_text


def test_temp_upload_returns_extracted_text_for_autofill(client):
    body = _upload_temp_asset(
        client,
        session_id=f"phase1-autofill-{uuid.uuid4().hex[:8]}",
        filename="goal.txt",
        content="주문 관리 화면 현대화\n".encode("utf-8"),
        content_type="text/plain",
    )
    assert body["filename"] == "goal.txt"
    assert body["extracted_text"].strip() == "주문 관리 화면 현대화"


def test_project_result_docx_download(client, monkeypatch, tmp_path):
    from types import SimpleNamespace
    from mellow_link import app_state

    user = _register(client, "phase1_docx")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-docx",
        filename="legacy.sql",
        content=b"SELECT * FROM orders;",
        project_name="DOCX 결과 패키지",
        client_name="OO카드",
    )

    output_path = tmp_path / "project_result.docx"
    output_path.write_bytes(b"fake-docx-content")

    class FakeDocService:
        def is_available(self):
            return True

        async def generate(self, request):
            assert request.output_type.value == "docx"
            assert request.filename == "DOCX_결과_패키지_result.docx"
            assert request.content.startswith("# 컨설팅 결과 - DOCX 결과 패키지")
            assert "## 컨설팅 개요" in request.content
            assert "## 컨설팅 설계" in request.content
            assert "## 참고 구조 비교" not in request.content
            return SimpleNamespace(output_path=output_path)

    monkeypatch.setattr(app_state, "doc_service", FakeDocService(), raising=False)

    res = client.get(f"/projects/{project_id}/result?format=docx", headers=user["headers"])
    assert res.status_code == 200, res.text
    assert "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in res.headers.get("content-type", "")
    assert "filename*=utf-8''docx_%ea%b2%b0%ea%b3%bc_%ed%8c%a8%ed%82%a4%ec%a7%80_result.docx" in res.headers.get("content-disposition", "").lower()
    assert res.content == b"fake-docx-content"


def test_project_result_docx_download_external_surface(client, monkeypatch, tmp_path):
    from types import SimpleNamespace
    from mellow_link import app_state

    user = _register(client, "phase1_docx_external")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-docx-external",
        filename="legacy.sql",
        content=b"SELECT * FROM orders;",
        project_name="DOCX 외부 설명 패키지",
        client_name="OO카드",
    )

    output_path = tmp_path / "project_result_external.docx"
    output_path.write_bytes(b"fake-docx-external-content")

    class FakeDocService:
        def is_available(self):
            return True

        async def generate(self, request):
            assert request.output_type.value == "docx"
            assert request.filename == "DOCX_외부_설명_패키지_external_result.docx"
            assert request.content.startswith("# 컨설팅 결과 - DOCX 외부 설명 패키지")
            assert "## 컨설팅 개요" in request.content
            assert "## 컨설팅 비전" in request.content
            assert "Review Diff" not in request.content
            assert "synthetic_signal_detected" not in request.content
            return SimpleNamespace(output_path=output_path)

    monkeypatch.setattr(app_state, "doc_service", FakeDocService(), raising=False)

    res = client.get(f"/projects/{project_id}/result?format=docx&surface_mode=external", headers=user["headers"])
    assert res.status_code == 200, res.text
    assert "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in res.headers.get("content-type", "")
    assert "filename*=utf-8''docx_%ec%99%b8%eb%b6%80_%ec%84%a4%eb%aa%85_%ed%8c%a8%ed%82%a4%ec%a7%80_external_result.docx" in res.headers.get("content-disposition", "").lower()
    assert res.content == b"fake-docx-external-content"


def test_project_result_pptx_download(client, monkeypatch, tmp_path):
    from types import SimpleNamespace
    from mellow_link import app_state

    user = _register(client, "phase1_pptx")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-pptx",
        filename="legacy.sql",
        content=b"SELECT * FROM orders;",
        project_name="PPTX 결과 패키지",
        client_name="OO카드",
    )

    output_path = tmp_path / "project_result.pptx"
    output_path.write_bytes(b"fake-pptx-content")

    class FakeDocService:
        def is_available(self):
            return True

        async def generate(self, request):
            assert request.output_type.value == "pptx"
            assert request.filename == "PPTX_결과_패키지_result.pptx"
            assert request.content == ""
            assert request.style_options["renderer"] == "slide_schema"
            assert request.payload["schema_version"] == "slide_schema.v1"
            assert request.payload["surface_mode"] == "internal"
            assert any(slide.get("title") == "컨설팅 개요" for slide in request.payload["slides"])
            return SimpleNamespace(output_path=output_path)

    monkeypatch.setattr(app_state, "doc_service", FakeDocService(), raising=False)

    res = client.get(f"/projects/{project_id}/result?format=pptx", headers=user["headers"])
    assert res.status_code == 200, res.text
    assert "application/vnd.openxmlformats-officedocument.presentationml.presentation" in res.headers.get("content-type", "")
    assert "filename*=utf-8''pptx_%ea%b2%b0%ea%b3%bc_%ed%8c%a8%ed%82%a4%ec%a7%80_result.pptx" in res.headers.get("content-disposition", "").lower()
    assert res.content == b"fake-pptx-content"


def test_project_result_pptx_download_external_surface(client, monkeypatch, tmp_path):
    from types import SimpleNamespace
    from mellow_link import app_state

    user = _register(client, "phase1_pptx_external")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-pptx-external",
        filename="legacy.sql",
        content=b"SELECT * FROM orders;",
        project_name="PPTX 외부 설명 패키지",
        client_name="OO카드",
    )

    output_path = tmp_path / "project_result_external.pptx"
    output_path.write_bytes(b"fake-pptx-external-content")

    class FakeDocService:
        def is_available(self):
            return True

        async def generate(self, request):
            assert request.output_type.value == "pptx"
            assert request.filename == "PPTX_외부_설명_패키지_external_result.pptx"
            assert request.content == ""
            assert request.style_options["renderer"] == "slide_schema"
            assert request.payload["schema_version"] == "slide_schema.v1"
            assert request.payload["surface_mode"] == "external"
            assert any(slide.get("title") == "컨설팅 개요" for slide in request.payload["slides"])
            return SimpleNamespace(output_path=output_path)

    monkeypatch.setattr(app_state, "doc_service", FakeDocService(), raising=False)

    res = client.get(f"/projects/{project_id}/result?format=pptx&surface_mode=external", headers=user["headers"])
    assert res.status_code == 200, res.text
    assert "application/vnd.openxmlformats-officedocument.presentationml.presentation" in res.headers.get("content-type", "")
    assert "filename*=utf-8''pptx_%ec%99%b8%eb%b6%80_%ec%84%a4%eb%aa%85_%ed%8c%a8%ed%82%a4%ec%a7%80_external_result.pptx" in res.headers.get("content-disposition", "").lower()
    assert res.content == b"fake-pptx-external-content"


def test_project_asset_download_and_assets_payload(client, monkeypatch):
    from mellow_link.infra import ProjectAsset, TempResource

    owner = _register(client, "phase1_asset_owner")
    other = _register(client, "phase1_asset_other")
    project_id, upload = _create_persisted_project(
        client,
        owner,
        monkeypatch,
        upload_session_id="phase1-download",
        filename="legacy.txt",
        content=b"legacy modernize me",
        project_name="다운로드 검증",
        client_name="OO증권",
    )

    detail = client.get(f"/projects/{project_id}?format=json", headers=owner["headers"])
    assert detail.status_code == 200, detail.text
    assets = detail.json()["assets"]
    assert len(assets) == 1
    assert assets[0]["project_asset_id"]
    assert assets[0]["download_url"].endswith("/download")
    assert assets[0]["extracted_chars"] > 0
    assert assets[0]["content_type"] == "text/plain"
    assert assets[0]["uploaded_at"]
    assert assets[0]["stage_status"] == "promoted"
    assert assets[0]["is_downloadable"] is True

    result = client.get(f"/projects/{project_id}/result?format=json", headers=owner["headers"])
    assert result.status_code == 200, result.text
    result_json = result.json()
    result_assets = result_json["assets"]
    assert result_assets[0]["content_type"] == "text/plain"
    assert result_assets[0]["uploaded_at"]
    assert result_assets[0]["stage_status"] == "promoted"
    assert result_assets[0]["is_downloadable"] is True
    provenance = result_json["provenance"]
    assert provenance["run_id"]
    assert provenance["module_id"] == "rebuild_assistant"
    assert provenance["run_kind"] == "rebuild_plan"
    assert provenance["app_version"] == "0.1.0"
    assert provenance["module_version"] == "0.1.0"
    assert provenance["run_status"] in ("pending", "running", "completed", "failed")
    assert provenance["generated_at"]
    assert provenance["generated_at"].endswith("Z")
    assert provenance["input_assets"] == result_assets

    download = client.get(assets[0]["download_url"], headers=owner["headers"])
    assert download.status_code == 200, download.text
    assert "attachment" in download.headers.get("content-disposition", "").lower()
    assert download.content == b"legacy modernize me"

    forbidden = client.get(assets[0]["download_url"], headers=other["headers"])
    assert forbidden.status_code == 403

    from mellow_link.infra.database import SessionLocal

    with SessionLocal() as db:
        temp = db.query(TempResource).filter(TempResource.temp_file_id == upload["temp_file_id"]).first()
        project_asset = db.query(ProjectAsset).filter(ProjectAsset.project_id == project_id).first()
        assert temp is not None
        assert temp.stage_status == "promoted"
        assert temp.promoted_to_project_id == project_id
        assert project_asset is not None


def test_project_assets_fallback_payload_keeps_structure_with_null_metadata(client):
    from mellow_link.services.scope_notice import PROJECT_SCOPE_NOTICE

    user = _register(client, "phase1_fallback")
    project_id = _create_fallback_project(user)

    detail = client.get(f"/projects/{project_id}?format=json", headers=user["headers"])
    assert detail.status_code == 200, detail.text
    assets = detail.json()["assets"]
    assert len(assets) == 1
    assert assets[0]["name"] == "legacy.jsp"
    assert assets[0]["content_type"] is None
    assert assets[0]["uploaded_at"] is None
    assert assets[0]["stage_status"] is None
    assert assets[0]["is_downloadable"] is False

    result = client.get(f"/projects/{project_id}/result?format=json", headers=user["headers"])
    assert result.status_code == 200, result.text
    result_json = result.json()
    result_assets = result_json["assets"]
    assert result_assets[0]["content_type"] is None
    assert result_assets[0]["is_downloadable"] is False
    provenance = result_json["provenance"]
    assert provenance["input_assets"] == result_assets
    assert provenance["app_version"] == "0.1.0"
    assert provenance["module_version"] == "0.1.0"
    scope_notice = result_json["scope_notice"]
    assert scope_notice == PROJECT_SCOPE_NOTICE


def test_failed_project_result_includes_scope_notice(client):
    from mellow_link.services.scope_notice import PROJECT_SCOPE_NOTICE

    user = _register(client, "phase1_scope_failed")
    project_id = _create_fallback_project(user, project_name="실패 범위 고지", status="failed")

    result = client.get(f"/projects/{project_id}/result?format=json", headers=user["headers"])
    assert result.status_code == 200, result.text
    body = result.json()
    assert body["project"]["status"] == "failed"
    assert body["scope_notice"] == PROJECT_SCOPE_NOTICE


def test_project_html_templates_include_asset_metadata_and_access_states(client):
    workspace = client.get("/projects/create", headers={"Accept": "text/html"})
    assert workspace.status_code == 200
    assert "단일 기능 / 단일 화면 V0" in workspace.text
    assert "자동 코드 치환" in workspace.text
    assert "설명 가능한 결과 패키지 제공" in workspace.text
    assert "PROJECT_SCOPE_NOTICE" in workspace.text
    # create page is not target here; use actual workspace/result static entry pages
    user_console = client.get("/projects/some-project", headers={"Accept": "text/html"})
    assert user_console.status_code == 200
    assert "업로드 시각" in user_console.text
    assert "재열람 불가" in user_console.text
    assert "Unknown" in user_console.text
    assert "추가 자료 업로드" in user_console.text
    assert "입력 익명화" in user_console.text
    assert "실행 이력" in user_console.text
    assert "재분석" in user_console.text

    result = client.get("/projects/some-project/result", headers={"Accept": "text/html"})
    assert result.status_code == 200
    assert "생성 이력" in result.text
    assert "실행 ID" in result.text
    assert "모듈 버전" in result.text
    assert "/static/anonymization_display.js" in result.text
    assert "표시 안내: [PERSON], [ORG], [PROJECT], [TERM] 등은 익명 처리된 민감정보입니다." in result.text
    assert "검토용 DOCX" in result.text
    assert "검토용 PPTX" in result.text
    assert "설명용 내보내기" in result.text
    assert "결정 요약" in result.text
    assert "판단 근거" in result.text
    assert "실행 조치" in result.text
    assert "유지 계약" in result.text
    assert "설계 선택지" in result.text
    assert "실행 계획" in result.text
    assert "범위 및 한계" in result.text
    assert "범위 안내 누락" in result.text
    assert "비지원 범위" in result.text
    assert "다운로드 가능" in result.text
    assert "재열람 불가" in result.text
    assert "업로드 시각" in result.text
    assert "STATIC_SCOPE_NOTICE" not in result.text

    admin = _register_admin(client, "phase1_dev_console")
    dev_console = client.get("/dev-console", headers={**admin["headers"], "Accept": "text/html"})
    assert dev_console.status_code == 200, dev_console.text
    assert "Anonymization" in dev_console.text
    assert "Summary + Validation + Preview" in dev_console.text
    assert "Review Order" in dev_console.text
    assert "기본 통과 기준은 validation.passed=true" in dev_console.text
    assert "가장 먼저 확인할 항목입니다." in dev_console.text
    assert "검증 보조용이며 canonical이 아닙니다." in dev_console.text


def test_project_result_markdown_includes_provenance_section(client, monkeypatch):
    from mellow_link.services.scope_notice import PROJECT_SCOPE_NOTICE

    user = _register(client, "phase1_md_provenance")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-md-prov",
        filename="legacy.sql",
        content=b"SELECT 1;",
        project_name="마크다운 provenance",
        client_name="OO금융",
    )
    res = client.get(f"/projects/{project_id}/result?format=md&internal_export_mode=full", headers=user["headers"])
    assert res.status_code == 200, res.text
    assert res.text.startswith("# 결과 패키지 - 마크다운 provenance")
    assert "## 컨설팅 개요" in res.text
    assert "## 참고 구조 비교" in res.text
    assert "참고 구조 비교 데이터가 부족해 변경 근거를 다시 모아야 합니다." in res.text
    assert "Run ID:" in res.text
    assert "앱 버전:" in res.text
    assert "모듈 버전:" in res.text
    assert "## 부록" in res.text
    assert PROJECT_SCOPE_NOTICE["summary"] in res.text


def test_build_result_package_sanitizes_forbidden_user_tokens():
    from mellow_link.infra import ModernizationProject
    from mellow_link.routers.projects import build_result_package
    from mellow_link.modules.rebuild_assistant.schemas import EvidenceRef, GroundedBusinessRule, StructuredRebuildResult

    project = ModernizationProject(
        id="proj_test_sanitize",
        user_id=1,
        session_id="sess_sanitize",
        run_id="run_sanitize",
        project_name="SANITIZE",
        client_name="OO",
        template_key="default_modernization_v1",
        template_mode="recommended",
        constraints_json="[]",
        upload_session_id="upload_sanitize",
        asset_manifest_json="[]",
        status="completed",
    )
    result = StructuredRebuildResult(
        one_line_conclusion="REDACTED_PATH 와 role/... 와 권한[] 표기는 제거되어야 합니다.",
        executive_summary_v2=["controller/... 표현과 [SAFE STRUCTURE: asset_deadbeef] 표시는 없이 설명해야 합니다."],
        grounded_business_rules=[
            GroundedBusinessRule(
                title="샘플 규칙",
                description="SAFE STRUCTURE 와 내부 토큰 없이 보여야 합니다.",
                evidence=[
                    EvidenceRef(
                        asset_name="legacy.jsp",
                        asset_type="ui",
                        locator="본문 키워드",
                        excerpt="[SAFE STRUCTURE: asset_deadbeef] node:table:TBL_001 role/... controller/... 권한[]",
                        evidence_kind="ui",
                    )
                ],
                design_targets=["정책 서비스"],
                confidence="확정",
                confidence_reason="schema.sql 과 query.sql 에서 직접 확인되었습니다.",
                needs_verification=False,
            )
        ],
    )

    pkg = build_result_package(project, {"status": "completed"}, result, assets=[], app_version="0.1.0")
    dumped = json.dumps(pkg, ensure_ascii=False)

    assert "REDACTED_PATH" not in dumped
    assert "role/..." not in dumped
    assert "controller/..." not in dumped
    assert "권한[]" not in dumped
    assert "SAFE STRUCTURE" not in dumped
    assert "controller/service/repository" not in dumped
    assert "query parameters" not in dumped
    assert "가장 적합합니다" not in dumped
    assert "가능합니다" not in dumped
    assert "검토가 필요합니다" not in dumped
    assert "evidence" not in (pkg["grounded_business_rules"][0].keys())
    assert pkg["grounded_business_rules"][0]["evidence_cards"]


def test_build_result_package_condition_summary_translates_source_and_sql_evidence():
    from mellow_link.infra import ModernizationProject
    from mellow_link.routers.projects import build_result_package
    from mellow_link.modules.rebuild_assistant.schemas import EvidenceRef, GroundedBusinessRule, StructuredRebuildResult

    project = ModernizationProject(
        id="proj_test_condition_summary",
        user_id=1,
        session_id="sess_condition_summary",
        run_id="run_condition_summary",
        project_name="조건 요약",
        client_name="OO",
        template_key="default_modernization_v1",
        template_mode="recommended",
        constraints_json="[]",
        upload_session_id="upload_condition_summary",
        asset_manifest_json="[]",
        status="completed",
    )
    result = StructuredRebuildResult(
        grounded_business_rules=[
            GroundedBusinessRule(
                title="수출 주문 고액건 REVIEW_REQUIRED",
                description="수출 주문의 고액 건은 즉시 마감하지 않고 REVIEW_REQUIRED 상태로 전환해야 합니다.",
                evidence=[
                    EvidenceRef(
                        asset_name="OrderCloseService.java",
                        asset_type="source",
                        locator="본문 키워드",
                        excerpt='if ("EXPORT".equals(order.getOrderType()) && order.getOrderAmount() >= 7000000) { order.setStatus("REVIEW_REQUIRED"); }',
                        evidence_kind="source",
                    ),
                    EvidenceRef(
                        asset_name="query.sql",
                        asset_type="sql",
                        locator="본문 키워드",
                        excerpt="AND o.status IN ('PAID', 'READY', 'REVIEW_REQUIRED')",
                        evidence_kind="sql",
                    ),
                ],
                design_targets=["상태 전이", "정책 서비스"],
                confidence="확정",
                confidence_reason="현재 자산에서 직접 확인되었습니다.",
                needs_verification=False,
            )
        ]
    )

    pkg = build_result_package(project, {"status": "completed"}, result, assets=[], app_version="0.1.0")
    cards = pkg["grounded_business_rules"][0]["evidence_cards"]

    assert any("REVIEW_REQUIRED 상태로 전이" in card["condition_summary"] for card in cards)
    assert any("상태값이 PAID, READY, REVIEW_REQUIRED" in card["condition_summary"] for card in cards)


def test_build_result_package_condition_summary_generalizes_role_and_status_conditions():
    from mellow_link.infra import ModernizationProject
    from mellow_link.routers.projects import build_result_package
    from mellow_link.modules.rebuild_assistant.schemas import EvidenceRef, GroundedBusinessRule, StructuredRebuildResult

    project = ModernizationProject(
        id="proj_test_condition_general",
        user_id=1,
        session_id="sess_condition_general",
        run_id="run_condition_general",
        project_name="조건 일반화",
        client_name="OO",
        template_key="default_modernization_v1",
        template_mode="recommended",
        constraints_json="[]",
        upload_session_id="upload_condition_general",
        asset_manifest_json="[]",
        status="completed",
    )
    result = StructuredRebuildResult(
        grounded_business_rules=[
            GroundedBusinessRule(
                title="고액 청구 전담 부서 처리",
                description="고액 청구는 전담 부서만 처리할 수 있습니다.",
                evidence=[
                    EvidenceRef(
                        asset_name="legacy_app.py",
                        asset_type="source",
                        locator="본문 키워드",
                        excerpt='if claim["claim_amount"] >= 10000000 and dept_code != "CLAIM_AUDIT": return "심사전담부서만 가능"',
                        evidence_kind="source",
                    ),
                    EvidenceRef(
                        asset_name="claim_adjustment.html",
                        asset_type="ui",
                        locator="본문 키워드",
                        excerpt='<c:if test="${status eq \'READY\'}"><button>조정</button></c:if>',
                        evidence_kind="ui",
                    ),
                ],
                design_targets=["정책 서비스", "검증 흐름"],
                confidence="확정",
                confidence_reason="현재 자산에서 직접 확인되었습니다.",
                needs_verification=False,
            )
        ]
    )

    pkg = build_result_package(project, {"status": "completed"}, result, assets=[], app_version="0.1.0")
    cards = pkg["grounded_business_rules"][0]["evidence_cards"]

    assert any("CLAIM_AUDIT 부서가 아니면 고액 청구를 처리할 수 없도록 제한" in card["condition_summary"] for card in cards)
    assert any("상태값이 READY일 때만 화면 액션" in card["condition_summary"] for card in cards)


def test_build_result_package_condition_summary_uses_rule_title_fallback_for_broken_excerpt():
    from mellow_link.infra import ModernizationProject
    from mellow_link.routers.projects import build_result_package
    from mellow_link.modules.rebuild_assistant.schemas import EvidenceRef, GroundedBusinessRule, StructuredRebuildResult

    project = ModernizationProject(
        id="proj_test_condition_title_fallback",
        user_id=1,
        session_id="sess_condition_title_fallback",
        run_id="run_condition_title_fallback",
        project_name="조건 타이틀 fallback",
        client_name="OO",
        template_key="default_modernization_v1",
        template_mode="recommended",
        constraints_json="[]",
        upload_session_id="upload_condition_title_fallback",
        asset_manifest_json="[]",
        status="completed",
    )
    result = StructuredRebuildResult(
        grounded_business_rules=[
            GroundedBusinessRule(
                title="대리점 고액 주문 본사 전용",
                description="대리점 채널의 고액 주문은 본사 권한으로만 마감할 수 있습니다.",
                evidence=[
                    EvidenceRef(
                        asset_name="OrderCloseService.java",
                        asset_type="source",
                        locator="본문 키워드",
                        excerpt='; } if ("AGENCY".equals(order.getChannelCode())',
                        evidence_kind="source",
                    )
                ],
                design_targets=["정책 서비스", "권한 모델", "API"],
                confidence="확정",
                confidence_reason="현재 자산에서 직접 확인되었습니다.",
                needs_verification=False,
            )
        ]
    )

    pkg = build_result_package(project, {"status": "completed"}, result, assets=[], app_version="0.1.0")
    cards = pkg["grounded_business_rules"][0]["evidence_cards"]

    assert any("대리점 채널의 고액 주문은 본사 승인 조건을 충족해야 처리되도록 제한" in card["condition_summary"] for card in cards)


def test_run_finished_preserves_large_structured_result_and_result_package_becomes_ready(tmp_path):
    import uuid

    from mellow_link.infra import ModernizationProject
    from mellow_link.infra.run_events import (
        EVENT_TYPE_RUN_FINISHED,
        create_run,
        emit_event,
        get_run_events,
    )
    from mellow_link.modules.rebuild_assistant.schemas import (
        DecisionItem,
        DesignOption,
        ExecutionPlanWeek,
        GroundedBusinessRule,
        RecommendedOption,
        RetainedContract,
        StructuredRebuildResult,
    )
    from mellow_link.routers.projects import _extract_structured_result, build_result_package

    large_lines = [f"핵심 요약 {index} " + ("세부 설명 " * 12) for index in range(12)]
    result = StructuredRebuildResult(
        one_line_conclusion="주문 마감 기능을 단계적으로 분리하는 것이 필요합니다.",
        executive_summary_v2=large_lines[:4],
        core_business_rules=large_lines[:5],
        grounded_business_rules=[
            GroundedBusinessRule(
                title="VIP 야간 마감 제한",
                description=large_lines[0],
                evidence=[],
                design_targets=["정책 서비스"],
                confidence="가정",
                confidence_reason="테스트용 대형 결과입니다.",
                needs_verification=True,
            )
        ],
        decision_items=[
            DecisionItem(statement="주문 마감 정책을 별도 서비스로 분리하는 것이 필요합니다.", rationale=large_lines[1])
        ],
        retained_contracts=[
            RetainedContract(item="주문 상태 코드는 유지하는 것이 필요합니다.", basis=large_lines[2], evidence=[])
        ],
        design_options=[
            DesignOption(
                name="옵션 A",
                structure_summary=large_lines[3],
                advantages=large_lines[:2],
                risks=large_lines[2:4],
                difficulty="MEDIUM",
                duration_weeks=4,
                recommended=True,
                selection_reason=large_lines[4],
            )
        ],
        recommended_option=RecommendedOption(
            name="옵션 A",
            structure_summary=large_lines[5],
            selection_reason=large_lines[6],
            expected_outcomes=large_lines[7:9],
        ),
        execution_plan=[
            ExecutionPlanWeek(
                week_label="1주차",
                goal=large_lines[7],
                tasks=large_lines[8:10],
                roles=["컨설턴트", "아키텍트"],
                duration_weeks=1,
                deliverables=["규칙 목록", "분리안"],
            )
        ],
        analysis_summary=large_lines,
        rebuild_strategy=large_lines,
        risks=large_lines[:4],
        recommended_directions=large_lines[:3],
    )

    project_id = f"proj_large_structured_{uuid.uuid4().hex[:8]}"
    upload_session_id = f"upload_large_structured_{uuid.uuid4().hex[:8]}"

    isolated_session, engine = _build_isolated_session_factory(tmp_path)
    try:
        with isolated_session() as db:
            user = _seed_isolated_user(db, username_prefix="phase1_large_structured")
            run_id = create_run(
                session_id="sess_large_structured",
                db=db,
                module_id="rebuild_assistant",
                run_kind="rebuild_plan",
            )
            project = ModernizationProject(
                id=project_id,
                user_id=user.id,
                session_id="sess_large_structured",
                run_id=run_id,
                project_name="대형 structured result",
                client_name="OO",
                template_key="default_modernization_v1",
                template_mode="recommended",
                constraints_json="[]",
                upload_session_id=upload_session_id,
                asset_manifest_json="[]",
                status="completed",
            )
            db.add(project)
            db.commit()
            emit_event(
                run_id,
                EVENT_TYPE_RUN_FINISHED,
                {
                    "success": True,
                    "summary": "summary " * 300,
                    "structured_result": result.model_dump(),
                    "primary_feature_mode": "status_permissions",
                    "secondary_feature_mode": "save_validation",
                    "confidence": result.confidence,
                    "needs_more_input": False,
                    "scope_limited": False,
                    "module_id": "rebuild_assistant",
                    "run_kind": "rebuild_plan",
                },
                db=db,
            )
            events = get_run_events(run_id, db=db)
            run_finished = next(event for event in events if event["type"] == "run_finished")
            assert isinstance(run_finished["payload"]["structured_result"], dict)

            extracted = _extract_structured_result(events)
            assert extracted is not None
            assert extracted.one_line_conclusion == result.one_line_conclusion
            pkg = build_result_package(
                project,
                {"status": "completed", "run_id": run_id},
                extracted,
                assets=[],
                app_version="0.1.0",
            )
            assert pkg["executive_summary"]["state"] == "ready"
            assert pkg["core_conclusion"] == result.one_line_conclusion
            assert pkg["confidence"] == result.confidence
    finally:
        engine.dispose()


def test_extract_project_insights_falls_back_to_rebuild_strategy_for_sparse_structured_result():
    from mellow_link.modules.rebuild_assistant.schemas import StructuredRebuildResult
    from mellow_link.routers.projects import _extract_project_insights

    result = StructuredRebuildResult(
        one_line_conclusion="주문 마감 기능을 분리합니다.",
        analysis_summary=["현행 분석 요약"],
        rebuild_strategy=[
            "주문 마감 정책을 별도 서비스로 분리합니다.",
            "상태 전이 검증을 API 경계 밖으로 이동합니다.",
            "기존 상태 코드는 유지합니다.",
        ],
    )

    insights = _extract_project_insights(
        [{"type": "run_finished", "payload": {"primary_feature_mode": "status_permissions"}}],
        result,
    )

    assert insights["feature_mode"] == "권한 및 상태 규칙"
    assert insights["findings"] == ["현행 분석 요약"]
    assert insights["rules"] == [
        "주문 마감 정책을 별도 서비스로 분리합니다.",
        "상태 전이 검증을 API 경계 밖으로 이동합니다.",
        "기존 상태 코드는 유지합니다.",
    ]


def test_build_result_package_exposes_customer_intent_separately(tmp_path):
    from mellow_link.infra import ModernizationProject
    from mellow_link.infra.run_events import create_run
    from mellow_link.modules.rebuild_assistant.schemas import StructuredRebuildResult
    from mellow_link.routers.projects import build_result_package

    result = StructuredRebuildResult(
        one_line_conclusion="고객 요청을 분리해 보여줍니다.",
        confidence=0.72,
    )

    isolated_session, engine = _build_isolated_session_factory(tmp_path)
    try:
        with isolated_session() as db:
            user = _seed_isolated_user(db, username_prefix="phase1_customer_intent")
            run_id = create_run(
                session_id="sess_customer_intent",
                db=db,
                module_id="rebuild_assistant",
                run_kind="rebuild_plan",
            )
            project = ModernizationProject(
                id=f"proj_customer_intent_{uuid.uuid4().hex[:8]}",
                user_id=user.id,
                session_id="sess_customer_intent",
                run_id=run_id,
                project_name="고객 요청 노출 확인",
                client_name="OO",
                template_key="default_modernization_v1",
                template_mode="recommended",
                constraints_json=json.dumps(
                    [
                        "기존 DB 계약 유지",
                        "고객 요청: 화면은 최대한 유지하고 백엔드만 분리",
                        "고객 요청: 단계적 전환을 원함",
                    ],
                    ensure_ascii=False,
                ),
                upload_session_id=f"upload_customer_intent_{uuid.uuid4().hex[:8]}",
                asset_manifest_json="[]",
                status="completed",
            )
            db.add(project)
            db.commit()

            pkg = build_result_package(
                project,
                {"status": "completed", "run_id": run_id},
                result,
                assets=[],
                app_version="0.1.0",
            )
            assert pkg["customer_intent"]["available"] is True
            assert pkg["customer_intent"]["items"] == [
                "화면은 최대한 유지하고 백엔드만 분리",
                "단계적 전환을 원함",
            ]
    finally:
        engine.dispose()


def test_run_finished_storage_preserves_polish_bundle_dict_without_emit_mock(client):
    from mellow_link.infra import User
    from mellow_link.infra.database import SessionLocal
    from mellow_link.infra.run_events import EVENT_TYPE_RUN_FINISHED, create_run, emit_event, get_run_events
    from mellow_link.routers.projects import _extract_polish_bundle, _extract_structured_result
    from mellow_link.routers.runs import _resolve_run_session_id

    user = _register(client, username_prefix="phase1_polish_storage")
    structured_result, polish_bundle = _build_large_rebuild_result_with_polish_bundle()

    with SessionLocal() as db:
        db_user = db.query(User).filter(User.username == user["username"]).first()
        session_id = _resolve_run_session_id(db, db_user, None)
        run_id = create_run(session_id=session_id, db=db, module_id="rebuild_assistant", run_kind="rebuild_plan")
        emit_event(
            run_id,
            EVENT_TYPE_RUN_FINISHED,
            {
                "success": True,
                "summary": "completed",
                "structured_result": structured_result,
                "polish_bundle": polish_bundle,
                "primary_feature_mode": "save_validation",
                "module_id": "rebuild_assistant",
                "run_kind": "rebuild_plan",
            },
            db=db,
        )
        events = get_run_events(run_id, db=db)

    run_finished = next(event for event in events if event["type"] == EVENT_TYPE_RUN_FINISHED)
    assert isinstance(run_finished["payload"]["polish_bundle"], dict)
    assert isinstance(run_finished["payload"]["polish_bundle"]["polished_sections"], list)
    assert run_finished["payload"]["polish_bundle"]["original_result"]["primary_judgment"] == structured_result["primary_judgment"]

    extracted = _extract_structured_result(events)
    polish = _extract_polish_bundle(events, extracted)

    assert extracted is not None
    assert isinstance(polish, dict)
    assert isinstance(polish["warnings"], list)


def test_extract_polish_bundle_rebuilds_when_event_payload_is_not_dict():
    from mellow_link.routers.projects import _extract_polish_bundle

    structured_result, _ = _build_large_rebuild_result_with_polish_bundle()
    events = [
        {
            "type": "run_finished",
            "payload": {
                "structured_result": structured_result,
                "polish_bundle": "[TRUNCATED_OBJECT]",
            },
        }
    ]

    bundle = _extract_polish_bundle(events)

    assert isinstance(bundle, dict)
    assert bundle["primary_judgment"] == structured_result["primary_judgment"]
    assert isinstance(bundle["polished_sections"], list)


def test_project_result_endpoint_keeps_polish_bundle_sections_after_truncate(client, monkeypatch):
    from mellow_link.infra import ModernizationProject
    from mellow_link.infra.database import SessionLocal
    from mellow_link.infra.run_events import EVENT_TYPE_RUN_FINISHED, emit_event

    user = _register(client, username_prefix="phase1_polish_result")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id=f"phase1-polish-{uuid.uuid4().hex[:8]}",
        filename="legacy.jsp",
        content=b"<% String sql = \"SELECT * FROM orders\"; %>",
        project_name="Polish Result Project",
        client_name="OO",
    )
    structured_result, polish_bundle = _build_large_rebuild_result_with_polish_bundle()

    with SessionLocal() as db:
        project = db.query(ModernizationProject).filter(ModernizationProject.id == project_id).first()
        emit_event(
            project.run_id,
            EVENT_TYPE_RUN_FINISHED,
            {
                "success": True,
                "summary": "completed",
                "structured_result": structured_result,
                "polish_bundle": polish_bundle,
                "primary_feature_mode": "save_validation",
                "module_id": "rebuild_assistant",
                "run_kind": "rebuild_plan",
            },
            db=db,
        )

    response = client.get(
        f"/projects/{project_id}/result",
        params={"format": "json", "surface_mode": "internal"},
        headers=user["headers"],
    )

    assert response.status_code == 200, response.text
    payload = response.json()
    assert isinstance(payload["polish_bundle"], dict)
    assert isinstance(payload["polish_bundle"]["polished_sections"], list)
    assert isinstance(payload["polish_bundle"]["preserved_facts"], list)
    assert isinstance(payload["polish_bundle"]["warnings"], list)


def test_project_result_json_persists_md_and_docx_archive(client, monkeypatch, tmp_path):
    from types import SimpleNamespace

    from mellow_link import app_state
    from mellow_link.infra import ModernizationProject
    from mellow_link.infra.database import SessionLocal
    from mellow_link.infra.run_events import EVENT_TYPE_RUN_FINISHED, emit_event

    user = _register(client, username_prefix="phase1_result_archive")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id=f"phase1-result-archive-{uuid.uuid4().hex[:8]}",
        filename="legacy.jsp",
        content=b"<% String sql = \"SELECT * FROM orders\"; %>",
        project_name="결과 저장 프로젝트",
        client_name="OO아카이브",
    )
    structured_result, polish_bundle = _build_large_rebuild_result_with_polish_bundle()

    with SessionLocal() as db:
        project = db.query(ModernizationProject).filter(ModernizationProject.id == project_id).first()
        run_id = project.run_id
        emit_event(
            run_id,
            EVENT_TYPE_RUN_FINISHED,
            {
                "success": True,
                "summary": "completed",
                "structured_result": structured_result,
                "polish_bundle": polish_bundle,
                "primary_feature_mode": "save_validation",
                "module_id": "rebuild_assistant",
                "run_kind": "rebuild_plan",
            },
            db=db,
        )

    archive_dir = _project_result_archive_dir(project_id, run_id)
    if archive_dir.exists():
        shutil.rmtree(archive_dir)

    generated_docx_path = tmp_path / "generated_archive.docx"
    generated_docx_path.write_bytes(b"fake-archive-docx")

    class FakeDocService:
        def is_available(self):
            return True

        async def generate(self, request):
            assert request.output_type.value == "docx"
            assert request.filename == "결과_저장_프로젝트_result.docx"
            return SimpleNamespace(output_path=generated_docx_path)

    monkeypatch.setattr(app_state, "doc_service", FakeDocService(), raising=False)

    response = client.get(
        f"/projects/{project_id}/result",
        params={"format": "json", "surface_mode": "internal"},
        headers=user["headers"],
    )
    assert response.status_code == 200, response.text

    markdown_path = archive_dir / "result.md"
    docx_path = archive_dir / "result.docx"
    assert markdown_path.exists()
    assert docx_path.exists()
    assert "# 결과 패키지 - 결과 저장 프로젝트" in markdown_path.read_text(encoding="utf-8")
    assert docx_path.read_bytes() == b"fake-archive-docx"


def test_manual_project_run_archives_previous_completed_result(client, monkeypatch, tmp_path):
    from types import SimpleNamespace

    from mellow_link import app_state
    from mellow_link.infra import ModernizationProject
    from mellow_link.infra.database import SessionLocal
    from mellow_link.infra.run_events import EVENT_TYPE_RUN_FINISHED, emit_event

    user = _register(client, username_prefix="phase1_rerun_archive")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id=f"phase1-rerun-archive-{uuid.uuid4().hex[:8]}",
        filename="legacy.sql",
        content=b"SELECT * FROM orders;",
        project_name="재실행 이전 결과 저장",
        client_name="OO리런",
    )
    structured_result, polish_bundle = _build_large_rebuild_result_with_polish_bundle()

    with SessionLocal() as db:
        project = db.query(ModernizationProject).filter(ModernizationProject.id == project_id).first()
        previous_run_id = project.run_id
        emit_event(
            previous_run_id,
            EVENT_TYPE_RUN_FINISHED,
            {
                "success": True,
                "summary": "completed",
                "structured_result": structured_result,
                "polish_bundle": polish_bundle,
                "primary_feature_mode": "save_validation",
                "module_id": "rebuild_assistant",
                "run_kind": "rebuild_plan",
            },
            db=db,
        )

    archive_dir = _project_result_archive_dir(project_id, previous_run_id)
    if archive_dir.exists():
        shutil.rmtree(archive_dir)

    generated_docx_path = tmp_path / "generated_rerun_archive.docx"
    generated_docx_path.write_bytes(b"fake-rerun-docx")

    class FakeDocService:
        def is_available(self):
            return True

        async def generate(self, request):
            assert request.output_type.value == "docx"
            assert request.filename == "재실행_이전_결과_저장_result.docx"
            return SimpleNamespace(output_path=generated_docx_path)

    monkeypatch.setattr(app_state, "doc_service", FakeDocService(), raising=False)

    response = client.post(f"/projects/{project_id}/run", headers=user["headers"])
    assert response.status_code == 200, response.text

    markdown_path = archive_dir / "result.md"
    docx_path = archive_dir / "result.docx"
    assert markdown_path.exists()
    assert docx_path.exists()
    assert "# 결과 패키지 - 재실행 이전 결과 저장" in markdown_path.read_text(encoding="utf-8")
    assert docx_path.read_bytes() == b"fake-rerun-docx"


def test_document_service_strips_duplicate_title_heading_from_markdown():
    from mellow_link.services.doc_service import DocumentService

    service = DocumentService()
    title = "결과 패키지 - 주문 관리 화면 현대화"
    content = "# 결과 패키지 - 주문 관리 화면 현대화\n\n## 결정 요약\n핵심 판단\n"

    normalized = service._strip_duplicate_title_heading(content, title)

    assert normalized.startswith("## 결정 요약")
    assert "# 결과 패키지 - 주문 관리 화면 현대화" not in normalized


def test_build_result_package_uses_unknown_app_version_and_null_generated_at():
    from mellow_link.infra import ModernizationProject
    from mellow_link.routers.projects import build_result_package

    project = ModernizationProject(
        id="proj_test_provenance",
        user_id=1,
        session_id="sess_test",
        run_id="run_test",
        project_name="프로비넌스",
        client_name="OO",
        template_key="default_modernization_v1",
        template_mode="recommended",
        constraints_json="[]",
        upload_session_id="upload_test",
        asset_manifest_json="[]",
        status="completed",
    )
    project.created_at = None
    snapshot = {
        "run_id": "run_test",
        "module_id": "rebuild_assistant",
        "run_kind": "rebuild_plan",
        "status": "completed",
        "created_at": None,
        "updated_at": None,
    }
    assets = [{"project_asset_id": None, "name": "legacy.jsp"}]

    pkg = build_result_package(project, snapshot, None, assets=assets, app_version=None)
    assert pkg["provenance"]["app_version"] == "unknown"
    assert pkg["provenance"]["generated_at"] is None
    assert pkg["provenance"]["input_assets"] == assets
    assert pkg["provenance"]["input_assets"] is not assets
    assert pkg["executive_summary"]["state"] == "pending"
    assert pkg["executive_summary"]["core_message"] == "분석 결과를 생성 중입니다."
    assert pkg["executive_summary"]["modernization_direction"] == []
    assert pkg["executive_summary"]["key_risks"] == []
    assert pkg["decision_items"] == []
    assert pkg["design_options"] == []
    assert pkg["executive_summary"]["next_steps"] == [
        "추천안을 기준으로 현대화 방향과 분리 우선순위를 확정",
        "단일 기능·단일 화면 기준으로 파일럿 범위를 고정",
        "추천안 기준으로 상세 설계 착수 여부와 후속 자산 확보 범위를 확정",
    ]


def test_build_result_package_normalizes_generated_at_to_utc_z():
    from mellow_link.infra import ModernizationProject
    from mellow_link.routers.projects import build_result_package
    from mellow_link.modules.rebuild_assistant.schemas import StructuredRebuildResult
    from mellow_link.services.scope_notice import PROJECT_SCOPE_NOTICE

    project = ModernizationProject(
        id="proj_test_utc",
        user_id=1,
        session_id="sess_utc",
        run_id="run_utc",
        project_name="UTC",
        client_name="OO",
        template_key="default_modernization_v1",
        template_mode="recommended",
        constraints_json="[]",
        upload_session_id="upload_utc",
        asset_manifest_json="[]",
        status="completed",
    )
    project.created_at = datetime(2026, 3, 24, 12, 34, 56)
    snapshot = {
        "run_id": "run_utc",
        "module_id": "rebuild_assistant",
        "run_kind": "rebuild_plan",
        "status": "completed",
        "created_at": None,
        "updated_at": datetime(2026, 3, 24, 12, 34, 56, tzinfo=timezone.utc).isoformat(),
    }

    result = StructuredRebuildResult(
        one_line_conclusion="주요 조회 기능을 분리 현대화하는 것이 적절합니다.",
        core_business_rules=["상태 전이 규칙을 우선 추출해야 합니다."],
        recommended_directions=["방향 1", "방향 2", "방향 3", "방향 4"],
        risks=["리스크 1", "리스크 2", "리스크 3", "리스크 4"],
        rebuild_strategy=["전략 1", "전략 2", "전략 3", "전략 4"],
    )
    pkg = build_result_package(project, snapshot, result, assets=[], app_version="0.1.0")
    assert pkg["provenance"]["generated_at"] == "2026-03-24T12:34:56Z"
    assert pkg["scope_notice"] == PROJECT_SCOPE_NOTICE
    summary = pkg["executive_summary"]
    assert summary["state"] == "ready"
    assert summary["core_message"] == "주요 조회 기능을 분리 현대화하는 것이 적절합니다."
    assert pkg["core_business_rules"] == ["상태 전이 규칙을 우선 추출해야 합니다."]
    assert summary["modernization_direction"] == ["방향 1", "방향 2", "방향 3"]
    assert summary["key_risks"] == ["리스크 1", "리스크 2", "리스크 3"]
    assert "summary_lines" in summary
    assert len(summary["next_steps"]) == 3


def test_build_result_package_partial_summary_uses_missing_context_message():
    from mellow_link.infra import ModernizationProject
    from mellow_link.routers.projects import build_result_package
    from mellow_link.modules.rebuild_assistant.schemas import MissingContextItem, StructuredRebuildResult

    project = ModernizationProject(
        id="proj_test_partial",
        user_id=1,
        session_id="sess_partial",
        run_id="run_partial",
        project_name="PARTIAL",
        client_name="OO",
        template_key="default_modernization_v1",
        template_mode="recommended",
        constraints_json="[]",
        upload_session_id="upload_partial",
        asset_manifest_json="[]",
        status="completed",
    )
    result = StructuredRebuildResult(
        rebuild_strategy=["전략 A", "전략 B", "전략 C", "전략 D"],
        missing_context_details=[
            MissingContextItem(required_material="화면 정의서", reason="업무 규칙 확인 필요"),
        ],
    )
    pkg = build_result_package(project, {"status": "completed"}, result, assets=[], app_version="0.1.0")
    summary = pkg["executive_summary"]
    assert summary["state"] == "partial"
    assert summary["core_message"] == "현재까지의 분석 결과를 기반으로 한 초안입니다."
    assert isinstance(summary["modernization_direction"], list)
    assert summary["next_steps"] == [
        "화면 정의서 자료를 확보해 확인 필요 항목을 확정",
        "단일 기능·단일 화면 기준으로 파일럿 범위를 고정",
        "추천안 기준으로 상세 설계 착수 여부와 후속 자산 확보 범위를 확정",
    ]


def test_build_result_package_fallback_detail_only_does_not_mark_ready():
    from mellow_link.infra import ModernizationProject
    from mellow_link.routers.projects import build_result_package
    from mellow_link.modules.rebuild_assistant.schemas import StructuredRebuildResult

    project = ModernizationProject(
        id="proj_test_fallback_only",
        user_id=1,
        session_id="sess_fallback_only",
        run_id="run_fallback_only",
        project_name="FALLBACK_ONLY",
        client_name="OO",
        template_key="default_modernization_v1",
        template_mode="recommended",
        constraints_json="[]",
        upload_session_id="upload_fallback_only",
        asset_manifest_json="[]",
        status="completed",
    )
    result = StructuredRebuildResult(
        rebuild_strategy=["전략 A", "전략 B", "전략 C"],
        risks=["리스크 A", "리스크 B", "리스크 C"],
    )
    pkg = build_result_package(project, {"status": "completed"}, result, assets=[], app_version="0.1.0")
    summary = pkg["executive_summary"]
    assert summary["state"] == "pending"
    assert summary["modernization_direction"] == []
    assert summary["key_risks"] == ["리스크 A", "리스크 B", "리스크 C"]


def test_result_html_does_not_fallback_scope_notice(client):
    result = client.get("/projects/some-project/result", headers={"Accept": "text/html"})
    assert result.status_code == 200
    assert "STATIC_SCOPE_NOTICE" not in result.text
    assert "범위 안내 누락" in result.text
    assert "한 줄 결론" in result.text
    assert "우선 조치" in result.text
    assert "판단 근거" in result.text
    assert "실행 단계" in result.text
    assert "요약 정보 누락" in result.text
    assert 'id="optionsDetails"' in result.text
    assert 'id="planDetails"' in result.text
    assert 'id="appendixDetails"' in result.text


def test_project_creation_rejects_empty_and_duplicate_asset_manifest(client):
    user = _register(client, "phase1_validation")

    empty = client.post(
        "/projects",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "project_name": "빈 자산",
            "client_name": "OO보험",
            "upload_session_id": "phase1-empty",
            "asset_manifest": [],
            "template_key": "default_modernization_v1",
            "constraints": [],
        },
    )
    assert empty.status_code == 400, empty.text

    upload = _upload_temp_asset(
        client,
        session_id="phase1-dup",
        filename="dup.txt",
        content=b"duplicate asset",
        headers=user["headers"],
    )
    duplicate = client.post(
        "/projects",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "project_name": "중복 자산",
            "client_name": "OO보험",
            "upload_session_id": "phase1-dup",
            "asset_manifest": [
                {"name": "dup.txt", "temp_file_id": upload["temp_file_id"], "size": 10},
                {"name": "dup.txt", "temp_file_id": upload["temp_file_id"], "size": 10},
            ],
            "template_key": "default_modernization_v1",
            "constraints": [],
        },
    )
    assert duplicate.status_code == 400, duplicate.text


def test_project_assets_follow_asset_manifest_order(client, monkeypatch):
    from mellow_link import app_state
    from mellow_link.routers import projects as projects_router

    monkeypatch.setattr(projects_router, "start_project_wrapped_run", lambda *args, **kwargs: None)
    user = _register(client, "phase1_order")
    first = _upload_temp_asset(
        client,
        session_id="phase1-order",
        filename="first.txt",
        content=b"first payload",
        headers=user["headers"],
    )
    second = _upload_temp_asset(
        client,
        session_id="phase1-order",
        filename="second.txt",
        content=b"second payload",
        headers=user["headers"],
    )

    create_res = client.post(
        "/projects",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "project_name": "순서 검증",
            "client_name": "OO은행",
            "upload_session_id": "phase1-order",
            "asset_manifest": [
                {"name": "second.txt", "temp_file_id": second["temp_file_id"], "size": 14},
                {"name": "first.txt", "temp_file_id": first["temp_file_id"], "size": 13},
            ],
            "template_key": "default_modernization_v1",
            "constraints": [],
        },
    )
    assert create_res.status_code == 200, create_res.text
    project_id = create_res.json()["project_id"]

    detail = client.get(f"/projects/{project_id}?format=json", headers=user["headers"])
    assert detail.status_code == 200, detail.text
    assets = detail.json()["assets"]
    assert [asset["name"] for asset in assets] == ["second.txt", "first.txt"]

    temp_context = app_state.TEMP_CONTEXT_STORE["phase1-order"]
    assert temp_context.index("===== ASSET: second.txt =====") < temp_context.index("===== ASSET: first.txt =====")


def test_new_project_creates_initial_run_history(client, monkeypatch):
    project_id, _ = _create_persisted_project(
        client,
        _register(client, "phase1_history_initial"),
        monkeypatch,
        upload_session_id="phase1-history-initial",
        filename="initial.txt",
        content=b"initial payload",
        project_name="히스토리 초기화",
        client_name="OO초기",
    )

    rows = _project_run_history_rows(project_id)
    assert len(rows) == 1
    assert rows[0].sequence_no == 1
    assert rows[0].trigger_kind == "initial"
    assert rows[0].created_at is not None


def test_project_reanalysis_without_new_assets_creates_new_run_and_history(client, monkeypatch):
    from mellow_link.routers import projects as projects_router

    user = _register(client, "phase1_reanalysis_empty")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-reanalysis-empty",
        filename="baseline.txt",
        content=b"baseline payload",
        project_name="동일 입력 재실행",
        client_name="OO루프",
    )

    res = client.post(
        f"/projects/{project_id}/reanalysis",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={"new_asset_manifest": []},
    )
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["promoted_asset_count"] == 0
    assert body["status"] == "pending"

    detail = client.get(f"/projects/{project_id}?format=json", headers=user["headers"])
    assert detail.status_code == 200, detail.text
    payload = detail.json()
    assert payload["project"]["run_id"] == body["run_id"]
    assert payload["project"]["status"] == "pending"
    history = payload["run_history"]
    assert len(history) == 2
    assert [item["sequence_no"] for item in history] == [1, 2]
    assert history[0]["trigger_kind"] == "initial"
    assert history[1]["trigger_kind"] == "reanalysis"
    assert history[1]["is_latest"] is True

    rows = _project_run_history_rows(project_id)
    assert [row.sequence_no for row in rows] == [1, 2]


def test_project_reanalysis_promotes_new_assets_and_updates_manifest(client, monkeypatch):
    from mellow_link import app_state

    user = _register(client, "phase1_reanalysis_new")
    project_id, first_upload = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-reanalysis-new",
        filename="first.txt",
        content=b"first payload",
        project_name="신규 자산 재실행",
        client_name="OO반복",
    )
    second_upload = _upload_temp_asset(
        client,
        session_id="phase1-reanalysis-new",
        filename="second.txt",
        content=b"second payload",
        headers=user["headers"],
    )

    res = client.post(
        f"/projects/{project_id}/reanalysis",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "new_asset_manifest": [
                {"name": "second.txt", "temp_file_id": second_upload["temp_file_id"], "size": 14, "category_hint": "supplement"}
            ]
        },
    )
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["promoted_asset_count"] == 1
    assert body["latest_asset_count"] == 2

    detail = client.get(f"/projects/{project_id}?format=json", headers=user["headers"])
    assert detail.status_code == 200, detail.text
    payload = detail.json()
    assert [item["temp_file_id"] for item in payload["project"]["asset_manifest"]] == [first_upload["temp_file_id"], second_upload["temp_file_id"]]
    assert [item["name"] for item in payload["assets"]] == ["first.txt", "second.txt"]
    assert "===== ASSET: first.txt =====" in app_state.TEMP_CONTEXT_STORE["phase1-reanalysis-new"]
    assert app_state.TEMP_CONTEXT_STORE["phase1-reanalysis-new"].index("===== ASSET: first.txt =====") < app_state.TEMP_CONTEXT_STORE["phase1-reanalysis-new"].index("===== ASSET: second.txt =====")


def test_project_reanalysis_rejects_duplicate_existing_manifest_asset(client, monkeypatch):
    user = _register(client, "phase1_reanalysis_existing_dup")
    project_id, upload = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-reanalysis-existing-dup",
        filename="dup.txt",
        content=b"dup payload",
        project_name="중복 재추가 금지",
        client_name="OO검증",
    )

    res = client.post(
        f"/projects/{project_id}/reanalysis",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={
            "new_asset_manifest": [
                {"name": "dup.txt", "temp_file_id": upload["temp_file_id"], "size": 11}
            ]
        },
    )
    assert res.status_code == 400, res.text


def test_project_reanalysis_bootstraps_legacy_history_and_increments_sequence(client, monkeypatch):
    user = _register(client, "phase1_reanalysis_bootstrap")
    project_id = _create_fallback_project(user, project_name="레거시 부트스트랩")
    upload = _upload_temp_asset(
        client,
        session_id="fallback-session",
        filename="new.jsp",
        content=b"<% new supplemental asset %>",
        headers=user["headers"],
    )

    first = client.post(
        f"/projects/{project_id}/reanalysis",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={"new_asset_manifest": [{"name": "new.jsp", "temp_file_id": upload["temp_file_id"], "size": 26}]},
    )
    assert first.status_code == 200, first.text

    second = client.post(
        f"/projects/{project_id}/reanalysis",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={"new_asset_manifest": []},
    )
    assert second.status_code == 200, second.text

    history = client.get(f"/projects/{project_id}?format=json", headers=user["headers"]).json()["run_history"]
    assert [item["sequence_no"] for item in history] == [1, 2, 3]
    assert [item["trigger_kind"] for item in history] == ["initial", "reanalysis", "reanalysis"]
    assert history[-1]["is_latest"] is True


def test_project_reanalysis_run_start_failure_keeps_manifest_and_marks_latest_failed(client, monkeypatch):
    from mellow_link.infra import TempResource
    from mellow_link.infra.database import SessionLocal
    from mellow_link.routers import projects as projects_router

    user = _register(client, "phase1_reanalysis_fail")
    project_id, _ = _create_persisted_project(
        client,
        user,
        monkeypatch,
        upload_session_id="phase1-reanalysis-fail",
        filename="base.txt",
        content=b"base payload",
        project_name="재실행 실패",
        client_name="OO실패",
    )
    upload = _upload_temp_asset(
        client,
        session_id="phase1-reanalysis-fail",
        filename="extra.txt",
        content=b"extra payload",
        headers=user["headers"],
    )

    monkeypatch.setattr(projects_router, "start_project_wrapped_run", lambda *args, **kwargs: (_ for _ in ()).throw(RuntimeError("boom")))
    res = client.post(
        f"/projects/{project_id}/reanalysis",
        headers={**user["headers"], "Content-Type": "application/json"},
        json={"new_asset_manifest": [{"name": "extra.txt", "temp_file_id": upload["temp_file_id"], "size": 13}]},
    )
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["status"] == "failed"
    assert body["promoted_asset_count"] == 1

    detail = client.get(f"/projects/{project_id}?format=json", headers=user["headers"])
    assert detail.status_code == 200, detail.text
    payload = detail.json()
    assert payload["project"]["status"] == "failed"
    assert [item["name"] for item in payload["assets"]] == ["base.txt", "extra.txt"]
    assert payload["run_history"][-1]["trigger_kind"] == "reanalysis"
    assert payload["run_history"][-1]["status"] == "failed"
    assert payload["run_history"][-1]["is_latest"] is True

    with SessionLocal() as db:
        temp = db.query(TempResource).filter(TempResource.temp_file_id == upload["temp_file_id"]).first()
        assert temp is not None
        assert temp.stage_status == "promoted"


def test_sse_allows_only_authenticated_owner(client):
    owner = _register(client, "phase1_sse_owner")
    other = _register(client, "phase1_sse_other")

    create_res = client.post("/runs", headers=owner["headers"])
    run_id = create_res.json()["run_id"]
    _emit_finished(run_id, success=True, summary="finished for sse")

    unauth = client.get(f"/runs/{run_id}/events")
    assert unauth.status_code == 401

    forbidden = client.get(
        f"/runs/{run_id}/events",
        params={"access_token": other["token"]},
        headers={"Accept": "text/event-stream"},
    )
    assert forbidden.status_code == 403

    allowed = client.get(
        f"/runs/{run_id}/events",
        params={"access_token": owner["token"]},
        headers={"Accept": "text/event-stream"},
    )
    assert allowed.status_code == 200
    assert "text/event-stream" in allowed.headers.get("content-type", "")


def test_anonymization_debug_event_is_hidden_from_user_surfaces_and_visible_to_admin(client):
    from mellow_link.infra import ModernizationProject, User
    from mellow_link.infra.database import SessionLocal
    from mellow_link.infra.run_events import (
        EVENT_TYPE_DEBUG_ANONYMIZATION_REPORT,
        EVENT_TYPE_RUN_FINISHED,
        EVENT_TYPE_RUN_STARTED,
        create_run,
        emit_event,
    )
    from mellow_link.routers.runs import _resolve_run_session_id

    owner = _register(client, "phase1_anon_owner")
    admin = _register_admin(client, "phase1_anon_admin")
    anonymization_summary = {
        "applied": True,
        "policy_version": "safe_bundle_exposure_v0",
        "masking_level": "FULL",
        "total_replacements": 3,
        "canonical_source_count": 1,
        "structure_count": 1,
        "asset_counts": [{"asset_id": "asset_001", "asset_name": "legacy.jsp", "replacement_count": 3}],
        "omitted_asset_count": 0,
        "validation_passed": True,
        "risk_flags": [],
    }
    debug_payload = {
        "policy_version": "safe_bundle_exposure_v0",
        "report_summary": anonymization_summary,
        "validation": {"passed": True, "shape_preserved": True, "user_surface_safe": True, "findings": []},
        "source_previews": [
            {
                "asset_id": "asset_001",
                "asset_name": "legacy.jsp",
                "language": "jsp",
                "replacement_count": 3,
                "preview_text": "class [CLASS] { return [STRING]; }",
            }
        ],
        "bundle_debug": {
            "canonical_source_count": 1,
            "structure_count": 1,
            "total_replacements": 3,
            "masking_level": "FULL",
            "policy_version": "safe_bundle_exposure_v0",
            "omitted_preview_count": 0,
            "validation_passed": True,
        },
    }

    with SessionLocal() as db:
        db_user = db.query(User).filter(User.username == owner["username"]).first()
        session_id = _resolve_run_session_id(db, db_user, None)
        run_id = create_run(session_id=session_id, db=db, module_id="rebuild_assistant", run_kind="rebuild_plan")
        emit_event(run_id, EVENT_TYPE_RUN_STARTED, {"user_input": "project anonymization", "mode": "fast", "session_id": session_id}, db=db)
        emit_event(run_id, EVENT_TYPE_DEBUG_ANONYMIZATION_REPORT, debug_payload, db=db)
        emit_event(
            run_id,
            EVENT_TYPE_RUN_FINISHED,
            {
                "success": True,
                "summary": "completed",
                "module_id": "rebuild_assistant",
                "run_kind": "rebuild_plan",
                "anonymization_summary": anonymization_summary,
            },
            db=db,
        )
        project = ModernizationProject(
            id=f"proj_{uuid.uuid4().hex[:12]}",
            user_id=db_user.id,
            session_id=session_id,
            run_id=run_id,
            project_name="익명화 분리 프로젝트",
            client_name="OO금융",
            template_key="default_modernization_v1",
            template_mode="recommended",
            constraints_json="[]",
            upload_session_id="anonymization-session",
            asset_manifest_json='[{"name":"legacy.jsp","temp_file_id":"temp_001","size":321}]',
            status="completed",
        )
        db.add(project)
        db.commit()
        project_id = project.id

    snapshot = client.get(f"/runs/{run_id}", headers=owner["headers"])
    assert snapshot.status_code == 200, snapshot.text
    assert snapshot.json()["anonymization_summary"] == anonymization_summary

    user_events = client.get(f"/runs/{run_id}/events?format=json", headers=owner["headers"])
    assert user_events.status_code == 200, user_events.text
    assert all(item["type"] != EVENT_TYPE_DEBUG_ANONYMIZATION_REPORT for item in user_events.json()["events"])

    owner_dev = client.get(f"/runs/{run_id}/dev", headers=owner["headers"])
    assert owner_dev.status_code == 200, owner_dev.text
    assert all(item["type"] != EVENT_TYPE_DEBUG_ANONYMIZATION_REPORT for item in owner_dev.json()["events"])

    project_detail = client.get(f"/projects/{project_id}?format=json", headers=owner["headers"])
    assert project_detail.status_code == 200, project_detail.text
    assert project_detail.json()["snapshot"]["anonymization_summary"] == anonymization_summary

    admin_events = client.get(f"/api/dev/runs/{run_id}/events", headers=admin["headers"])
    assert admin_events.status_code == 200, admin_events.text
    assert any(item["type"] == EVENT_TYPE_DEBUG_ANONYMIZATION_REPORT for item in admin_events.json()["events"])

    admin_raw = client.get(f"/api/dev/runs/{run_id}/raw", headers=admin["headers"])
    assert admin_raw.status_code == 200, admin_raw.text
    assert any(item["type"] == EVENT_TYPE_DEBUG_ANONYMIZATION_REPORT for item in admin_raw.json()["events"])


def test_owned_runs_only_and_orphan_runs_hidden(client):
    user = _register(client, "phase1_visible")

    owned = client.post("/runs", headers=user["headers"]).json()
    _emit_finished(owned["run_id"], success=True, summary="owned run")

    from mellow_link.infra.database import SessionLocal, AgentRun
    from datetime import datetime

    orphan_id = f"run_orphan_{uuid.uuid4().hex[:8]}"
    with SessionLocal() as db:
        db.add(AgentRun(run_id=orphan_id, session_id=None, status="completed", created_at=datetime.utcnow(), updated_at=datetime.utcnow(), summary="orphan"))
        db.commit()

    list_res = client.get("/runs", headers=user["headers"])
    assert list_res.status_code == 200
    ids = {item["run_id"] for item in list_res.json()["runs"]}
    assert owned["run_id"] in ids
    assert orphan_id not in ids


def test_pilot_security_notice_doc_exists_and_covers_current_facts():
    doc_path = MELLOW_LINK_ROOT / "docs" / "PILOT_SECURITY_AND_OPERATIONS_NOTICE.md"
    assert doc_path.exists()
    text = doc_path.read_text(encoding="utf-8")
    assert "단일 기능 / 단일 화면" in text
    assert "partial" in text
    assert "제품 워크스페이스 기준 상대 경로 구조" in text
    assert "mellow_link/data/aventurine_v3.db" in text
    assert "mellow_link/data/temp_uploads" in text
    assert "mellow_link/data/project_assets" in text
    assert "Bearer JWT" in text
    assert "프로젝트 소유자 기준" in text
    assert "프로젝트 종료 시점까지 보관" in text
    assert "자동 만료 정책은 현재 구현에 포함되어 있지 않다." in text
    assert "자동 삭제는 미지원이다." in text
    assert "수동 관리 대상" in text
    assert "현재 구현 기준" in text
    assert "현재 비지원 / 한계" in text


def test_docs_index_lists_pilot_security_notice():
    readme_path = MELLOW_LINK_ROOT / "docs" / "README.md"
    text = readme_path.read_text(encoding="utf-8")
    assert "PILOT_SECURITY_AND_OPERATIONS_NOTICE.md" in text
    assert "고객 제출/설명용 문서" in text



