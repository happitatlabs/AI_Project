"""
Document Service - CPU-based Document Generation

This module provides document generation capabilities for Word documents.
Runs ENTIRELY on CPU in a thread pool to avoid blocking the async event loop.

CRITICAL: All CPU-bound operations MUST use loop.run_in_executor() to prevent
blocking the async event loop while GPU tasks (LLM/Image) are running.

Libraries:
    - Word: python-docx
    - PDF: reportlab (optional - basic implementation)

Design:
    - All CPU-bound operations run in ThreadPoolExecutor
    - Non-blocking async interface
    - Does NOT touch the GPU
"""

import asyncio
import logging
import time
from concurrent.futures import ThreadPoolExecutor
from typing import Optional, Dict, Any, List
from dataclasses import dataclass, field
from enum import Enum, auto
from pathlib import Path
from datetime import datetime

logger = logging.getLogger(__name__)


# =============================================================================
# Enums and Data Classes
# =============================================================================

class DocumentType(Enum):
    """Supported document output types."""
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    HTML = "html"
    MARKDOWN = "md"


class DocumentStatus(Enum):
    """Document service status."""
    IDLE = auto()
    PROCESSING = auto()
    ERROR = auto()


@dataclass
class DocumentRequest:
    """
    Request structure for document generation.

    Attributes:
        content: Main content (text or Markdown)
        output_type: Desired output format (default: DOCX)
        title: Document title
        metadata: Additional metadata (author, date, etc.)
        style_options: Formatting options (fonts, margins, etc.)
        filename: Optional custom filename
    """
    content: str
    output_type: DocumentType = DocumentType.DOCX
    title: str = "Untitled"
    metadata: Dict[str, str] = field(default_factory=dict)
    style_options: Dict[str, Any] = field(default_factory=dict)
    filename: Optional[str] = None
    payload: Any = None


@dataclass
class DocumentResult:
    """
    Result structure from document generation.

    Attributes:
        output_path: Path to generated document
        output_type: Type of document generated
        page_count: Estimated number of pages
        file_size_bytes: Size of generated file
        generation_time_ms: Time to generate
    """
    output_path: Path
    output_type: DocumentType
    page_count: int = 0
    file_size_bytes: int = 0
    generation_time_ms: float = 0.0


class DocumentGenerationError(Exception):
    """Exception for document generation failures."""
    pass


# =============================================================================
# Document Service Class
# =============================================================================

class DocumentService:
    """
    Service for CPU-based document generation.

    CRITICAL DESIGN:
        All CPU-intensive operations run in a ThreadPoolExecutor via
        loop.run_in_executor(). This ensures the async event loop is
        NEVER blocked, allowing GPU tasks to continue uninterrupted.

    This service:
        - Generates Word documents (python-docx)
        - Generates basic PDF documents (reportlab)
        - Generates HTML from Markdown
        - Does NOT use the GPU

    Usage:
        service = DocumentService()
        await service.initialize()

        request = DocumentRequest(content="Hello World", title="Test")
        result = await service.generate(request)

        await service.shutdown()
    """

    DEFAULT_OUTPUT_DIR: Path = Path("./outputs/documents")
    DEFAULT_THREAD_WORKERS: int = 2  # Keep low - CPU only

    def __init__(
        self,
        output_dir: Optional[Path] = None,
        max_workers: int = DEFAULT_THREAD_WORKERS
    ):
        """
        Initialize Document Service.

        Args:
            output_dir: Directory for generated documents
            max_workers: Maximum thread pool workers (default: 2)
        """
        self.output_dir = output_dir or self.DEFAULT_OUTPUT_DIR
        self.max_workers = max_workers

        self._status: DocumentStatus = DocumentStatus.IDLE
        self._executor: Optional[ThreadPoolExecutor] = None
        self._is_initialized: bool = False

    # ==================== Lifecycle ====================

    async def initialize(self) -> None:
        """
        Initialize the document service.

        Creates:
            - Output directory
            - Thread pool for CPU-bound work
        """
        logger.info("[DocumentService] Initializing...")

        # Create output directory
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Initialize thread pool for CPU-bound operations
        self._executor = ThreadPoolExecutor(
            max_workers=self.max_workers,
            thread_name_prefix="doc_worker"
        )

        self._is_initialized = True
        self._status = DocumentStatus.IDLE
        logger.info(f"[DocumentService] Initialized with {self.max_workers} workers")

    async def shutdown(self) -> None:
        """Shutdown the document service and cleanup thread pool."""
        logger.info("[DocumentService] Shutting down...")

        if self._executor:
            self._executor.shutdown(wait=True)
            self._executor = None

        self._is_initialized = False
        logger.info("[DocumentService] Shutdown complete")

    async def connect(self) -> bool:
        """Connect alias for orchestrator compatibility."""
        await self.initialize()
        return True

    async def disconnect(self) -> None:
        """Disconnect alias for orchestrator compatibility."""
        await self.shutdown()

    def get_status(self) -> DocumentStatus:
        """Get current service status."""
        return self._status

    def is_ready(self) -> bool:
        """Check if service is ready."""
        return self._is_initialized and self._status == DocumentStatus.IDLE

    def is_available(self) -> bool:
        """Check if service is available (orchestrator compatibility)."""
        return self._is_initialized

    # ==================== Document Generation ====================

    async def generate(self, request: DocumentRequest) -> DocumentResult:
        """
        Generate a document from request.

        CRITICAL: This method delegates CPU work to run_in_executor()
        to avoid blocking the async event loop.

        Args:
            request: DocumentRequest with content and options

        Returns:
            DocumentResult with generated file path

        Raises:
            DocumentGenerationError: If generation fails
        """
        if not self._is_initialized:
            raise DocumentGenerationError("DocumentService not initialized")

        self._status = DocumentStatus.PROCESSING
        start_time = time.time()

        try:
            # Generate output path
            output_path = self._get_output_path(
                request.filename or request.title,
                request.output_type
            )

            # Generate based on type - ALL use run_in_executor
            if request.output_type == DocumentType.DOCX:
                result = await self._generate_docx(
                    request.content,
                    output_path,
                    request.title,
                    request.style_options
                )
            elif request.output_type == DocumentType.PPTX:
                result = await self._generate_pptx(
                    request.content,
                    output_path,
                    request.title,
                    request.style_options,
                    request.payload,
                )
            elif request.output_type == DocumentType.PDF:
                result = await self._generate_pdf(
                    request.content,
                    output_path,
                    request.title,
                    request.style_options
                )
            elif request.output_type == DocumentType.HTML:
                result = await self._generate_html(
                    request.content,
                    output_path,
                    request.title
                )
            else:
                # Markdown - simple write
                result = await self._generate_markdown(
                    request.content,
                    output_path
                )

            result.generation_time_ms = (time.time() - start_time) * 1000

            logger.info(
                f"[DocumentService] Generated {request.output_type.value}: "
                f"{output_path.name} ({result.generation_time_ms:.0f}ms)"
            )
            return result

        except Exception as e:
            logger.error(f"[DocumentService] Generation failed: {e}")
            raise DocumentGenerationError(str(e))

        finally:
            self._status = DocumentStatus.IDLE

    async def _generate_docx(
        self,
        content: str,
        output_path: Path,
        title: str,
        options: Dict[str, Any]
    ) -> DocumentResult:
        """
        Generate Word document using python-docx.

        CRITICAL: Uses run_in_executor to run in thread pool.
        """
        loop = asyncio.get_event_loop()

        def _create_docx() -> int:
            """CPU-bound DOCX creation - runs in thread pool."""
            try:
                from docx import Document
                from docx.shared import Pt
                from docx.enum.text import WD_ALIGN_PARAGRAPH
            except ImportError:
                raise DocumentGenerationError(
                    "python-docx not installed. Run: pip install python-docx"
                )

            doc = Document()

            normalized_content = self._strip_duplicate_title_heading(content, title)

            # Add title
            doc.add_heading(title, level=0)

            # Configure style
            style = doc.styles['Normal']
            font = style.font
            font.name = options.get("font_name", "Calibri")
            font.size = Pt(options.get("font_size", 11))

            # Add structured content from markdown-like result package
            for block in self._parse_markdown_outline(normalized_content):
                if block["kind"] == "heading":
                    doc.add_heading(block["text"], level=min(block["level"], 4))
                    continue
                if block["kind"] == "paragraph":
                    p = doc.add_paragraph(block["text"])
                    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    continue
                if block["kind"] == "bullets":
                    for item in block["items"]:
                        doc.add_paragraph(item, style="List Bullet")

            doc.save(str(output_path))

            # Estimate page count
            return max(1, len(normalized_content) // 3000)

        # CRITICAL: Run in executor to avoid blocking async loop
        page_count = await loop.run_in_executor(self._executor, _create_docx)

        return DocumentResult(
            output_path=output_path,
            output_type=DocumentType.DOCX,
            page_count=page_count,
            file_size_bytes=output_path.stat().st_size
        )

    async def _generate_pptx(
        self,
        content: str,
        output_path: Path,
        title: str,
        options: Dict[str, Any],
        payload: Any = None,
    ) -> DocumentResult:
        """
        Generate PPTX document using python-pptx.

        CRITICAL: Uses run_in_executor to run in thread pool.
        """
        loop = asyncio.get_event_loop()

        def _create_pptx() -> int:
            """CPU-bound PPTX creation - runs in thread pool."""
            try:
                from pptx import Presentation
                from pptx.util import Pt
            except ImportError:
                raise DocumentGenerationError(
                    "python-pptx not installed. Run: pip install python-pptx"
                )

            prs = Presentation()

            if options.get("renderer") == "slide_schema" and isinstance(payload, dict):
                self._render_slide_schema_pptx(prs, payload, title=title, font_size=options.get("font_size", 20))
                prs.save(str(output_path))
                return max(1, len(prs.slides))

            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = title
            subtitle = title_slide.placeholders[1]
            subtitle.text = options.get("subtitle", "결과 패키지 요약 프레젠테이션")

            normalized_content = self._strip_duplicate_title_heading(content, title)
            sections = self._markdown_sections_for_slides(normalized_content)
            for section_title, lines in sections:
                chunks = self._chunk_slide_lines(lines, size=7)
                for index, chunk in enumerate(chunks):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = section_title if index == 0 else f"{section_title} ({index + 1})"
                    text_frame = slide.placeholders[1].text_frame
                    text_frame.clear()
                    for line_index, line in enumerate(chunk):
                        paragraph = text_frame.paragraphs[0] if line_index == 0 else text_frame.add_paragraph()
                        paragraph.text = line
                        paragraph.font.size = Pt(options.get("font_size", 20))
                        paragraph.level = 1 if line.startswith("- ") else 0

            prs.save(str(output_path))
            return max(1, len(prs.slides))

        page_count = await loop.run_in_executor(self._executor, _create_pptx)

        return DocumentResult(
            output_path=output_path,
            output_type=DocumentType.PPTX,
            page_count=page_count,
            file_size_bytes=output_path.stat().st_size
        )

    async def _generate_pdf(
        self,
        content: str,
        output_path: Path,
        title: str,
        options: Dict[str, Any]
    ) -> DocumentResult:
        """
        Generate PDF document using reportlab.

        CRITICAL: Uses run_in_executor to run in thread pool.
        """
        loop = asyncio.get_event_loop()

        def _create_pdf() -> int:
            """CPU-bound PDF creation - runs in thread pool."""
            try:
                from reportlab.lib.pagesizes import A4
                from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                from reportlab.lib.units import cm
                from reportlab.lib.enums import TA_JUSTIFY
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            except ImportError:
                raise DocumentGenerationError(
                    "reportlab not installed. Run: pip install reportlab"
                )

            # Create document
            doc = SimpleDocTemplate(
                str(output_path),
                pagesize=A4,
                leftMargin=2.5 * cm,
                rightMargin=2.5 * cm,
                topMargin=2.5 * cm,
                bottomMargin=2.5 * cm
            )

            # Styles
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=30
            )
            body_style = ParagraphStyle(
                'CustomBody',
                parent=styles['Normal'],
                fontSize=options.get("font_size", 11),
                alignment=TA_JUSTIFY,
                spaceBefore=6,
                spaceAfter=6,
                leading=14
            )

            # Build content
            story = []
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 12))

            # Add paragraphs
            for para in content.split('\n\n'):
                if para.strip():
                    para_text = para.replace('\n', '<br/>')
                    story.append(Paragraph(para_text, body_style))

            doc.build(story)

            # Estimate page count
            return max(1, output_path.stat().st_size // 3000)

        # CRITICAL: Run in executor to avoid blocking async loop
        page_count = await loop.run_in_executor(self._executor, _create_pdf)

        return DocumentResult(
            output_path=output_path,
            output_type=DocumentType.PDF,
            page_count=page_count,
            file_size_bytes=output_path.stat().st_size
        )

    async def _generate_html(
        self,
        content: str,
        output_path: Path,
        title: str
    ) -> DocumentResult:
        """
        Generate HTML document.

        CRITICAL: Uses run_in_executor to run in thread pool.
        """
        loop = asyncio.get_event_loop()

        def _create_html():
            """CPU-bound HTML creation - runs in thread pool."""
            # Try to convert markdown if available
            html_content = content
            try:
                import markdown
                html_content = markdown.markdown(
                    content,
                    extensions=['tables', 'fenced_code', 'nl2br']
                )
            except ImportError:
                # Basic fallback
                html_content = content.replace('\n\n', '</p><p>')
                html_content = f"<p>{html_content}</p>"

            html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        body {{
            font-family: 'Segoe UI', Tahoma, sans-serif;
            max-width: 800px;
            margin: 0 auto;
            padding: 20px;
            line-height: 1.6;
        }}
        h1 {{ color: #333; }}
        p {{ text-align: justify; }}
        code {{ background: #f4f4f4; padding: 2px 6px; border-radius: 3px; }}
        pre {{ background: #f4f4f4; padding: 15px; overflow-x: auto; }}
    </style>
</head>
<body>
    <h1>{title}</h1>
    {html_content}
</body>
</html>"""
            output_path.write_text(html, encoding='utf-8')

        # CRITICAL: Run in executor to avoid blocking async loop
        await loop.run_in_executor(self._executor, _create_html)

        return DocumentResult(
            output_path=output_path,
            output_type=DocumentType.HTML,
            page_count=1,
            file_size_bytes=output_path.stat().st_size
        )

    async def _generate_markdown(
        self,
        content: str,
        output_path: Path
    ) -> DocumentResult:
        """
        Save content as Markdown file.

        Even simple I/O uses run_in_executor for consistency.
        """
        loop = asyncio.get_event_loop()

        def _write_file():
            output_path.write_text(content, encoding='utf-8')

        # CRITICAL: Run in executor to avoid blocking async loop
        await loop.run_in_executor(self._executor, _write_file)

        return DocumentResult(
            output_path=output_path,
            output_type=DocumentType.MARKDOWN,
            page_count=1,
            file_size_bytes=output_path.stat().st_size
        )

    # ==================== Orchestrator Interface ====================

    async def execute(self, request_data: Dict[str, Any]) -> DocumentResult:
        """
        Execute method for orchestrator compatibility.

        Args:
            request_data: Dict with generation parameters

        Returns:
            DocumentResult
        """
        doc_type_str = request_data.get("output_type", "docx").lower()
        doc_type_map = {
            "pdf": DocumentType.PDF,
            "docx": DocumentType.DOCX,
            "pptx": DocumentType.PPTX,
            "html": DocumentType.HTML,
            "md": DocumentType.MARKDOWN,
            "markdown": DocumentType.MARKDOWN,
        }

        request = DocumentRequest(
            content=request_data.get("content", ""),
            output_type=doc_type_map.get(doc_type_str, DocumentType.DOCX),
            title=request_data.get("title", "Document"),
            metadata=request_data.get("metadata", {}),
            style_options=request_data.get("style_options", {}),
            filename=request_data.get("filename"),
            payload=request_data.get("payload"),
        )
        return await self.generate(request)

    # ==================== Utilities ====================

    def _get_output_path(self, filename: str, doc_type: DocumentType) -> Path:
        """Generate unique output path for document."""
        # Clean filename
        safe_name = "".join(c for c in filename if c.isalnum() or c in (' ', '-', '_'))
        safe_name = safe_name.strip() or "document"

        # Add timestamp for uniqueness
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        return self.output_dir / f"{safe_name}_{timestamp}.{doc_type.value}"

    def _parse_markdown_outline(self, content: str) -> List[Dict[str, Any]]:
        """Parse simple markdown headings and bullet groups into outline blocks."""
        blocks: List[Dict[str, Any]] = []
        paragraph_lines: List[str] = []
        bullet_lines: List[str] = []

        def flush_paragraph() -> None:
            nonlocal paragraph_lines
            if paragraph_lines:
                blocks.append({"kind": "paragraph", "text": " ".join(paragraph_lines).strip()})
                paragraph_lines = []

        def flush_bullets() -> None:
            nonlocal bullet_lines
            if bullet_lines:
                blocks.append({"kind": "bullets", "items": bullet_lines[:]})
                bullet_lines = []

        for raw_line in content.splitlines():
            line = raw_line.strip()
            if not line:
                flush_paragraph()
                flush_bullets()
                continue
            if line.startswith("#"):
                flush_paragraph()
                flush_bullets()
                level = len(line) - len(line.lstrip("#"))
                blocks.append({"kind": "heading", "level": max(1, level), "text": line[level:].strip()})
                continue
            if line.startswith("- "):
                flush_paragraph()
                bullet_lines.append(line[2:].strip())
                continue
            if line in ("[필요 자료]", "[이유]"):
                flush_paragraph()
                flush_bullets()
                blocks.append({"kind": "heading", "level": 4, "text": line})
                continue
            paragraph_lines.append(line)

        flush_paragraph()
        flush_bullets()
        return blocks

    def _strip_duplicate_title_heading(self, content: str, title: str) -> str:
        """Remove the first H1 when it duplicates the explicit document title."""
        lines = content.splitlines()
        if not lines:
            return content
        normalized_title = (title or "").strip()
        for index, raw_line in enumerate(lines):
            line = raw_line.strip()
            if not line:
                continue
            if line.startswith("# ") and line[2:].strip() == normalized_title:
                return "\n".join(lines[:index] + lines[index + 1 :]).lstrip("\n")
            break
        return content

    def _markdown_sections_for_slides(self, content: str) -> List[tuple[str, List[str]]]:
        """Group markdown-like content into slide sections using ## headings."""
        sections: List[tuple[str, List[str]]] = []
        current_title = "요약"
        current_lines: List[str] = []

        for raw_line in content.splitlines():
            line = raw_line.strip()
            if not line or line.startswith("# "):
                continue
            if line.startswith("## "):
                if current_lines:
                    sections.append((current_title, current_lines[:]))
                current_title = line[3:].strip()
                current_lines = []
                continue
            if line.startswith("### "):
                current_lines.append(line[4:].strip() + ":")
                continue
            current_lines.append(line)

        if current_lines:
            sections.append((current_title, current_lines))
        return sections

    def _chunk_slide_lines(self, lines: List[str], size: int = 7) -> List[List[str]]:
        """Split section lines into slide-sized chunks."""
        if not lines:
            return [["내용 없음"]]
        return [lines[index:index + size] for index in range(0, len(lines), size)]

    def _render_slide_schema_pptx(
        self,
        prs: Any,
        payload: Dict[str, Any],
        *,
        title: str,
        font_size: int,
    ) -> None:
        from pptx.util import Inches

        from mellow_link.modules.rebuild_assistant.postprocess.schemas import SlideSchemaDeck

        deck = SlideSchemaDeck.model_validate(payload)
        blank_layout = prs.slide_layouts[6]
        for slide_schema in deck.slides:
            slide = prs.slides.add_slide(blank_layout)
            self._render_slide_schema_header(slide, slide_schema, fallback_title=title)
            if slide_schema.slide_type == "overview":
                self._render_overview_slide(slide, slide_schema, font_size=font_size)
            elif slide_schema.slide_type == "as_is_gap":
                self._render_as_is_gap_slide(slide, slide_schema, font_size=font_size)
            elif slide_schema.slide_type == "flow":
                self._render_flow_slide(slide, slide_schema, font_size=font_size)
            elif slide_schema.slide_type == "design":
                self._render_design_slide(slide, slide_schema, font_size=font_size)
            elif slide_schema.slide_type == "vision":
                self._render_vision_slide(slide, slide_schema, font_size=font_size)
            if slide_schema.absorbed_summary_text:
                self._render_absorbed_summary_note(slide, slide_schema)

    def _render_slide_schema_header(self, slide: Any, slide_schema: Any, *, fallback_title: str) -> None:
        from pptx.util import Inches

        theme = self._slide_type_theme(slide_schema.slide_type)
        rendered_title = slide_schema.title or fallback_title
        if slide_schema.is_continuation or slide_schema.sequence > 1:
            rendered_title = f"{rendered_title} ({slide_schema.sequence})"
        self._add_textbox(
            slide,
            Inches(0.45),
            Inches(0.3),
            Inches(6.6),
            Inches(0.45),
            rendered_title,
            font_size=28,
            bold=True,
            text_rgb=theme["accent"],
        )
        if slide_schema.is_continuation:
            self._render_continuation_badge(slide, slide_schema, theme=theme)
        if slide_schema.headline:
            self._render_headline_band(slide, slide_schema.headline, theme=theme)

    def _render_continuation_badge(self, slide: Any, slide_schema: Any, *, theme: Dict[str, tuple[int, int, int]]) -> None:
        from pptx.util import Inches

        badge_text = "독립 연속 장표" if slide_schema.continuation_value == "retain" else "보조 연속 장표"
        self._add_textbox(
            slide,
            Inches(6.95),
            Inches(0.26),
            Inches(1.9),
            Inches(0.44),
            badge_text,
            font_size=11,
            bold=True,
            fill_rgb=theme["accent"],
            text_rgb=(255, 255, 255),
            line_rgb=theme["accent"],
            center=True,
        )

    def _render_headline_band(self, slide: Any, headline: str, *, theme: Dict[str, tuple[int, int, int]]) -> None:
        from pptx.util import Inches

        self._add_textbox(
            slide,
            Inches(0.45),
            Inches(0.84),
            Inches(8.35),
            Inches(0.52),
            headline,
            font_size=16,
            fill_rgb=theme["panel"],
            line_rgb=theme["border"],
            text_rgb=theme["accent"],
        )

    def _render_overview_slide(self, slide: Any, slide_schema: Any, *, font_size: int) -> None:
        from pptx.util import Inches

        items = list(slide_schema.context_bullets or [])
        scope_items = list(slide_schema.scope_bullets or [])
        constraint_items = list(slide_schema.constraint_bullets or [])
        layout = slide_schema.layout_hint or "overview_context_only"
        if layout == "overview_context_scope_constraints":
            self._add_bullet_box(slide, Inches(0.55), Inches(1.4), Inches(5.0), Inches(4.8), "현행 요약", items, font_size=font_size, style_key="overview")
            self._add_bullet_box(slide, Inches(5.8), Inches(1.4), Inches(3.0), Inches(2.15), "범위", scope_items, font_size=max(14, font_size - 2), style_key="overview_meta")
            self._add_bullet_box(slide, Inches(5.8), Inches(3.8), Inches(3.0), Inches(2.15), "제약", constraint_items, font_size=max(14, font_size - 2), style_key="overview_meta")
            return
        if layout == "overview_context_scope":
            self._add_bullet_box(slide, Inches(0.55), Inches(1.4), Inches(5.6), Inches(4.8), "현행 요약", items, font_size=font_size, style_key="overview")
            self._add_bullet_box(
                slide,
                Inches(6.35),
                Inches(1.4),
                Inches(2.7),
                Inches(4.8),
                "범위 / 제약",
                scope_items + constraint_items,
                font_size=max(14, font_size - 2),
                style_key="overview_meta",
            )
            return
        if layout == "overview_meta_split":
            self._add_bullet_box(slide, Inches(0.55), Inches(1.4), Inches(4.1), Inches(4.8), "추가 배경", items, font_size=max(14, font_size - 1), style_key="overview")
            self._add_bullet_box(slide, Inches(4.9), Inches(1.4), Inches(1.95), Inches(4.8), "범위", scope_items, font_size=max(13, font_size - 3), style_key="overview_meta")
            self._add_bullet_box(slide, Inches(7.0), Inches(1.4), Inches(2.0), Inches(4.8), "제약", constraint_items, font_size=max(13, font_size - 3), style_key="overview_meta")
            return
        self._add_bullet_box(slide, Inches(0.55), Inches(1.4), Inches(8.4), Inches(4.9), "현행 요약", items, font_size=font_size, style_key="overview")

    def _render_as_is_gap_slide(self, slide: Any, slide_schema: Any, *, font_size: int) -> None:
        from pptx.util import Inches

        if slide_schema.layout_hint == "gap_risk_continuation":
            self._add_bullet_box(
                slide,
                Inches(0.55),
                Inches(1.4),
                Inches(5.0),
                Inches(4.2),
                "핵심 GAP",
                list(slide_schema.gap_bullets or []),
                font_size=max(14, font_size - 2),
                style_key="approach_core",
            )
            self._add_bullet_box(
                slide,
                Inches(5.8),
                Inches(1.4),
                Inches(3.1),
                Inches(4.2),
                "검토 포인트",
                list(slide_schema.risk_bullets or []),
                font_size=13,
                style_key="approach_risk",
            )
            if slide_schema.decision_message:
                self._add_textbox(slide, Inches(0.65), Inches(5.8), Inches(8.1), Inches(0.55), slide_schema.decision_message, font_size=13, bold=True, fill_rgb=(247, 244, 240), line_rgb=(190, 172, 154))
            return
        self._add_bullet_box(
            slide,
            Inches(0.45),
            Inches(1.4),
            Inches(2.7),
            Inches(3.65),
            "AS-IS",
            list(slide_schema.as_is_bullets or []),
            font_size=max(14, font_size - 2),
            style_key="approach_as_is",
        )
        self._add_bullet_box(
            slide,
            Inches(3.35),
            Inches(1.4),
            Inches(2.7),
            Inches(3.65),
            "GAP",
            list(slide_schema.gap_bullets or []),
            font_size=max(14, font_size - 2),
            style_key="approach_core",
        )
        self._add_bullet_box(
            slide,
            Inches(6.25),
            Inches(1.4),
            Inches(2.7),
            Inches(3.65),
            "TO-BE",
            list(slide_schema.to_be_bullets or []),
            font_size=max(14, font_size - 2),
            style_key="approach_to_be",
        )
        if slide_schema.risk_bullets:
            self._add_bullet_box(
                slide,
                Inches(0.55),
                Inches(5.25),
                Inches(8.5),
                Inches(1.35),
                "검토 포인트",
                list(slide_schema.risk_bullets or []),
                font_size=13,
                style_key="approach_risk",
            )

    def _render_flow_slide(self, slide: Any, slide_schema: Any, *, font_size: int) -> None:
        from pptx.util import Inches

        steps = list(slide_schema.steps or [])
        if slide_schema.layout_hint == "timeline_horizontal" and steps:
            card_width = Inches(1.95)
            left_positions = [Inches(0.55), Inches(2.72), Inches(4.89), Inches(7.06)]
            for index, step in enumerate(steps[:4]):
                self._add_bullet_box(
                    slide,
                    left_positions[index],
                    Inches(1.65),
                    card_width,
                    Inches(2.9),
                    step.step_label,
                    [step.step_text],
                    font_size=max(13, font_size - 3),
                    style_key="flow_step",
                )
            if slide_schema.action_bullets:
                self._add_bullet_box(slide, Inches(0.75), Inches(4.95), Inches(8.0), Inches(1.35), "중점 실행 과제", list(slide_schema.action_bullets or []), font_size=13, style_key="flow_actions")
            return
        if slide_schema.layout_hint in ("stacked_flow", "flow_continuation"):
            step_lines = [f"{step.step_label} | {step.step_text}" for step in steps]
            self._add_bullet_box(
                slide,
                Inches(0.55),
                Inches(1.4),
                Inches(5.6),
                Inches(4.85),
                "단계별 추진 흐름",
                step_lines,
                font_size=max(14, font_size - 1),
                style_key="flow_stack",
            )
            if slide_schema.action_bullets:
                self._add_bullet_box(slide, Inches(6.4), Inches(1.4), Inches(2.55), Inches(2.6), "중점 실행 과제", list(slide_schema.action_bullets or []), font_size=13, style_key="flow_actions")
            if slide_schema.footer_note:
                self._add_textbox(slide, Inches(6.4), Inches(4.35), Inches(2.45), Inches(0.9), slide_schema.footer_note, font_size=12, bold=False, fill_rgb=(242, 245, 248), line_rgb=(176, 188, 201))
            return
        step_lines = [f"{step.step_label} | {step.step_text}" for step in steps]
        self._add_bullet_box(slide, Inches(0.55), Inches(1.4), Inches(8.4), Inches(4.5), "단계별 추진 흐름", step_lines, font_size=max(14, font_size - 1), style_key="flow_stack")

    def _render_design_slide(self, slide: Any, slide_schema: Any, *, font_size: int) -> None:
        from pptx.util import Inches

        if not list(slide_schema.rule_cards or []) and slide_schema.flow_bullets:
            self._add_bullet_box(
                slide,
                Inches(0.75),
                Inches(1.75),
                Inches(7.95),
                Inches(3.8),
                "계산 / 전표 / GL 흐름",
                list(slide_schema.flow_bullets or []),
                font_size=14,
                style_key="design_flow",
            )
            return
        use_meta_sidebar = slide_schema.layout_hint == "rule_cards_with_meta_sidebar"
        positions = [
            (Inches(0.55), Inches(1.4)),
            (Inches(3.9), Inches(1.4)),
            (Inches(0.55), Inches(3.4)),
            (Inches(3.9), Inches(3.4)),
        ] if use_meta_sidebar else [
            (Inches(0.55), Inches(1.4)),
            (Inches(4.85), Inches(1.4)),
            (Inches(0.55), Inches(3.4)),
            (Inches(4.85), Inches(3.4)),
        ]
        card_width = Inches(2.95) if use_meta_sidebar else Inches(3.7)
        for index, card in enumerate(slide_schema.rule_cards or []):
            if index >= len(positions):
                break
            left, top = positions[index]
            self._add_bullet_box(
                slide,
                left,
                top,
                card_width,
                Inches(1.55),
                card.title,
                [card.body],
                font_size=13,
                style_key="design_card",
            )
        if slide_schema.layout_hint == "rule_cards_with_meta_sidebar":
            self._add_bullet_box(
                slide,
                Inches(7.15),
                Inches(1.4),
                Inches(1.7),
                Inches(2.1),
                "엔터티",
                list(slide_schema.entity_blocks or []),
                font_size=11,
                style_key="design_meta",
            )
            self._add_bullet_box(
                slide,
                Inches(7.15),
                Inches(3.75),
                Inches(1.7),
                Inches(2.1),
                "연계",
                list(slide_schema.interface_points or []),
                font_size=11,
                style_key="design_meta",
            )
            if slide_schema.flow_bullets:
                self._add_bullet_box(
                    slide,
                    Inches(0.55),
                    Inches(5.25),
                    Inches(6.3),
                    Inches(1.2),
                    "계산 / 전표 / GL 흐름",
                    list(slide_schema.flow_bullets or []),
                    font_size=12,
                    style_key="design_flow",
                )
            return
        if slide_schema.flow_bullets:
            self._add_bullet_box(slide, Inches(0.55), Inches(5.25), Inches(8.3), Inches(1.2), "계산 / 전표 / GL 흐름", list(slide_schema.flow_bullets or []), font_size=13, style_key="design_flow")

    def _render_vision_slide(self, slide: Any, slide_schema: Any, *, font_size: int) -> None:
        from pptx.util import Inches

        items = list(slide_schema.future_state_bullets or [])
        if slide_schema.layout_hint == "future_state_with_effect":
            self._add_bullet_box(slide, Inches(0.6), Inches(1.6), Inches(4.1), Inches(3.9), "적용 방향", items, font_size=max(14, font_size - 1), style_key="vision_core")
            self._add_bullet_box(slide, Inches(4.95), Inches(1.6), Inches(3.95), Inches(3.9), "기대 효과", list(slide_schema.effect_bullets or []), font_size=13, style_key="vision_effect")
            if slide_schema.closing_statement:
                self._add_textbox(slide, Inches(0.8), Inches(5.8), Inches(8.0), Inches(0.45), slide_schema.closing_statement, font_size=13, bold=True, fill_rgb=(243, 244, 241), line_rgb=(181, 188, 176))
            return
        if slide_schema.layout_hint == "future_state_pillars" and len(items) <= 3 and items:
            widths = [Inches(2.55), Inches(2.55), Inches(2.55)]
            for index, item in enumerate(items):
                self._add_bullet_box(
                    slide,
                    Inches(0.6 + (index * 2.85)),
                    Inches(1.75),
                    widths[index],
                    Inches(3.8),
                    f"방향 {index + 1}",
                    [item],
                    font_size=14,
                    style_key="vision_pillar",
                )
            if slide_schema.closing_statement:
                self._add_textbox(slide, Inches(0.85), Inches(5.8), Inches(7.9), Inches(0.45), slide_schema.closing_statement, font_size=12, bold=True, fill_rgb=(243, 244, 241), line_rgb=(181, 188, 176))
            return
        self._add_bullet_box(slide, Inches(0.6), Inches(1.6), Inches(8.3), Inches(4.6), "적용 방향", items, font_size=max(14, font_size - 2), style_key="vision_core")
        if slide_schema.closing_statement:
            self._add_textbox(slide, Inches(0.8), Inches(5.8), Inches(8.0), Inches(0.45), slide_schema.closing_statement, font_size=12, bold=True, fill_rgb=(243, 244, 241), line_rgb=(181, 188, 176))

    def _render_absorbed_summary_note(self, slide: Any, slide_schema: Any) -> None:
        from pptx.util import Inches

        theme = self._slide_type_theme(slide_schema.slide_type)
        self._add_textbox(
            slide,
            Inches(0.55),
            Inches(6.62),
            Inches(8.25),
            Inches(0.42),
            f"보강 메모 | {slide_schema.absorbed_summary_text}",
            font_size=10,
            bold=False,
            fill_rgb=theme["panel"],
            line_rgb=theme["border"],
            text_rgb=theme["accent"],
        )

    def _slide_type_theme(self, slide_type: str) -> Dict[str, tuple[int, int, int]]:
        palette = {
            "overview": {"accent": (78, 92, 108), "panel": (244, 246, 248), "border": (185, 194, 203)},
            "as_is_gap": {"accent": (132, 88, 61), "panel": (248, 243, 239), "border": (202, 182, 168)},
            "flow": {"accent": (77, 103, 133), "panel": (241, 246, 250), "border": (179, 193, 206)},
            "design": {"accent": (76, 103, 86), "panel": (241, 245, 242), "border": (181, 193, 184)},
            "vision": {"accent": (96, 90, 78), "panel": (246, 245, 241), "border": (193, 189, 181)},
        }
        return palette.get(slide_type, palette["overview"])

    def _add_bullet_box(
        self,
        slide: Any,
        left: Any,
        top: Any,
        width: Any,
        height: Any,
        heading: str,
        items: List[str],
        *,
        font_size: int,
        style_key: str = "",
    ) -> None:
        from pptx.util import Pt

        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame
        text_frame.clear()
        self._apply_text_frame_style(text_frame, style_key=style_key)
        fill_rgb, line_rgb, heading_rgb = self._box_style(style_key)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(fill_rgb)
        shape.line.color.rgb = self._rgb(line_rgb)
        if heading:
            first = text_frame.paragraphs[0]
            first.text = heading
            first.font.size = Pt(font_size + 1)
            first.font.bold = True
            first.font.color.rgb = self._rgb(heading_rgb)
            self._apply_paragraph_style(first, role="heading", style_key=style_key, font_size=font_size + 1)
        for item in items:
            paragraph = text_frame.add_paragraph()
            paragraph.text = f"- {item}"
            paragraph.font.size = Pt(font_size if not self._is_meta_style(style_key) else max(11, font_size - 1))
            paragraph.level = 0
            self._apply_paragraph_style(paragraph, role="body", style_key=style_key, font_size=font_size)

    def _add_textbox(
        self,
        slide: Any,
        left: Any,
        top: Any,
        width: Any,
        height: Any,
        text: str,
        *,
        font_size: int,
        bold: bool = False,
        fill_rgb: tuple[int, int, int] | None = None,
        line_rgb: tuple[int, int, int] | None = None,
        text_rgb: tuple[int, int, int] | None = None,
        center: bool = False,
    ) -> None:
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt

        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame
        text_frame.clear()
        self._apply_text_frame_style(text_frame, style_key="textbox")
        if fill_rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(fill_rgb)
        if line_rgb:
            shape.line.color.rgb = self._rgb(line_rgb)
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        self._apply_paragraph_style(paragraph, role="textbox", style_key="textbox", font_size=font_size)
        if text_rgb:
            paragraph.font.color.rgb = self._rgb(text_rgb)
        if center:
            paragraph.alignment = PP_ALIGN.CENTER

    def _apply_text_frame_style(self, text_frame: Any, *, style_key: str) -> None:
        from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
        from pptx.util import Inches

        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        compact = self._is_meta_style(style_key)
        text_frame.margin_left = Inches(0.1 if compact else 0.12)
        text_frame.margin_right = Inches(0.08 if compact else 0.1)
        text_frame.margin_top = Inches(0.05 if compact else 0.07)
        text_frame.margin_bottom = Inches(0.04 if compact else 0.06)

    def _apply_paragraph_style(self, paragraph: Any, *, role: str, style_key: str, font_size: int) -> None:
        from pptx.util import Pt

        if role == "heading":
            paragraph.space_after = Pt(5)
            paragraph.space_before = Pt(0)
            return
        if role == "textbox":
            paragraph.space_after = Pt(0)
            paragraph.space_before = Pt(0)
            paragraph.line_spacing = 1.02
            return
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(2 if self._is_meta_style(style_key) else 3)
        paragraph.line_spacing = 1.05 if self._is_meta_style(style_key) else 1.08

    def _is_meta_style(self, style_key: str) -> bool:
        return style_key in {"overview_meta", "design_meta"} or style_key.endswith("_meta")

    def _box_style(self, style_key: str) -> tuple[tuple[int, int, int], tuple[int, int, int], tuple[int, int, int]]:
        styles = {
            "overview": ((244, 246, 248), (185, 194, 203), (78, 92, 108)),
            "overview_meta": ((249, 250, 251), (198, 205, 212), (78, 92, 108)),
            "approach_as_is": ((247, 244, 240), (206, 190, 178), (120, 93, 72)),
            "approach_core": ((248, 240, 235), (192, 165, 146), (132, 88, 61)),
            "approach_to_be": ((243, 247, 242), (181, 196, 182), (84, 110, 88)),
            "approach_risk": ((250, 244, 241), (207, 180, 171), (132, 88, 61)),
            "flow_step": ((241, 246, 250), (179, 193, 206), (77, 103, 133)),
            "flow_actions": ((246, 248, 250), (190, 201, 211), (77, 103, 133)),
            "flow_stack": ((242, 247, 251), (179, 193, 206), (77, 103, 133)),
            "design_card": ((241, 245, 242), (181, 193, 184), (76, 103, 86)),
            "design_meta": ((247, 249, 247), (190, 199, 192), (76, 103, 86)),
            "design_flow": ((244, 248, 244), (181, 193, 184), (76, 103, 86)),
            "vision_core": ((246, 245, 241), (193, 189, 181), (96, 90, 78)),
            "vision_effect": ((248, 247, 243), (203, 198, 189), (96, 90, 78)),
            "vision_pillar": ((246, 245, 241), (193, 189, 181), (96, 90, 78)),
        }
        return styles.get(style_key, ((248, 248, 248), (204, 204, 204), (80, 80, 80)))

    def _rgb(self, color: tuple[int, int, int]) -> Any:
        from pptx.dml.color import RGBColor

        return RGBColor(*color)

    async def health_check(self) -> Dict[str, Any]:
        """Perform health check."""
        return {
            "healthy": self._is_initialized,
            "status": self._status.name,
            "output_dir": str(self.output_dir),
            "max_workers": self.max_workers,
        }

    def to_dict(self) -> Dict[str, Any]:
        """Export service state as dictionary."""
        return {
            "status": self._status.name,
            "is_initialized": self._is_initialized,
            "output_dir": str(self.output_dir),
            "max_workers": self.max_workers,
        }


# =============================================================================
# Factory Function
# =============================================================================

def create_document_service(
    output_dir: Optional[Path] = None,
    max_workers: int = 2
) -> DocumentService:
    """
    Factory function to create DocumentService.

    Args:
        output_dir: Output directory for documents
        max_workers: Thread pool size (default: 2)

    Returns:
        Configured DocumentService instance
    """
    return DocumentService(
        output_dir=output_dir,
        max_workers=max_workers
    )
