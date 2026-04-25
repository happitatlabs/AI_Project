from __future__ import annotations

import logging
import re
import struct
import tempfile
import zipfile
from io import BytesIO
from pathlib import Path
from typing import Iterable
from xml.etree import ElementTree as ET

logger = logging.getLogger(__name__)

_PLACEHOLDER_PREFIXES = (
    "click to ",
    "클릭하여 ",
    "텍스트를 입력",
)
_LEGACY_MASTER_PLACEHOLDERS = (
    "click to edit master title style",
    "click to edit master text styles",
    "마스터 제목 스타일 편집",
    "마스터 텍스트 스타일을 편집합니다",
)
_RT_SLIDE_PERSIST_ATOM = 1011
_RT_TEXT_HEADER_ATOM = 3999
_RT_TEXT_CHARS_ATOM = 4000
_RT_TEXT_BYTES_ATOM = 4008
_TEXT_TYPE_TITLE = {0, 6}
_TEXT_TYPE_NOTES = {2}
_LEGACY_NOISE_VALUES = {"*", "•", "-", "―"}


def extract_presentation_sml(file_path: Path, content_bytes: bytes | None = None) -> str:
    """Extract PPT/PPTX into slide-aware SML text for anonymization and analysis input.

    SML is a canonical text representation used only as a pre-analysis intermediate.
    The original PPT/PPTX binary should not be passed to downstream analysis engines.
    """

    suffix = file_path.suffix.lower()
    if suffix == ".pptx":
        return _extract_pptx_sml(file_path, content_bytes, source_name=file_path.name)
    if suffix == ".ppt":
        return _extract_legacy_ppt_sml(file_path, content_bytes)
    return ""


def _extract_pptx_sml(
    file_path: Path,
    content_bytes: bytes | None = None,
    *,
    source_name: str | None = None,
) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("[PresentationExtraction] python-pptx not installed")
        return ""

    try:
        presentation = Presentation(BytesIO(content_bytes)) if content_bytes else Presentation(str(file_path))
        return _render_sml(presentation, source_name=source_name or file_path.name)
    except Exception as exc:
        fallback = _extract_pptx_sml_from_zip(file_path, content_bytes, source_name=source_name or file_path.name)
        if fallback:
            logger.info("[PresentationExtraction] PPTX zip fallback used for %s after parser failure: %s", file_path, exc)
        else:
            logger.error("[PresentationExtraction] PPTX extraction failed for %s: %s", file_path, exc)
        return fallback


def _extract_pptx_sml_from_zip(
    file_path: Path,
    content_bytes: bytes | None = None,
    *,
    source_name: str,
) -> str:
    try:
        package = BytesIO(content_bytes) if content_bytes is not None else str(file_path)
        with zipfile.ZipFile(package) as archive:
            slide_names = sorted(
                (
                    name
                    for name in archive.namelist()
                    if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
                ),
                key=_pptx_slide_sort_key,
            )
            if not slide_names:
                return ""
            lines = [
                "[SML v1]",
                f"presentation_file: {source_name or '-'}",
                f"slide_count: {len(slide_names)}",
                "extraction_mode: pptx_zip_fallback",
            ]
            for index, slide_name in enumerate(slide_names, start=1):
                slide_xml = archive.read(slide_name)
                slide_lines = _extract_pptx_text_lines(slide_xml)
                title = slide_lines[0] if slide_lines else ""
                body = slide_lines[1:] if len(slide_lines) > 1 else []
                notes_name = f"ppt/notesSlides/notesSlide{index}.xml"
                notes = _extract_pptx_text_lines(archive.read(notes_name)) if notes_name in archive.namelist() else []
                lines.extend(["", f"[SLIDE {index}]"])
                if title:
                    lines.append(f"title: {title}")
                if body:
                    lines.append("texts:")
                    lines.extend(f"- {item}" for item in body)
                elif title:
                    lines.append("texts:")
                    lines.append(f"- {title}")
                if notes:
                    lines.append("notes:")
                    lines.extend(f"- {item}" for item in notes)
                if not title and not body and not notes:
                    lines.append("content: [NO_EXTRACTABLE_TEXT]")
            return "\n".join(lines).strip()
    except Exception as exc:
        logger.error("[PresentationExtraction] PPTX zip fallback failed for %s: %s", file_path, exc)
        return ""


def _pptx_slide_sort_key(name: str) -> int:
    match = re.search(r"slide(\d+)\.xml$", name)
    return int(match.group(1)) if match else 0


def _extract_pptx_text_lines(xml_bytes: bytes) -> list[str]:
    try:
        root = ET.fromstring(xml_bytes)
    except Exception:
        return []
    paragraphs: list[str] = []
    current_runs: list[str] = []
    for element in root.iter():
        if element.tag.endswith("}br"):
            current_runs.append("\n")
            continue
        if element.tag.endswith("}t"):
            value = _clean_text(element.text or "")
            if value:
                current_runs.append(value)
            continue
        if element.tag.endswith("}p"):
            paragraph = _clean_text("".join(current_runs).replace("\n", " ").strip())
            if paragraph:
                paragraphs.append(paragraph)
            current_runs = []
    if current_runs:
        paragraph = _clean_text("".join(current_runs).replace("\n", " ").strip())
        if paragraph:
            paragraphs.append(paragraph)
    return _dedupe_preserve_order(paragraphs)


def _extract_legacy_ppt_sml(file_path: Path, content_bytes: bytes | None = None) -> str:
    with tempfile.TemporaryDirectory(prefix="ppt_extract_") as temp_dir:
        temp_root = Path(temp_dir)
        input_path = file_path
        if content_bytes is not None:
            input_path = temp_root / (file_path.name or "legacy.ppt")
            input_path.write_bytes(content_bytes)

        converted_path = temp_root / f"{input_path.stem or 'converted'}.pptx"
        if not _convert_ppt_to_pptx_with_powerpoint(input_path, converted_path):
            stream_sml = _extract_legacy_ppt_sml_from_structured_storage(input_path, source_name=file_path.name)
            if stream_sml:
                logger.info("[PresentationExtraction] Legacy PPT stream fallback used for %s", file_path)
                return stream_sml
            binary_sml = _extract_legacy_ppt_sml_from_binary_scan(
                input_path,
                content_bytes=content_bytes,
                source_name=file_path.name,
            )
            if binary_sml:
                logger.info("[PresentationExtraction] Legacy PPT binary scan fallback used for %s", file_path)
                return binary_sml
            logger.warning("[PresentationExtraction] Legacy PPT conversion unavailable for %s", file_path)
            return ""
        return _extract_pptx_sml(converted_path, None, source_name=file_path.name)


def _convert_ppt_to_pptx_with_powerpoint(input_path: Path, output_path: Path) -> bool:
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        return False

    app = None
    presentation = None
    try:
        pythoncom.CoInitialize()
        app = win32com.client.Dispatch("PowerPoint.Application")
        app.Visible = 0
        presentation = app.Presentations.Open(str(input_path), ReadOnly=1, Untitled=0, WithWindow=0)
        presentation.SaveAs(str(output_path), 24)
        return output_path.exists()
    except Exception as exc:
        logger.info("[PresentationExtraction] PowerPoint conversion unavailable for %s: %s", input_path, exc)
        return False
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def _extract_legacy_ppt_sml_from_structured_storage(
    file_path: Path,
    *,
    source_name: str,
) -> str:
    try:
        import pythoncom
        from win32com.storagecon import STGM_DIRECT, STGM_READ, STGM_SHARE_EXCLUSIVE
    except ImportError:
        return ""

    try:
        storage = pythoncom.StgOpenStorage(
            str(file_path),
            None,
            STGM_DIRECT | STGM_READ | STGM_SHARE_EXCLUSIVE,
        )
        stream = storage.OpenStream(
            "PowerPoint Document",
            None,
            STGM_READ | STGM_SHARE_EXCLUSIVE,
        )
    except Exception as exc:
        logger.debug("[PresentationExtraction] Structured storage fallback failed for %s: %s", file_path, exc)
        return ""

    chunks: list[bytes] = []
    try:
        while True:
            chunk = stream.Read(65536)
            if not chunk:
                break
            chunks.append(chunk)
    except Exception as exc:
        logger.debug("[PresentationExtraction] Failed reading PowerPoint Document stream for %s: %s", file_path, exc)
        return ""

    stream_bytes = b"".join(chunks)
    if not stream_bytes:
        return ""
    return _render_legacy_stream_sml(stream_bytes, source_name=source_name)


def _render_legacy_stream_sml(stream_bytes: bytes, *, source_name: str) -> str:
    slides: list[dict[str, list[str] | str]] = []
    current_slide: dict[str, list[str] | str] | None = None
    current_text_type: int | None = None
    slide_persist_count = 0

    for record in _iter_legacy_ppt_records(stream_bytes):
        if record["type"] == _RT_SLIDE_PERSIST_ATOM:
            slide_persist_count += 1
            continue
        if record["type"] == _RT_TEXT_HEADER_ATOM:
            payload = record["payload"]
            if len(payload) == 4:
                current_text_type = struct.unpack_from("<I", payload, 0)[0]
            continue
        if record["type"] not in {_RT_TEXT_CHARS_ATOM, _RT_TEXT_BYTES_ATOM} or current_text_type is None:
            continue

        decoded = _decode_legacy_text_payload(record["type"], record["payload"])
        lines = [_clean_text(line) for line in _split_candidate_lines(decoded)]
        lines = [line for line in lines if _is_useful_legacy_text(line)]
        if not lines:
            continue

        if current_text_type in _TEXT_TYPE_TITLE:
            title = lines[0]
            if current_slide and _legacy_slide_has_content(current_slide):
                if _should_drop_sparse_leading_slide(current_slide):
                    slides.pop()
                current_slide = None
            if current_slide is None:
                current_slide = {"title": title, "texts": [], "notes": []}
                slides.append(current_slide)
            else:
                current_slide["title"] = title
            for line in lines[1:]:
                _append_unique(current_slide["texts"], line)
            continue

        if current_slide is None:
            current_slide = {"title": "", "texts": [], "notes": []}
            slides.append(current_slide)
        target = current_slide["notes"] if current_text_type in _TEXT_TYPE_NOTES else current_slide["texts"]
        for line in lines:
            _append_unique(target, line)

    slides = [slide for slide in slides if _legacy_slide_has_content(slide)]
    if not slides:
        return ""

    lines = [
        "[SML v1]",
        f"presentation_file: {source_name or '-'}",
        f"slide_count: {len(slides)}",
        "extraction_mode: legacy_binary_stream_fallback",
    ]
    if slide_persist_count and slide_persist_count != len(slides):
        lines.append(f"source_slide_count_estimate: {slide_persist_count}")

    for index, slide in enumerate(slides, start=1):
        title = str(slide.get("title") or "").strip()
        body_texts = _dedupe_preserve_order(slide.get("texts") or [])
        notes = _dedupe_preserve_order(slide.get("notes") or [])
        lines.extend(["", f"[SLIDE {index}]"])
        if title:
            lines.append(f"title: {title}")
        if body_texts:
            lines.append("texts:")
            lines.extend(f"- {item}" for item in body_texts)
        if notes:
            lines.append("notes:")
            lines.extend(f"- {item}" for item in notes)
        if not any((title, body_texts, notes)):
            lines.append("content: [NO_EXTRACTABLE_TEXT]")
    return "\n".join(lines).strip()


def _iter_legacy_ppt_records(data: bytes, start: int = 0, end: int | None = None) -> Iterable[dict[str, int | bytes]]:
    limit = len(data) if end is None else min(len(data), end)
    offset = start
    while offset + 8 <= limit:
        try:
            rec_ver_inst, rec_type, rec_len = struct.unpack_from("<HHI", data, offset)
        except struct.error:
            return
        payload_start = offset + 8
        payload_end = payload_start + rec_len
        if payload_end > limit:
            return
        rec_ver = rec_ver_inst & 0x000F
        payload = data[payload_start:payload_end]
        yield {
            "offset": offset,
            "type": rec_type,
            "version": rec_ver,
            "payload": payload,
        }
        if rec_ver == 0x0F and payload:
            yield from _iter_legacy_ppt_records(data, payload_start, payload_end)
        offset = payload_end


def _decode_legacy_text_payload(record_type: int, payload: bytes) -> str:
    if record_type == _RT_TEXT_CHARS_ATOM:
        return payload.decode("utf-16le", errors="ignore")
    if record_type == _RT_TEXT_BYTES_ATOM:
        return payload.decode("latin1", errors="ignore")
    return ""


def _legacy_slide_has_content(slide: dict[str, list[str] | str]) -> bool:
    return bool(str(slide.get("title") or "").strip() or list(slide.get("texts") or []) or list(slide.get("notes") or []))


def _should_drop_sparse_leading_slide(slide: dict[str, list[str] | str]) -> bool:
    title = str(slide.get("title") or "").strip()
    texts = [str(item).strip() for item in list(slide.get("texts") or []) if str(item).strip()]
    notes = [str(item).strip() for item in list(slide.get("notes") or []) if str(item).strip()]
    if title:
        return False
    combined = texts + notes
    if not combined:
        return True
    return len(combined) == 1 and len(combined[0]) <= 64


def _append_unique(target: list[str], value: str) -> None:
    normalized = _clean_text(value)
    if normalized and normalized not in target:
        target.append(normalized)


def _is_useful_legacy_text(value: str) -> bool:
    normalized = _clean_text(value)
    if not normalized:
        return False
    lowered = normalized.lower()
    if any(lowered.startswith(prefix) for prefix in _PLACEHOLDER_PREFIXES):
        return False
    if lowered in _LEGACY_MASTER_PLACEHOLDERS:
        return False
    if normalized in _LEGACY_NOISE_VALUES:
        return False
    if normalized.isdigit() and len(normalized) <= 2:
        return False
    if len(normalized) == 1 and not normalized.isalpha():
        return False
    return True


def _extract_legacy_ppt_sml_from_binary_scan(
    file_path: Path,
    *,
    content_bytes: bytes | None = None,
    source_name: str,
) -> str:
    try:
        raw_bytes = content_bytes if content_bytes is not None else file_path.read_bytes()
    except Exception:
        return ""

    extracted_lines = _dedupe_preserve_order(
        [
            *_extract_utf16le_string_candidates(raw_bytes),
            *_extract_single_byte_string_candidates(raw_bytes),
        ]
    )
    extracted_lines = [line for line in extracted_lines if _is_useful_legacy_text(line)]
    if not extracted_lines:
        return ""

    lines = [
        "[SML v1]",
        f"presentation_file: {source_name or '-'}",
        "slide_count: 1",
        "extraction_mode: legacy_binary_scan_fallback",
        "",
        "[SLIDE 1]",
        "texts:",
    ]
    lines.extend(f"- {item}" for item in extracted_lines[:120])
    return "\n".join(lines).strip()


def _extract_utf16le_string_candidates(data: bytes) -> list[str]:
    candidates: list[str] = []
    for start_offset in (0, 1):
        current: list[str] = []
        for index in range(start_offset, len(data) - 1, 2):
            code_unit = data[index] | (data[index + 1] << 8)
            char = chr(code_unit)
            if _is_reasonable_legacy_char(char):
                current.append(char)
                continue
            candidate = _clean_text("".join(current))
            if _is_reasonable_legacy_string(candidate):
                candidates.append(candidate)
            current = []
        candidate = _clean_text("".join(current))
        if _is_reasonable_legacy_string(candidate):
            candidates.append(candidate)
    return candidates


def _extract_single_byte_string_candidates(data: bytes) -> list[str]:
    candidates: list[str] = []
    for match in re.finditer(rb"[ -~]{4,}", data):
        candidate = _clean_text(match.group(0).decode("latin1", errors="ignore"))
        if _is_reasonable_legacy_string(candidate):
            candidates.append(candidate)
    return candidates


def _is_reasonable_legacy_char(char: str) -> bool:
    if not char or char in {"\x00", "\uffff"}:
        return False
    if char in {"\r", "\n", "\t", "\v"}:
        return True
    code = ord(char)
    if 0x20 <= code <= 0x7E:
        return True
    if 0xAC00 <= code <= 0xD7A3:
        return True
    if 0x3131 <= code <= 0x318E:
        return True
    if 0x4E00 <= code <= 0x9FFF:
        return True
    return False


def _is_reasonable_legacy_string(value: str) -> bool:
    if not value or len(value) < 4:
        return False
    lowered = value.lower()
    if lowered in {"powerpoint document", "current user", "summaryinformation", "documentsummaryinformation"}:
        return False
    printable_count = sum(1 for char in value if _is_reasonable_legacy_char(char))
    if printable_count / max(len(value), 1) < 0.85:
        return False
    return True


def _render_sml(presentation, *, source_name: str) -> str:
    lines = [
        "[SML v1]",
        f"presentation_file: {source_name or '-'}",
        f"slide_count: {len(presentation.slides)}",
    ]
    for index, slide in enumerate(presentation.slides, start=1):
        title = _clean_text(getattr(getattr(slide.shapes, "title", None), "text", ""))
        body_texts = _dedupe_preserve_order(_iter_slide_texts(slide, exclude={title} if title else set()))
        table_blocks = list(_iter_table_blocks(slide))
        chart_blocks = list(_iter_chart_blocks(slide))
        notes = _dedupe_preserve_order(_iter_notes_text(slide, exclude={title, *body_texts}))
        visual_markers = _dedupe_preserve_order(_iter_visual_markers(slide))
        layout_name = _clean_text(getattr(getattr(slide, "slide_layout", None), "name", ""))

        lines.extend(["", f"[SLIDE {index}]"])
        if layout_name:
            lines.append(f"layout: {layout_name}")
        if title:
            lines.append(f"title: {title}")
        if body_texts:
            lines.append("texts:")
            lines.extend(f"- {item}" for item in body_texts)
        if table_blocks:
            lines.append("tables:")
            lines.extend(table_blocks)
        if chart_blocks:
            lines.append("charts:")
            lines.extend(chart_blocks)
        if notes:
            lines.append("notes:")
            lines.extend(f"- {item}" for item in notes)
        if visual_markers and not any((body_texts, table_blocks, chart_blocks)):
            lines.append("visual_elements:")
            lines.extend(f"- {item}" for item in visual_markers)
        if not any((title, body_texts, table_blocks, chart_blocks, notes, visual_markers)):
            lines.append("content: [NO_EXTRACTABLE_TEXT]")
    return "\n".join(lines).strip()


def _iter_slide_texts(slide, *, exclude: set[str]) -> Iterable[str]:
    for shape in getattr(slide, "shapes", []):
        yield from _iter_shape_texts(shape, exclude=exclude)


def _iter_shape_texts(shape, *, exclude: set[str]) -> Iterable[str]:
    if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
        for child in shape.shapes:
            yield from _iter_shape_texts(child, exclude=exclude)
        return
    if getattr(shape, "has_table", False) or getattr(shape, "has_chart", False):
        return

    raw_text = ""
    if hasattr(shape, "text_frame") and getattr(shape, "has_text_frame", False):
        raw_text = getattr(shape.text_frame, "text", "")
    elif hasattr(shape, "text"):
        raw_text = getattr(shape, "text", "")

    for line in _split_candidate_lines(raw_text):
        if line in exclude:
            continue
        yield line


def _iter_table_blocks(slide) -> Iterable[str]:
    for shape in getattr(slide, "shapes", []):
        if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
            for child in shape.shapes:
                yield from _iter_table_blocks_from_shape(child)
            continue
        yield from _iter_table_blocks_from_shape(shape)


def _iter_table_blocks_from_shape(shape) -> Iterable[str]:
    if not getattr(shape, "has_table", False):
        return
    rows: list[str] = []
    for row in shape.table.rows:
        cell_values = [_clean_text(cell.text) for cell in row.cells]
        if any(cell_values):
            rows.append("| " + " | ".join(value or "-" for value in cell_values) + " |")
    if not rows:
        return
    yield f"- table_rows: {len(rows)}"
    for row in rows:
        yield f"  {row}"


def _iter_chart_blocks(slide) -> Iterable[str]:
    for shape in getattr(slide, "shapes", []):
        if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
            for child in shape.shapes:
                yield from _iter_chart_blocks_from_shape(child)
            continue
        yield from _iter_chart_blocks_from_shape(shape)


def _iter_chart_blocks_from_shape(shape) -> Iterable[str]:
    if not getattr(shape, "has_chart", False):
        return
    chart = shape.chart
    chart_title = ""
    if getattr(chart, "has_title", False):
        chart_title = _clean_text(getattr(chart.chart_title.text_frame, "text", ""))
    series_names = [_clean_text(getattr(series, "name", "")) for series in getattr(chart, "series", [])]
    series_names = [item for item in series_names if item]
    summary = []
    if chart_title:
        summary.append(f"title={chart_title}")
    if series_names:
        summary.append("series=" + ", ".join(series_names))
    yield "- " + ("; ".join(summary) if summary else "chart")


def _iter_notes_text(slide, *, exclude: set[str]) -> Iterable[str]:
    try:
        notes_text = getattr(slide.notes_slide.notes_text_frame, "text", "")
    except Exception:
        return
    for line in _split_candidate_lines(notes_text):
        if line in exclude:
            continue
        yield line


def _iter_visual_markers(slide) -> Iterable[str]:
    for shape in getattr(slide, "shapes", []):
        if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
            for child in shape.shapes:
                yield from _shape_visual_markers(child)
            continue
        yield from _shape_visual_markers(shape)


def _shape_visual_markers(shape) -> Iterable[str]:
    if getattr(shape, "has_chart", False):
        yield "chart"
    elif getattr(shape, "has_table", False):
        yield "table"
    elif hasattr(shape, "image"):
        yield "image"
    elif hasattr(shape, "shape_type"):
        marker = str(shape.shape_type).strip()
        if marker:
            yield marker.lower()


def _split_candidate_lines(text: str) -> list[str]:
    candidates: list[str] = []
    for raw_line in (text or "").replace("\v", "\n").splitlines():
        normalized = _clean_text(raw_line)
        if not normalized:
            continue
        lowered = normalized.lower()
        if any(lowered.startswith(prefix) for prefix in _PLACEHOLDER_PREFIXES):
            continue
        candidates.append(normalized)
    return candidates


def _clean_text(value: str) -> str:
    return " ".join((value or "").replace("\r", "\n").split())


def _dedupe_preserve_order(items: Iterable[str]) -> list[str]:
    seen: set[str] = set()
    deduped: list[str] = []
    for item in items:
        normalized = _clean_text(item)
        if not normalized or normalized in seen:
            continue
        seen.add(normalized)
        deduped.append(normalized)
    return deduped
